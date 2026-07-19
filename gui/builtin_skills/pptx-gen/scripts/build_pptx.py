"""Build a .pptx deck from a JSON outline using python-pptx.

Input JSON shape:
    {"title": "...", "subtitle": "...(optional)",
     "slides": [{"title": "...", "bullets": ["..."], "notes": "...(optional)"}]}

Usage:
    python build_pptx.py --input outline.json --output deck.pptx
    python build_pptx.py --input outline.json --output ./out_dir
    echo '{...}' | python build_pptx.py --output deck.pptx   (stdin fallback)
"""
import argparse
import json
import os
import re
import sys

# XML 1.0 allows: #x9 | #xA | #xD | [#x20-#xD7FF] | [#xE000-#xFFFD] | [#x10000-#x10FFFF]
# Remove all characters outside this set to prevent lxml ValueError.
_ILLEGAL_XML_CHARS_RE = re.compile(
    r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x84\x86-\x9f'
    r'\ud800-\udfff\ufdd0-\ufdef\ufffe\uffff]'
)

# Default font: widely available on Windows/macOS and renders Chinese well.
# python-pptx sets the latin typeface via font.name; East Asian text needs the
# a:ea typeface attribute too, otherwise PowerPoint falls back per-run.
_DEFAULT_FONT = "Microsoft YaHei"


def _clean_xml_text(text):
    """Remove characters that are not valid in XML 1.0."""
    if not text:
        return text
    return _ILLEGAL_XML_CHARS_RE.sub('', str(text))


def _set_run_font(run, size=None, bold=None):
    from pptx.oxml.ns import qn
    run.font.name = _DEFAULT_FONT
    if size is not None:
        run.font.size = size
    if bold is not None:
        run.font.bold = bold
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', _DEFAULT_FONT)


def _load_outline(args):
    if args.input:
        if not os.path.isfile(args.input):
            print(f"ERROR: Input JSON file not found: {args.input}", file=sys.stderr)
            sys.exit(1)
        with open(args.input, "r", encoding="utf-8") as f:
            raw = f.read()
    else:
        raw = sys.stdin.read()
    try:
        data = json.loads(raw)
    except json.JSONDecodeError as exc:
        print(f"ERROR: Invalid outline JSON: {exc}", file=sys.stderr)
        sys.exit(1)
    if not isinstance(data, dict):
        print("ERROR: Outline JSON must be an object with 'title' and 'slides'.", file=sys.stderr)
        sys.exit(1)
    slides = data.get("slides")
    if not isinstance(slides, list) or not slides:
        print("ERROR: Outline JSON must contain a non-empty 'slides' array.", file=sys.stderr)
        sys.exit(1)
    return data


def _resolve_output_path(output, title):
    if os.path.isdir(output):
        safe = re.sub(r'[\\/:*?"<>|]+', "_", (title or "presentation").strip()) or "presentation"
        output = os.path.join(output, f"{safe}.pptx")
    elif not output.lower().endswith('.pptx'):
        output += '.pptx'
    os.makedirs(os.path.dirname(os.path.abspath(output)), exist_ok=True)
    return output


def build_pptx(outline, output_path):
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except ImportError:
        print("ERROR: python-pptx not installed. Run: pip install python-pptx", file=sys.stderr)
        sys.exit(1)

    title = _clean_xml_text(outline.get("title") or "演示文稿")
    subtitle = _clean_xml_text(outline.get("subtitle") or "")
    slides = outline["slides"]

    prs = Presentation()

    # Title slide (layout 0: title + subtitle).
    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = title
    for run in slide.shapes.title.text_frame.paragraphs[0].runs:
        _set_run_font(run, size=Pt(32), bold=True)
    if subtitle and len(slide.placeholders) > 1:
        sub = slide.placeholders[1]
        sub.text = subtitle
        for para in sub.text_frame.paragraphs:
            for run in para.runs:
                _set_run_font(run, size=Pt(18))

    # Content slides (layout 1: title + bullet body).
    body_layout = prs.slide_layouts[1]
    for idx, item in enumerate(slides, start=1):
        if not isinstance(item, dict):
            continue
        slide = prs.slides.add_slide(body_layout)
        slide_title = _clean_xml_text(item.get("title") or f"第 {idx} 页")
        slide.shapes.title.text = slide_title
        for run in slide.shapes.title.text_frame.paragraphs[0].runs:
            _set_run_font(run, size=Pt(24), bold=True)

        bullets = item.get("bullets") or []
        if isinstance(bullets, str):
            bullets = [bullets]
        body = slide.placeholders[1].text_frame
        body.clear()
        for i, bullet in enumerate(bullets):
            text = _clean_xml_text(bullet)
            if not text:
                continue
            para = body.paragraphs[0] if i == 0 else body.add_paragraph()
            run = para.add_run()
            run.text = text
            _set_run_font(run, size=Pt(16))

        notes = _clean_xml_text(item.get("notes") or "")
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    prs.save(output_path)
    return output_path


def main():
    parser = argparse.ArgumentParser(description="Build a .pptx deck from a JSON outline")
    parser.add_argument("--input", default="", help="Outline JSON file path (omit to read from stdin)")
    parser.add_argument("--output", default=".", help="Output .pptx file path or directory")
    args = parser.parse_args()

    outline = _load_outline(args)
    output_path = _resolve_output_path(args.output, outline.get("title"))
    result = build_pptx(outline, output_path)
    print(f"PPTX 生成完成: {result}")
    print(f"   页数: {len(outline['slides']) + 1}")
    print(f"   输出: {os.path.abspath(result)}")


if __name__ == "__main__":
    main()
