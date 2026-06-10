#!/usr/bin/env python3
"""Step 5: Generate group meeting PPT from review content and images."""

import argparse, json, os, sys, glob, re
from typing import Optional

# PPT template colors
SLIDE_TITLE_COLOR = "1B3A5C"      # Deep navy
SLIDE_SUBTITLE_COLOR = "5B8DB8"   # Steel blue
CONTENT_TEXT_COLOR = "333333"     # Dark gray
ACCENT_COLOR = "C0504D"           # Academic red
BG_COLOR = "F5F5F0"              # Warm white
TABLE_HEADER_BG = "1B3A5C"
TABLE_ROW_ALT = "EBF0F5"

SLIDE_WIDTH = 13.333  # 16:9
SLIDE_HEIGHT = 7.5

def create_ppt(review: dict, images_dir: str, output_path: str, lang: str = "zh"):
    """Generate group meeting PPT with python-pptx."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.enum.shapes import MSO_SHAPE
    except ImportError:
        print("ERROR: python-pptx not installed. Run: pip install python-pptx", file=sys.stderr)
        sys.exit(1)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH)
    prs.slide_height = Inches(SLIDE_HEIGHT)

    # Use blank layout
    blank_layout = prs.slide_layouts[6]  # Blank

    def add_bg(slide, color=BG_COLOR):
        """Set slide background color."""
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(
            int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
        )

    def add_text_box(slide, left, top, width, height, text,
                     font_size=18, bold=False, color=CONTENT_TEXT_COLOR,
                     alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
        """Add a text box to slide."""
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = RGBColor(
            int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
        )
        p.font.name = font_name
        p.alignment = alignment
        return txBox

    def add_multiline_text(slide, left, top, width, height, lines_data,
                           font_name="Microsoft YaHei"):
        """Add multiline text box. lines_data: list of (text, font_size, bold, color)"""
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, item in enumerate(lines_data):
            text = item[0]
            fsize = item[1] if len(item) > 1 else 16
            bold = item[2] if len(item) > 2 else False
            color = item[3] if len(item) > 3 else CONTENT_TEXT_COLOR

            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = text
            p.font.size = Pt(fsize)
            p.font.bold = bold
            p.font.color.rgb = RGBColor(
                int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
            )
            p.font.name = font_name
            p.space_after = Pt(6)

        return txBox

    def add_image_safe(slide, img_path, left, top, width=None, height=None):
        """Add image if it exists, with resize."""
        if not os.path.exists(img_path):
            return False
        try:
            if width and height:
                slide.shapes.add_picture(img_path, Inches(left), Inches(top),
                                          Inches(width), Inches(height))
            elif width:
                slide.shapes.add_picture(img_path, Inches(left), Inches(top),
                                          Inches(width))
            elif height:
                slide.shapes.add_picture(img_path, Inches(left), Inches(top),
                                          height=Inches(height))
            else:
                slide.shapes.add_picture(img_path, Inches(left), Inches(top))
            return True
        except Exception as e:
            print(f"  WARNING: Could not add image {img_path}: {e}", file=sys.stderr)
            return False

    def add_accent_line(slide, top=0.15):
        """Add thin accent line at top."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(top),
            Inches(SLIDE_WIDTH), Inches(0.04)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(
            int(ACCENT_COLOR[0:2], 16), int(ACCENT_COLOR[2:4], 16), int(ACCENT_COLOR[4:6], 16)
        )
        shape.line.fill.background()

    def add_bottom_bar(slide, page_num, total):
        """Add page number at bottom right."""
        add_text_box(slide, 11.5, 7.0, 1.5, 0.4,
                     f"{page_num}/{total}", font_size=10,
                     color="999999", alignment=PP_ALIGN.RIGHT)

    # Collect available images
    available_images = []
    if os.path.isdir(images_dir):
        for ext in ["*.png", "*.jpg", "*.jpeg"]:
            available_images.extend(sorted(glob.glob(os.path.join(images_dir, ext))))

    n_pages = 15  # Total number of slides
    page = 0

    # --- Slide 1: Cover ---
    page += 1
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_accent_line(slide, 0.1)

    # Title area
    paper_title = os.path.basename(output_path).replace("_presentation.pptx", "").replace("_", " ")
    add_text_box(slide, 1.0, 1.5, 11.3, 2.0, paper_title,
                 font_size=32, bold=True, color=SLIDE_TITLE_COLOR, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, 1.0, 3.8, 11.3, 0.8, "璁烘枃娣卞害瑙ｈ 路 缁勪細鍒嗕韩",
                 font_size=20, color=SLIDE_SUBTITLE_COLOR, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, 1.0, 5.5, 11.3, 0.5, "Generated by Paper Review Skill | RapidAI Research",
                 font_size=12, color="999999", alignment=PP_ALIGN.CENTER)
    add_bottom_bar(slide, page, n_pages)

    # --- Slide 2: Problem & Motivation ---
    page += 1
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_accent_line(slide)
    add_text_box(slide, 0.6, 0.3, 5, 0.6, "馃挕 鐮旂┒鑳屾櫙涓庡姩鏈?, font_size=26, bold=True, color=SLIDE_TITLE_COLOR)
    add_text_box(slide, 0.6, 1.1, 12, 6.0,
                 review.get("problem", "N/A"),
                 font_size=16, color=CONTENT_TEXT_COLOR)
    add_bottom_bar(slide, page, n_pages)

    # --- Slide 3: Problem Statement ---
    page += 1
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_accent_line(slide)
    add_text_box(slide, 0.6, 0.3, 5, 0.6, "馃幆 瑕佽В鍐崇殑闂", font_size=26, bold=True, color=SLIDE_TITLE_COLOR)

    problem = review.get("problem", "")
    # Split into key points
    points = re.split(r'[.銆?锛?锛焅n]', problem)
    points = [p.strip() for p in points if len(p.strip()) > 10][:6]

    lines = []
    for i, p in enumerate(points, 1):
        lines.append((f"鈻?{p}", 14, False))

    add_multiline_text(slide, 0.6, 1.1, 12, 6.0, lines)
    add_bottom_bar(slide, page, n_pages)

    # --- Slide 4: Key Innovations ---
    page += 1
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_accent_line(slide)
    add_text_box(slide, 0.6, 0.3, 5, 0.6, "鉁?鏍稿績鍒涙柊鐐?, font_size=26, bold=True, color=SLIDE_TITLE_COLOR)

    novelties = review.get("novelty", [])
    lines = []
    for i, n in enumerate(novelties, 1):
        lines.append((f"馃専 鍒涙柊鐐?{i}:", 18, True, ACCENT_COLOR))
        lines.append((n, 14, False))
        lines.append(("", 6, False))

    add_multiline_text(slide, 0.6, 1.1, 12, 6.0, lines)
    add_bottom_bar(slide, page, n_pages)

    # --- Slide 5: Method Overview ---
    page += 1
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_accent_line(slide)
    add_text_box(slide, 0.6, 0.3, 5, 0.6, "馃彈锔?鏂规硶鎬昏", font_size=26, bold=True, color=SLIDE_TITLE_COLOR)
    add_text_box(slide, 0.6, 1.1, 5.5, 3.0,
                 review.get("method_essence", "")[:500],
                 font_size=14, color=CONTENT_TEXT_COLOR)

    # Try to add an architecture image
    added_img = False
    for img in available_images:
        if "page_001" in img or "page_002" in img or "page_003" in img:
            if add_image_safe(slide, img, 6.5, 0.8, width=6.0):
                added_img = True
                break

    if not added_img and available_images:
        add_image_safe(slide, available_images[0], 6.5, 0.8, width=6.0)

    # Method essence continuation
    if len(review.get("method_essence", "")) > 500:
        add_text_box(slide, 0.6, 4.3, 12, 3.0,
                     review.get("method_essence", "")[500:1000],
                     font_size=14, color=CONTENT_TEXT_COLOR)

    add_bottom_bar(slide, page, n_pages)

    # --- Slides 6-7: Method Detail ---
    method_detail = review.get("method_detail", "")
    detail_parts = []
    if len(method_detail) > 1500:
        detail_parts = [method_detail[:len(method_detail)//2], method_detail[len(method_detail)//2:]]
    else:
        detail_parts = [method_detail]

    for part_idx, part_text in enumerate(detail_parts):
        page += 1
        slide = prs.slides.add_slide(blank_layout)
        add_bg(slide)
        add_accent_line(slide)

        sub_title = "馃敩 鏂规硶璇﹁В" if len(detail_parts) == 1 else f"馃敩 鏂规硶璇﹁В ({chr(65+part_idx)})"
        add_text_box(slide, 0.6, 0.3, 5, 0.6, sub_title, font_size=26, bold=True, color=SLIDE_TITLE_COLOR)
        add_text_box(slide, 0.6, 1.1, 12, 6.0, part_text, font_size=14, color=CONTENT_TEXT_COLOR)
        add_bottom_bar(slide, page, n_pages)

    # --- Slide 8: Core Principle ---
    page += 1
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_accent_line(slide)
    add_text_box(slide, 0.6, 0.3, 5, 0.6, "馃 鏍稿績鍘熺悊鍓栨瀽", font_size=26, bold=True, color=SLIDE_TITLE_COLOR)
    add_text_box(slide, 0.6, 1.1, 12, 6.0,
                 review.get("method_essence", "N/A"),
                 font_size=15, color=CONTENT_TEXT_COLOR)
    add_bottom_bar(slide, page, n_pages)

    # --- Slide 9: Experiment Setup ---
    page += 1
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_accent_line(slide)
    add_text_box(slide, 0.6, 0.3, 5, 0.6, "馃搳 瀹為獙璁剧疆", font_size=26, bold=True, color=SLIDE_TITLE_COLOR)

    exp_text = review.get("experiments", "")
    # Extract dataset and setup part
    setup_text = exp_text[:len(exp_text)//2]
    add_text_box(slide, 0.6, 1.1, 12, 6.0, setup_text, font_size=15, color=CONTENT_TEXT_COLOR)
    add_bottom_bar(slide, page, n_pages)

    # --- Slide 10: Experimental Results ---
    page += 1
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_accent_line(slide)
    add_text_box(slide, 0.6, 0.3, 5, 0.6, "馃搱 瀹為獙缁撴灉", font_size=26, bold=True, color=SLIDE_TITLE_COLOR)

    exp_text2 = review.get("experiments", "")
    result_text = exp_text2[len(exp_text2)//2:] if len(exp_text2) > len(exp_text2)//2 else exp_text2
    add_text_box(slide, 0.6, 1.1, 5.5, 6.0, result_text, font_size=14, color=CONTENT_TEXT_COLOR)

    # Try to add a result image
    for img in available_images:
        page_num_match = re.search(r'page_0(4|5|6|7|8)', os.path.basename(img))
        if page_num_match:
            add_image_safe(slide, img, 6.5, 0.8, width=6.0)
            break
    else:
        # Use any available image
        for img in available_images:
            if add_image_safe(slide, img, 6.5, 0.8, width=6.0):
                break

    add_bottom_bar(slide, page, n_pages)

    # --- Slide 11: Ablation Study ---
    page += 1
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_accent_line(slide)
    add_text_box(slide, 0.6, 0.3, 5, 0.6, "馃攳 娑堣瀺瀹為獙鍒嗘瀽", font_size=26, bold=True, color=SLIDE_TITLE_COLOR)

    # Extract ablation-related content
    exp_full = review.get("experiments", "")
    ablation_section = ""
    if "ablation" in exp_full.lower() or "娑堣瀺" in exp_full:
        parts = re.split(r'(?i)(ablation study|娑堣瀺瀹為獙|娑堣瀺鍒嗘瀽)', exp_full)
        if len(parts) > 1:
            ablation_section = " ".join(parts[1:3])

    if not ablation_section:
        ablation_section = exp_full[-len(exp_full)//3:] if len(exp_full) > 500 else exp_full

    add_text_box(slide, 0.6, 1.1, 12, 6.0, ablation_section[:1200],
                 font_size=15, color=CONTENT_TEXT_COLOR)
    add_bottom_bar(slide, page, n_pages)

    # --- Slide 12: Visualization Analysis ---
    page += 1
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_accent_line(slide)
    add_text_box(slide, 0.6, 0.3, 5, 0.6, "馃憗锔?鍙鍖栧垎鏋?, font_size=26, bold=True, color=SLIDE_TITLE_COLOR)

    # Add up to 4 images in a grid
    img_positions = [
        (0.6, 1.0, 5.8, 3.0),
        (6.8, 1.0, 5.8, 3.0),
        (0.6, 4.2, 5.8, 3.0),
        (6.8, 4.2, 5.8, 3.0),
    ]

    img_added_count = 0
    for img in available_images:
        if img_added_count >= 4:
            break
        l, t, w, h = img_positions[img_added_count]
        if add_image_safe(slide, img, l, t, width=w):
            img_added_count += 1

    if img_added_count == 0:
        add_text_box(slide, 0.6, 2.5, 12, 2.0,
                     "(璁烘枃鍥剧墖绱犳潗鏈彁鍙栨垚鍔燂紝寤鸿鎵嬪姩琛ュ厖璁烘枃涓殑鍙鍖栫粨鏋滃浘)",
                     font_size=16, color="999999", alignment=PP_ALIGN.CENTER)

    add_bottom_bar(slide, page, n_pages)

    # --- Slide 13: Discussion ---
    page += 1
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_accent_line(slide)
    add_text_box(slide, 0.6, 0.3, 5, 0.6, "馃挰 璁ㄨ涓庡眬闄愭€?, font_size=26, bold=True, color=SLIDE_TITLE_COLOR)

    discussion = review.get("conclusion", "")
    # Split into discussion and conclusion
    parts = re.split(r'(?i)(limitation|limitations|灞€闄愭€涓嶈冻|future work|鏈潵宸ヤ綔)', discussion)
    if len(parts) > 2:
        limitations = " ".join(parts[1:3]) if len(parts) >= 3 else discussion[len(discussion)//2:]
    else:
        limitations = discussion[len(discussion)//2:] if len(discussion) > 300 else discussion

    add_text_box(slide, 0.6, 1.1, 12, 6.0, limitations[:1000],
                 font_size=15, color=CONTENT_TEXT_COLOR)
    add_bottom_bar(slide, page, n_pages)

    # --- Slide 14: Summary ---
    page += 1
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_accent_line(slide)
    add_text_box(slide, 0.6, 0.3, 5, 0.6, "馃搵 鎬荤粨涓庡惎绀?, font_size=26, bold=True, color=SLIDE_TITLE_COLOR)

    conclusion = review.get("conclusion", "")
    summary = conclusion[:len(conclusion)//2] if len(conclusion) > 300 else conclusion
    add_text_box(slide, 0.6, 1.1, 12, 3.0, summary, font_size=15, color=CONTENT_TEXT_COLOR)

    # Key takeaways box
    add_text_box(slide, 0.6, 4.5, 12, 2.5,
                 "馃挕 **瀵规垜浠殑鍚ず**\n\n鑰冭檻濡備綍灏嗘湰鏂囨柟娉曞簲鐢ㄥ埌鎴戜滑鐨勭爺绌惰棰樹腑銆俓n鍏虫敞鏂规硶鐨勪紭缂虹偣锛屾€濊€冩敼杩涚┖闂淬€?,
                 font_size=14, color=SLIDE_SUBTITLE_COLOR)
    add_bottom_bar(slide, page, n_pages)

    # --- Slide 15: Q&A ---
    page += 1
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_accent_line(slide, 2.5)
    add_text_box(slide, 1.0, 2.5, 11.3, 1.5, "Thank You! 馃檹",
                 font_size=40, bold=True, color=SLIDE_TITLE_COLOR, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1.0, 4.2, 11.3, 1.0, "Q & A",
                 font_size=28, color=SLIDE_SUBTITLE_COLOR, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1.0, 5.8, 11.3, 0.5, "鍙傝€冮摼鎺?& 鑷磋阿",
                 font_size=14, color="999999", alignment=PP_ALIGN.CENTER)
    add_bottom_bar(slide, page, n_pages)

    prs.save(output_path)
    print(f"PPT_SAVED: {output_path}")
    return True


def main():
    parser = argparse.ArgumentParser(description="Generate PPT from paper review")
    parser.add_argument("--review", required=True, help="Review JSON file path or glob")
    parser.add_argument("--images", default="./images", help="Images directory")
    parser.add_argument("--output", default=".", help="Output directory")
    parser.add_argument("--lang", default="zh", help="PPT language (zh/en)")
    args = parser.parse_args()

    json_files = glob.glob(args.review)
    if not json_files and os.path.exists(args.review):
        json_files = [args.review]

    if not json_files:
        print("ERROR: No review JSON files found", file=sys.stderr)
        sys.exit(1)

    os.makedirs(args.output, exist_ok=True)

    for json_path in json_files:
        json_path = os.path.abspath(json_path)
        with open(json_path, "r", encoding="utf-8") as f:
            review = json.load(f)

        base = os.path.splitext(os.path.basename(json_path))[0].replace("_review", "")
        ppt_path = os.path.join(args.output, f"{base}_presentation.pptx")

        print(f"Generating PPT: {ppt_path}")
        create_ppt(review, args.images, ppt_path, args.lang)

        print(f"PPT_PATH: {ppt_path}")

if __name__ == "__main__":
    main()
