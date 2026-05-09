# -*- coding: utf-8 -*-
"""
MaClaw（码卡龙）产品介绍 PPT 生成脚本 — 32 页完整版 v2
面向用户，突出功能特点、核心优势、竞品对比（含 Hermes Agent）
增强：知识库(外脑)、MIS结构化数据、三层记忆、自进化
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ========== 配色 ==========
PRIMARY      = RGBColor(0x0A, 0x0E, 0x27)
ACCENT_BLUE  = RGBColor(0x00, 0x7B, 0xFF)
ACCENT_CYAN  = RGBColor(0x00, 0xD4, 0xFF)
ACCENT_PURPLE= RGBColor(0x8B, 0x5C, 0xF6)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xB0, 0xB8, 0xCC)
MID_GRAY     = RGBColor(0x6B, 0x7B, 0x94)
DARK_CARD    = RGBColor(0x12, 0x18, 0x38)
DARKER_CARD  = RGBColor(0x0E, 0x13, 0x30)
ORANGE       = RGBColor(0xFF, 0x6B, 0x35)
GREEN        = RGBColor(0x00, 0xC8, 0x53)
RED          = RGBColor(0xFF, 0x44, 0x44)
YELLOW       = RGBColor(0xFF, 0xD6, 0x00)
PINK         = RGBColor(0xFF, 0x6B, 0xA0)
TEAL         = RGBColor(0x14, 0xB8, 0xA6)

# ========== 工具函数 ==========
def bg(slide, color=PRIMARY):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def rect(slide, l, t, w, h, fill=None, line=None, st=MSO_SHAPE.ROUNDED_RECTANGLE):
    s = slide.shapes.add_shape(st, l, t, w, h); s.shadow.inherit = False
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else: s.fill.background()
    if line: s.line.color.rgb = line; s.line.width = Pt(1)
    else: s.line.fill.background()
    return s

def tb(slide, l, t, w, h, text, sz=18, c=WHITE, b=False, a=PP_ALIGN.LEFT, fn='Microsoft YaHei'):
    bx = slide.shapes.add_textbox(l, t, w, h); tf = bx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz); p.font.color.rgb = c
    p.font.bold = b; p.font.name = fn; p.alignment = a; return bx

def mp(slide, l, t, w, h, lines, sz=14, c=WHITE, a=PP_ALIGN.LEFT, fn='Microsoft YaHei'):
    bx = slide.shapes.add_textbox(l, t, w, h); tf = bx.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(line, tuple):
            p.text = line[0]; p.font.size = Pt(line[1] if len(line)>1 else sz)
            p.font.color.rgb = line[2] if len(line)>2 else c
            p.font.bold = line[3] if len(line)>3 else False
        else:
            p.text = line; p.font.size = Pt(sz); p.font.color.rgb = c
        p.font.name = fn; p.alignment = a; p.space_after = Pt(sz * 0.3)

def line(slide, l, t, w, c=ACCENT_BLUE, h=Pt(3)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background(); s.shadow.inherit = False
    return s

def card(slide, x, y, w, h, title, items, icon="", tc=ACCENT_CYAN, ac=ACCENT_BLUE):
    rect(slide, x, y, w, h, fill=DARK_CARD, line=RGBColor(0x1E,0x2A,0x50))
    if icon: tb(slide, x+Inches(0.2), y+Inches(0.12), w-Inches(0.4), Inches(0.35), icon, sz=18, c=ac)
    tb(slide, x+Inches(0.2), y+Inches(0.45), w-Inches(0.4), Inches(0.35), title, sz=15, c=tc, b=True)
    line(slide, x+Inches(0.2), y+Inches(0.82), Inches(0.6), ac, Pt(2))
    mp(slide, x+Inches(0.2), y+Inches(0.95), w-Inches(0.4), h-Inches(1.1),
       [f"•  {it}" for it in items], sz=10, c=LIGHT_GRAY)

def section_header(slide, label, title, subtitle=""):
    tb(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.4), label, sz=11, c=ACCENT_CYAN, b=True)
    tb(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.7), title, sz=34, c=WHITE, b=True)
    line(slide, Inches(0.8), Inches(1.5), Inches(2.2), ACCENT_BLUE)
    if subtitle:
        tb(slide, Inches(0.8), Inches(1.7), Inches(11), Inches(0.4), subtitle, sz=14, c=LIGHT_GRAY)

def new_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); return s

def table_cell(slide, x, y, w, h, text, sz=10, c=WHITE, b=False, fill=None, a=PP_ALIGN.CENTER):
    rect(slide, x, y, w, h, fill=fill, line=RGBColor(0x1E,0x2A,0x50), st=MSO_SHAPE.RECTANGLE)
    tb(slide, x+Inches(0.06), y+Inches(0.04), w-Inches(0.12), h-Inches(0.08), text, sz=sz, c=c, b=b, a=a)

print("工具函数准备完成")

# ================================================================
# P1: 封面
# ================================================================
s = new_slide()
bg(s, RGBColor(0x06,0x09,0x1C))
line(s, Inches(0), Inches(0), Inches(13.333), ACCENT_BLUE, Pt(4))
tb(s, Inches(1), Inches(1.0), Inches(11), Inches(0.7), "MaClaw  码卡龙", sz=58, c=WHITE, b=True, a=PP_ALIGN.CENTER)
line(s, Inches(4.8), Inches(1.95), Inches(3.7), ACCENT_CYAN, Pt(3))
tb(s, Inches(1.5), Inches(2.4), Inches(10), Inches(0.8), "通用可自进化智能体平台", sz=38, c=ACCENT_CYAN, a=PP_ALIGN.CENTER)
tb(s, Inches(2), Inches(3.4), Inches(9), Inches(0.6), "你的个人数智工作伙伴 — 不只是聊天，而是替你干活", sz=22, c=LIGHT_GRAY, a=PP_ALIGN.CENTER)
tags = ["长期记忆", "知识库(外脑)", "自进化", "SSH远程", "GUI自动化", "MIS数据", "19种工作流"]
tw = Inches(1.55); th = Inches(0.45); total = len(tags)*tw.inches + (len(tags)-1)*0.18
sx = Inches((13.333-total)/2)
for i, tag in enumerate(tags):
    x = sx + Inches(i*(tw.inches+0.18)); y = Inches(4.6)
    rect(s, x, y, tw, th, fill=DARK_CARD, line=ACCENT_BLUE)
    tb(s, x, y+Inches(0.06), tw, th-Inches(0.12), tag, sz=12, c=ACCENT_CYAN, a=PP_ALIGN.CENTER)
tb(s, Inches(0), Inches(6.5), Inches(13.333), Inches(0.4), "github.com/RapidAI/MaClaw  ·  maclaw.top", sz=14, c=MID_GRAY, a=PP_ALIGN.CENTER)
tb(s, Inches(0), Inches(6.9), Inches(13.333), Inches(0.3), "© 2025 Dr. Daniel", sz=11, c=MID_GRAY, a=PP_ALIGN.CENTER)
print("P1 封面 完成")

# ================================================================
# P2: 目录
# ================================================================
s = new_slide()
section_header(s, "CONTENTS", "目录")
toc = [
    ("01", "产品定位与痛点", "AI 助手的困境与 MaClaw 的回答", ACCENT_CYAN),
    ("02", "核心能力总览", "六大能力支柱 · 全链路架构", ACCENT_BLUE),
    ("03", "长期记忆系统", "内脑 · 中脑 · 外脑 · 三层架构", ACCENT_PURPLE),
    ("04", "知识库（外脑）", "12种来源 · 蒸馏管线 · 知识图谱", ORANGE),
    ("05", "MIS 结构化数据", "企业级数据管理 · 连接器 · 治理", GREEN),
    ("06", "执行能力", "SSH · 浏览器 · GUI · 软件开发", TEAL),
    ("07", "自进化与技能系统", "经验积累 · 能力缺口 · 越用越强", YELLOW),
    ("08", "竞品对比", "MaClaw vs Hermes/ChatGPT/Copilot", RED),
    ("09", "应用场景与快速上手", "6大行业 · 4步开始", PINK),
]
for i, (num, title, desc, color) in enumerate(toc):
    y = Inches(1.9)+i*Inches(0.6)
    rect(s, Inches(1.0), y, Inches(0.55), Inches(0.45), fill=color, st=MSO_SHAPE.RECTANGLE)
    tb(s, Inches(1.0), y+Inches(0.06), Inches(0.55), Inches(0.35), num, sz=13, c=WHITE, b=True, a=PP_ALIGN.CENTER)
    tb(s, Inches(1.8), y+Inches(0.02), Inches(3.0), Inches(0.25), title, sz=14, c=WHITE, b=True)
    tb(s, Inches(1.8), y+Inches(0.26), Inches(8.0), Inches(0.2), desc, sz=10, c=MID_GRAY)
    if i < len(toc)-1:
        line(s, Inches(1.0), y+Inches(0.55), Inches(11.3), RGBColor(0x1E,0x2A,0x50), Pt(1))
print("P2 目录 完成")

# ================================================================
# P3: 产品定位
# ================================================================
s = new_slide()
section_header(s, "POSITIONING", "MaClaw 是什么？",
               "一个通用可自进化智能体平台 —— 你的个人数智工作伙伴")
rect(s, Inches(0.8), Inches(2.2), Inches(11.7), Inches(1.5), fill=DARK_CARD, line=ACCENT_BLUE)
mp(s, Inches(1.2), Inches(2.35), Inches(11.0), Inches(1.2), [
    ("不只是聊天机器人，不是 IDE 插件，也不是简单的 API Wrapper", 18, ACCENT_CYAN, True),
    ("", 6),
    ("MaClaw 是一个能记住你、学习你、替你干活的自主智能体平台。", 16, WHITE, False),
    ("基于 Wails + Go + React 构建，集结构化工作流、长期记忆、知识库、技能扩展、多通道协作于一体。", 13, LIGHT_GRAY, False),
])
features = [
    ("🧠", "能记", "内中外脑\n三层记忆", ACCENT_PURPLE),
    ("📚", "能学", "知识库+经验\n持续积累", ACCENT_BLUE),
    ("🛠️", "能做", "19种工作流\nSSH/GUI/浏览器", GREEN),
    ("🧬", "能进化", "能力缺口检测\n技能自修复", ORANGE),
    ("🏢", "企业级", "MIS结构化数据\n审计合规", TEAL),
]
fw = Inches(2.1); fh = Inches(2.5)
for i, (icon, title, desc, color) in enumerate(features):
    x = Inches(0.8)+i*(fw+Inches(0.15)); y = Inches(4.2)
    rect(s, x, y, fw, fh, fill=DARK_CARD, line=color)
    tb(s, x, y+Inches(0.15), fw, Inches(0.4), icon, sz=26, c=color, a=PP_ALIGN.CENTER)
    tb(s, x, y+Inches(0.55), fw, Inches(0.3), title, sz=18, c=color, b=True, a=PP_ALIGN.CENTER)
    line(s, x+Inches(0.3), y+Inches(0.9), fw-Inches(0.6), color, Pt(2))
    tb(s, x+Inches(0.2), y+Inches(1.05), fw-Inches(0.4), Inches(1.2), desc, sz=11, c=LIGHT_GRAY, a=PP_ALIGN.CENTER)
print("P3 产品定位 完成")

# ================================================================
# P4: 痛点分析
# ================================================================
s = new_slide()
section_header(s, "PAIN POINTS", "现有 AI 助手的六大不足")
pains = [
    ("🧠", "记忆缺失", "每次对话都是陌生人\n重复解释偏好、上下文", RED),
    ("📚", "知识孤岛", "无法沉淀和积累知识\n同样的错误反复犯", ORANGE),
    ("🔧", "只会说不会做", "给出建议但不能执行\n你还得自己动手操作", YELLOW),
    ("📱", "入口单一", "只能在浏览器或 IDE 使用\n无法随时随地访问", ACCENT_PURPLE),
    ("🔄", "无法进化", "能力固定不变\n用一年和用一天没有区别", MID_GRAY),
    ("🔒", "数据安全隐患", "数据上传云端\n企业敏感信息泄露风险", PINK),
]
pw = Inches(3.65); ph = Inches(1.9)
for i, (icon, title, desc, color) in enumerate(pains):
    col = i%3; row = i//3
    x = Inches(0.8)+col*(pw+Inches(0.2)); y = Inches(2.2)+row*(ph+Inches(0.2))
    rect(s, x, y, pw, ph, fill=DARK_CARD, line=color)
    tb(s, x+Inches(0.2), y+Inches(0.1), pw-Inches(0.4), Inches(0.35), icon+"  "+title, sz=16, c=color, b=True)
    line(s, x+Inches(0.2), y+Inches(0.5), Inches(0.8), color, Pt(2))
    tb(s, x+Inches(0.2), y+Inches(0.65), pw-Inches(0.4), Inches(1.1), desc, sz=12, c=LIGHT_GRAY)
print("P4 痛点 完成")

# ================================================================
# P5: 设计哲学
# ================================================================
s = new_slide()
section_header(s, "DESIGN PHILOSOPHY", "三大核心理念")
philosophies = [
    ("自主而非被动", ACCENT_BLUE, [
        "不是 Copilot 辅助你操作，而是 Agent 替你执行",
        "从需求理解到成果交付，全链路自主推进",
        "19 种结构化工作流，每种都遵循质量闭环",
        "你说想法，它出成果",
    ]),
    ("积累而非遗忘", ACCENT_PURPLE, [
        "三层记忆架构：内脑(对话) + 中脑(长期) + 外脑(知识库)",
        "7 种记忆类型覆盖身份、偏好、经验、知识",
        "Ebbinghaus 遗忘曲线 + 语义图谱 + 冲突检测",
        "用得越久，它越懂你",
    ]),
    ("进化而非固定", ORANGE, [
        "能力缺口检测 → 自动搜索 SkillHub → 一键安装",
        "Combee 经验蒸馏 → 保守学习 → 权重自适应",
        "技能自修复 → 错误分析 → 自动修补",
        "Nudge 推荐系统 → 重复操作自动封装为技能",
    ]),
]
cw = Inches(3.7); ch = Inches(4.5)
for i, (title, color, items) in enumerate(philosophies):
    x = Inches(0.8)+i*(cw+Inches(0.2)); y = Inches(2.0)
    rect(s, x, y, cw, ch, fill=DARK_CARD, line=color)
    tb(s, x+Inches(0.2), y+Inches(0.15), cw-Inches(0.4), Inches(0.35), title, sz=18, c=color, b=True)
    line(s, x+Inches(0.2), y+Inches(0.55), Inches(1.2), color, Pt(2))
    mp(s, x+Inches(0.2), y+Inches(0.75), cw-Inches(0.4), ch-Inches(1.0),
       [f"•  {it}" for it in items], sz=11, c=LIGHT_GRAY)
print("P5 设计哲学 完成")

# ================================================================
# P6: 核心能力总览
# ================================================================
s = new_slide()
section_header(s, "CORE CAPABILITIES", "六大核心能力支柱")
caps = [
    ("🧠", "三层记忆系统", ACCENT_PURPLE,
     ["对话记忆 · 工作记忆","长期记忆 · 语义图谱","Ebbinghaus衰减 + 冲突检测","7种记忆类型覆盖全场景"]),
    ("📚", "知识库（外脑）", ACCENT_BLUE,
     ["12种来源导入","Card/Fact/Node三层产出","蒸馏Pipeline · 知识图谱","质量管控 · 敏感数据过滤"]),
    ("🛠️", "197+ 动态工具", GREEN,
     ["SSH远程 · SFTP文件传输","浏览器CDP自动化","桌面GUI双引擎","渐进暴露 · 智能路由"]),
    ("🧬", "自进化引擎", ORANGE,
     ["Combee经验蒸馏","能力缺口检测 · 技能自修复","Nudge推荐 · craft_tool转化","工具路由自适应"]),
    ("🏗️", "MIS 结构化数据", TEAL,
     ["多租户 · 多工作域","Schema渐进演化 · 连接器","业务规则引擎 · 审批流","审计日志 · 治理证据包"]),
    ("🌐", "多通道协作", ACCENT_CYAN,
     ["桌面GUI · 终端TUI","微信/飞书/QQ/Telegram","REST API 多租户","AgentNet P2P网络"]),
]
cw = Inches(3.8); ch = Inches(2.0)
for i, (icon, title, color, items) in enumerate(caps):
    col = i%3; row = i//3
    x = Inches(0.8)+col*(cw+Inches(0.25)); y = Inches(2.0)+row*(ch+Inches(0.2))
    card(s, x, y, cw, ch, title, items, icon, tc=color, ac=color)
print("P6 核心能力 完成")

# ================================================================
# P7: 能力全景图
# ================================================================
s = new_slide()
section_header(s, "CAPABILITY MAP", "能力全景——从感知到执行的全链路")
stages = [
    ("感知", "意图理解\n三层UIC", ACCENT_CYAN),
    ("记忆", "内中外脑\n长期记忆", ACCENT_BLUE),
    ("规划", "工作流引擎\n19种模板", ACCENT_PURPLE),
    ("执行", "SSH/GUI\n工具调用", GREEN),
    ("进化", "经验蒸馏\n技能沉淀", ORANGE),
]
sw = Inches(2.1); sh = Inches(1.2)
total = len(stages)*sw.inches + (len(stages)-1)*0.3
sx = Inches((13.333-total)/2)
for i, (name, desc, color) in enumerate(stages):
    x = sx + Inches(i*(sw.inches+0.3)); y = Inches(2.2)
    rect(s, x, y, sw, sh, fill=DARK_CARD, line=color)
    tb(s, x, y+Inches(0.05), sw, Inches(0.35), name, sz=16, c=color, b=True, a=PP_ALIGN.CENTER)
    line(s, x+Inches(0.3), y+Inches(0.42), sw-Inches(0.6), color, Pt(2))
    tb(s, x+Inches(0.1), y+Inches(0.5), sw-Inches(0.2), Inches(0.6), desc, sz=11, c=LIGHT_GRAY, a=PP_ALIGN.CENTER)
    if i < len(stages)-1:
        tb(s, x+sw, y+Inches(0.3), Inches(0.3), Inches(0.4), "→", sz=18, c=ACCENT_BLUE, a=PP_ALIGN.CENTER)
details = [
    ("感知层", "关键词规则(<1ms) → BM25语义检索(<5ms) → LLM确认(10-30s)  渐进式意图识别", ACCENT_CYAN),
    ("记忆层", "对话记忆 → 工作记忆 → 长期记忆 → 知识库(外脑)  三维度存储 + 语义图谱 + 时序树", ACCENT_BLUE),
    ("规划层", "结构化工作流 · 任务分解 · TDD验收  需求→设计→编码→测试→集成→验收", ACCENT_PURPLE),
    ("执行层", "197+动态工具 · SSH远程 · 浏览器CDP · 桌面GUI  工具按上下文按需激活", GREEN),
    ("进化层", "Combee经验聚合 · 能力缺口检测 · 技能自修复 · Nudge  保守学习原则", ORANGE),
]
for i, (title, desc, color) in enumerate(details):
    y = Inches(3.8) + i*Inches(0.72)
    tb(s, Inches(1.0), y, Inches(1.2), Inches(0.35), title, sz=12, c=color, b=True)
    tb(s, Inches(2.3), y, Inches(10.0), Inches(0.6), desc, sz=10, c=LIGHT_GRAY)
    if i < len(details)-1:
        line(s, Inches(1.0), y+Inches(0.62), Inches(11.3), RGBColor(0x1E,0x2A,0x50), Pt(1))
print("P7 能力全景 完成")

# ================================================================
# P8: 核心差异化
# ================================================================
s = new_slide()
section_header(s, "DIFFERENTIATION", "MaClaw 的核心差异化——一句话说清楚")
rect(s, Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.8), fill=DARK_CARD, line=ACCENT_BLUE)
mp(s, Inches(2.0), Inches(2.15), Inches(9.3), Inches(1.5), [
    ("能记  ·  能学  ·  能做  ·  能进化", 34, ACCENT_CYAN, True),
    ("", 6),
    ("当其他 AI 助手还在「回答问题」时，MaClaw 已经在「替你干活」了", 16, LIGHT_GRAY, False),
])
diffs = [
    ("能记", "三层记忆 + 7种类型", "对话结束后仍然记得你的偏好、项目背景、历史决策\n不会每次都从零开始", ACCENT_PURPLE),
    ("能学", "知识库 + 经验蒸馏", "12种来源导入知识，Combee蒸馏成功经验\n用得越久，知识库越丰富", ACCENT_BLUE),
    ("能做", "197+工具 + SSH/GUI", "不只是给建议，而是直接SSH到服务器操作\n浏览器自动化 + 桌面GUI自动化", GREEN),
    ("能进化", "技能自修复 + 能力缺口", "发现能力不足自动搜索安装技能\n成功经验自动沉淀为可复用能力", ORANGE),
]
dw = Inches(5.7); dh = Inches(2.0)
for i, (title, subtitle, desc, color) in enumerate(diffs):
    col = i%2; row = i//2
    x = Inches(0.8)+col*(dw+Inches(0.3)); y = Inches(4.2)+row*(dh+Inches(0.2))
    rect(s, x, y, dw, dh, fill=DARK_CARD, line=color)
    tb(s, x+Inches(0.2), y+Inches(0.1), Inches(1.0), Inches(0.35), title, sz=20, c=color, b=True)
    tb(s, x+Inches(1.3), y+Inches(0.15), Inches(4.0), Inches(0.25), subtitle, sz=12, c=YELLOW)
    line(s, x+Inches(0.2), y+Inches(0.55), Inches(0.8), color, Pt(2))
    tb(s, x+Inches(0.2), y+Inches(0.7), dw-Inches(0.4), Inches(1.1), desc, sz=11, c=LIGHT_GRAY)
print("P8 差异化 完成")
