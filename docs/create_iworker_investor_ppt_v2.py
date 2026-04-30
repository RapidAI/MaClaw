# -*- coding: utf-8 -*-
"""
iWorker 投资人宣传 PPT v2 — 面向投资人的升级版
基于最新白皮书、客户宣传稿、executive architecture/skills 文档
增加：市场规模、竞争壁垒、商业模式、融资计划、投资亮点
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ========== 配色方案 ==========
PRIMARY = RGBColor(0x08, 0x0C, 0x22)       # 极深蓝黑
ACCENT_BLUE = RGBColor(0x00, 0x7B, 0xFF)
ACCENT_CYAN = RGBColor(0x00, 0xD4, 0xFF)
ACCENT_PURPLE = RGBColor(0x7C, 0x3A, 0xFF) # 紫色点缀
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xB8, 0xCC)
MID_GRAY = RGBColor(0x6B, 0x7B, 0x94)
DARK_CARD = RGBColor(0x0F, 0x14, 0x32)
CARD_BORDER = RGBColor(0x1A, 0x25, 0x4A)
ORANGE = RGBColor(0xFF, 0x6B, 0x35)
GREEN = RGBColor(0x00, 0xC8, 0x53)
RED = RGBColor(0xFF, 0x3B, 0x5C)
GOLD = RGBColor(0xFF, 0xC1, 0x07)

def set_slide_bg(slide, color=PRIMARY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.shadow.inherit = False
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline_text(slide, left, top, width, height, lines, font_size=16, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei', line_spacing=1.3):
    """多行文本，lines 是字符串列表"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_text in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line_text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(font_size * 0.3)
    return txBox

def add_decorative_line(slide, left, top, width, color=ACCENT_BLUE, height=Pt(3)):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    line.shadow.inherit = False
    return line

def add_icon_circle(slide, left, top, size, color=ACCENT_BLUE):
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    circle.shadow.inherit = False
    return circle

def add_gradient_bar(slide, left, top, width, height=Pt(4)):
    """渐变装饰条（模拟：蓝→紫→青）"""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_BLUE
    bar.line.fill.background()
    bar.shadow.inherit = False
    return bar

def add_section_header(slide, label, title, y_start=0.4):
    """标准化章节头部"""
    add_text_box(slide, Inches(0.8), Inches(y_start), Inches(10), Inches(0.4),
                 label, font_size=11, color=ACCENT_CYAN, bold=True)
    add_text_box(slide, Inches(0.8), Inches(y_start + 0.35), Inches(10), Inches(0.8),
                 title, font_size=36, color=WHITE, bold=True)
    add_gradient_bar(slide, Inches(0.8), Inches(y_start + 1.1), Inches(2.5))

def add_card(slide, x, y, w, h, fill=DARK_CARD, border=CARD_BORDER):
    """标准卡片"""
    card = add_shape(slide, x, y, w, h, fill_color=fill, line_color=border)
    return card

print("工具函数加载完成")
