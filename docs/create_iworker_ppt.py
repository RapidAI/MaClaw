# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 创建演示文稿 (16:9)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ========== 配色方案 ==========
PRIMARY = RGBColor(0x0A, 0x0E, 0x27)       # 深蓝黑背景
ACCENT_BLUE = RGBColor(0x00, 0x7B, 0xFF)   # 亮蓝
ACCENT_CYAN = RGBColor(0x00, 0xD4, 0xFF)   # 青色
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xB8, 0xCC)
MID_GRAY = RGBColor(0x6B, 0x7B, 0x94)
DARK_CARD = RGBColor(0x12, 0x18, 0x38)     # 卡片背景
GRADIENT_END = RGBColor(0x1A, 0x1A, 0x2E)
ORANGE = RGBColor(0xFF, 0x6B, 0x35)
GREEN = RGBColor(0x00, 0xC8, 0x53)

def set_slide_bg(slide, color=PRIMARY):
    """设置幻灯片背景色"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    """添加形状"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.shadow.inherit = False
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_decorative_line(slide, left, top, width, color=ACCENT_BLUE, height=Pt(3)):
    """添加装饰线"""
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    line.shadow.inherit = False
    return line

def add_icon_circle(slide, left, top, size, color=ACCENT_BLUE):
    """添加圆形图标占位"""
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    circle.shadow.inherit = False
    return circle

# ========== 第1页：封面 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
set_slide_bg(slide)

# 装饰元素 - 右上角光晕
add_shape(slide, Inches(8), Inches(-1), Inches(6), Inches(6), 
          fill_color=RGBColor(0x00, 0x2A, 0x5E), shape_type=MSO_SHAPE.OVAL)

# 左侧蓝色装饰条
add_shape(slide, Inches(0.6), Inches(2.2), Pt(4), Inches(2.5), fill_color=ACCENT_BLUE)

# 主标题
add_text_box(slide, Inches(1.0), Inches(1.8), Inches(7), Inches(1.2),
             "iWorker", font_size=72, color=ACCENT_CYAN, bold=True)

# 副标题
add_text_box(slide, Inches(1.0), Inches(3.0), Inches(8), Inches(1.0),
             "AI Native 企业组织能力系统", font_size=36, color=WHITE, bold=True)

# Slogan
add_text_box(slide, Inches(1.0), Inches(4.2), Inches(8), Inches(0.6),
             "让企业能力沉淀在系统中，而不是绑定在少数人身上", 
             font_size=20, color=LIGHT_GRAY, bold=False)

# 装饰线
add_decorative_line(slide, Inches(1.0), Inches(5.0), Inches(3), ACCENT_BLUE)

# 底部信息
add_text_box(slide, Inches(1.0), Inches(6.2), Inches(5), Inches(0.5),
             "企业AI转型 · 从工具到组织能力的跨越", font_size=14, color=MID_GRAY)

# 右下角大字
add_text_box(slide, Inches(9.5), Inches(5.5), Inches(3.5), Inches(1.5),
             "AI\nNative", font_size=48, color=RGBColor(0x15, 0x20, 0x45), bold=True, alignment=PP_ALIGN.RIGHT)

print("封面完成")

# ========== 第2页：行业痛点 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

# 顶部标题
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
             "PROBLEM", font_size=12, color=ACCENT_CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.8),
             "企业AI应用的现实困境", font_size=36, color=WHITE, bold=True)
add_decorative_line(slide, Inches(0.8), Inches(1.55), Inches(2), ACCENT_BLUE)

# 5个痛点卡片
pain_points = [
    ("01", "能力无法沉淀", "有人会用，有人不会用\n有场景能提效，有场景无法接入\n结果惊艳但难以复用"),
    ("02", "安全不可控", "员工各自使用外部工具\n企业数据外流风险高\n缺乏统一的合规审计"),
    ("03", "成果难复制", "做了不少试点\n始终进不了核心流程\n试点成果无法规模推广"),
    ("04", "成本不可控", "各部门各自订阅AI工具\n重复采购，费用失控\n缺乏统一管理策略"),
    ("05", "流程不闭环", "AI只停留在问答和生成\n接不上审批、归档、追踪\n最终还是回到人工处理"),
]

card_width = Inches(2.2)
card_height = Inches(4.5)
start_x = Inches(0.6)
gap = Inches(0.25)

for i, (num, title, desc) in enumerate(pain_points):
    x = start_x + i * (card_width + gap)
    y = Inches(2.2)
    
    # 卡片背景
    card = add_shape(slide, x, y, card_width, card_height, fill_color=DARK_CARD)
    
    # 序号
    add_text_box(slide, x + Inches(0.25), y + Inches(0.3), Inches(1), Inches(0.6),
                 num, font_size=32, color=ACCENT_BLUE, bold=True)
    
    # 分隔线
    add_decorative_line(slide, x + Inches(0.25), y + Inches(1.0), Inches(0.8), ACCENT_CYAN, Pt(2))
    
    # 标题
    add_text_box(slide, x + Inches(0.25), y + Inches(1.2), card_width - Inches(0.5), Inches(0.6),
                 title, font_size=18, color=WHITE, bold=True)
    
    # 描述
    tb = add_text_box(slide, x + Inches(0.25), y + Inches(2.0), card_width - Inches(0.5), Inches(2.2),
                 desc, font_size=12, color=LIGHT_GRAY)

print("痛点页完成")

# ========== 第3页：核心理念 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
             "INSIGHT", font_size=12, color=ACCENT_CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.8),
             "企业需要的不是更多AI工具，而是组织能力系统", font_size=32, color=WHITE, bold=True)
add_decorative_line(slide, Inches(0.8), Inches(1.55), Inches(2), ACCENT_BLUE)

# 左侧 - 旧范式
left_card = add_shape(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.5), fill_color=DARK_CARD)
add_text_box(slide, Inches(1.2), Inches(2.5), Inches(4.8), Inches(0.5),
             "❌  过去：信息化", font_size=20, color=ORANGE, bold=True)
add_text_box(slide, Inches(1.2), Inches(3.2), Inches(4.8), Inches(0.5),
             "解决「信息能不能流动」的问题", font_size=16, color=LIGHT_GRAY)

old_items = ["场景经验", "岗位知识", "处理流程", "协作方式", "输出标准"]
for i, item in enumerate(old_items):
    add_text_box(slide, Inches(1.2), Inches(4.0 + i * 0.45), Inches(4.8), Inches(0.4),
                 f"⚙  {item}  —  绑定在个人身上，无法复制", font_size=13, color=MID_GRAY)

# 右侧 - 新范式
right_card = add_shape(slide, Inches(6.8), Inches(2.2), Inches(5.7), Inches(4.5), fill_color=DARK_CARD)
add_text_box(slide, Inches(7.2), Inches(2.5), Inches(5.0), Inches(0.5),
             "✅  现在：AI组织能力系统", font_size=20, color=GREEN, bold=True)
add_text_box(slide, Inches(7.2), Inches(3.2), Inches(5.0), Inches(0.5),
             "解决「能力能不能复制」的问题", font_size=16, color=LIGHT_GRAY)

new_items = ["数字员工", "业务流程", "组织记忆", "权限安全", "AI基础设施"]
for i, item in enumerate(new_items):
    add_text_box(slide, Inches(7.2), Inches(4.0 + i * 0.45), Inches(5.0), Inches(0.4),
                 f"◆  {item}  —  沉淀在系统中，可复制可扩展", font_size=13, color=ACCENT_CYAN)

print("理念页完成")

# ========== 第4页：iWorker定位 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
             "PRODUCT", font_size=12, color=ACCENT_CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
             "iWorker — AI Native 企业组织能力系统", font_size=36, color=WHITE, bold=True)
add_decorative_line(slide, Inches(0.8), Inches(1.55), Inches(2), ACCENT_BLUE)

# 核心定位描述
add_text_box(slide, Inches(0.8), Inches(2.0), Inches(11), Inches(0.8),
             "帮助企业把关键能力沉淀进系统，构建可落地、可治理、可复制的 AI 生产力体系",
             font_size=18, color=LIGHT_GRAY)

# 5大能力模块 - 横排卡片
modules = [
    ("👥", "数字员工", "可分工协作的\nAI数字同事团队"),
    ("🔄", "业务流程", "可编排运行的\n自动化流程引擎"),
    ("🧠", "组织记忆", "可共享复用的\n企业知识资产"),
    ("🔒", "权限安全", "可审计追溯的\n治理与安全体系"),
    ("🏗️", "AI基础设施", "可持续扩展的\n模型与算力管理"),
]

mod_width = Inches(2.2)
mod_height = Inches(3.2)
mod_start_x = Inches(0.6)
mod_gap = Inches(0.25)

for i, (icon, title, desc) in enumerate(modules):
    x = mod_start_x + i * (mod_width + mod_gap)
    y = Inches(3.2)
    
    card = add_shape(slide, x, y, mod_width, mod_height, fill_color=DARK_CARD)
    
    # 图标圆
    add_icon_circle(slide, x + Inches(0.7), y + Inches(0.3), Inches(0.7), 
                    RGBColor(0x00, 0x1F, 0x45))
    add_text_box(slide, x + Inches(0.7), y + Inches(0.35), Inches(0.7), Inches(0.6),
                 icon, font_size=24, color=WHITE, alignment=PP_ALIGN.CENTER)
    
    # 标题
    add_text_box(slide, x + Inches(0.2), y + Inches(1.2), mod_width - Inches(0.4), Inches(0.5),
                 title, font_size=18, color=ACCENT_CYAN, bold=True, alignment=PP_ALIGN.CENTER)
    
    # 分隔线
    add_decorative_line(slide, x + Inches(0.6), y + Inches(1.75), Inches(1), ACCENT_BLUE, Pt(2))
    
    # 描述
    add_text_box(slide, x + Inches(0.2), y + Inches(2.0), mod_width - Inches(0.4), Inches(1.0),
                 desc, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

print("产品定位页完成")

# ========== 第5页：系统架构 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
             "ARCHITECTURE", font_size=12, color=ACCENT_CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.8),
             "iWorker 核心能力架构", font_size=36, color=WHITE, bold=True)
add_decorative_line(slide, Inches(0.8), Inches(1.55), Inches(2), ACCENT_BLUE)

# 架构层级 - 从上到下
layers = [
    ("用户交互层", "对话式工作台 · 多端接入 · 企业IM集成 · 自然语言交互", ACCENT_CYAN),
    ("数字员工层", "iWorkerCenter 组织中枢 · 角色定义 · 能力配置 · 任务分派 · 协作网络", ACCENT_BLUE),
    ("流程编排层", "任务流程 · 审批流转 · 条件路由 · 结果归档 · 闭环追踪", RGBColor(0x00, 0x96, 0xC7)),
    ("组织记忆层", "共享记忆 · 业务知识 · 模板库 · 案例库 · 按需检索复用", RGBColor(0x00, 0x5F, 0x99)),
    ("AI基础设施层", "模型路由 · 能力包管理 · 成本策略 · 安全审计 · 混合部署", RGBColor(0x00, 0x3D, 0x66)),
]

for i, (title, desc, color) in enumerate(layers):
    y = Inches(2.2 + i * 1.0)
    w = Inches(11.5 - i * 0.3)
    x = Inches(0.8 + i * 0.15)
    
    layer_shape = add_shape(slide, x, y, w, Inches(0.8), fill_color=color)
    
    add_text_box(slide, x + Inches(0.3), y + Inches(0.05), Inches(2.2), Inches(0.35),
                 title, font_size=15, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.38), w - Inches(0.6), Inches(0.35),
                 desc, font_size=11, color=RGBColor(0xCC, 0xDD, 0xEE))

print("架构页完成")

# ========== 第6页：五大解决方案 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
             "SOLUTION", font_size=12, color=ACCENT_CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.8),
             "用系统化思维替代零散试点", font_size=36, color=WHITE, bold=True)
add_decorative_line(slide, Inches(0.8), Inches(1.55), Inches(2), ACCENT_BLUE)

solutions = [
    ("数字员工", "替代零散工具", "把常用岗位能力封装为\n数字同事，统一管理", "👤"),
    ("组织中枢", "替代分散试用", "所有角色、能力、流程、\n记忆统一进入系统管理", "🏢"),
    ("共享记忆", "替代重复摸索", "模板、案例、知识、经验\n按范围沉淀复用", "💾"),
    ("统一路由", "替代各自直连", "统一管理模型供应商、\n调用规则与成本策略", "🔀"),
    ("流程编排", "替代孤立对话", "AI进入任务、审批、\n归档等真实业务环节", "📋"),
]

sol_width = Inches(2.2)
sol_height = Inches(4.2)

for i, (title, sub, desc, icon) in enumerate(solutions):
    x = Inches(0.6 + i * 2.45)
    y = Inches(2.3)
    
    card = add_shape(slide, x, y, sol_width, sol_height, fill_color=DARK_CARD)
    
    # 图标
    add_icon_circle(slide, x + Inches(0.7), y + Inches(0.3), Inches(0.7), RGBColor(0x00, 0x1F, 0x45))
    add_text_box(slide, x + Inches(0.7), y + Inches(0.35), Inches(0.7), Inches(0.6),
                 icon, font_size=22, color=WHITE, alignment=PP_ALIGN.CENTER)
    
    # 标题
    add_text_box(slide, x + Inches(0.2), y + Inches(1.2), sol_width - Inches(0.4), Inches(0.4),
                 title, font_size=18, color=ACCENT_CYAN, bold=True, alignment=PP_ALIGN.CENTER)
    
    # 替代说明
    add_text_box(slide, x + Inches(0.2), y + Inches(1.65), sol_width - Inches(0.4), Inches(0.4),
                 sub, font_size=13, color=ORANGE, alignment=PP_ALIGN.CENTER)
    
    add_decorative_line(slide, x + Inches(0.5), y + Inches(2.1), Inches(1.2), ACCENT_BLUE, Pt(1))
    
    # 描述
    add_text_box(slide, x + Inches(0.2), y + Inches(2.3), sol_width - Inches(0.4), Inches(1.5),
                 desc, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

print("解决方案页完成")

# ========== 第7页：核心价值主张 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
             "VALUE PROPOSITION", font_size=12, color=ACCENT_CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.8),
             "核心价值：能力沉淀 > 工具堆砌", font_size=36, color=WHITE, bold=True)
add_decorative_line(slide, Inches(0.8), Inches(1.55), Inches(2), ACCENT_BLUE)

# 中间大标语
big_quote = add_shape(slide, Inches(1.5), Inches(2.2), Inches(10.3), Inches(1.5), fill_color=DARK_CARD)
add_text_box(slide, Inches(2.0), Inches(2.35), Inches(9.3), Inches(1.2),
             "\"人员流动不应中断组织基本运行\n企业缺少某个具体人，不等于企业缺少这项能力\"",
             font_size=24, color=ACCENT_CYAN, bold=True, alignment=PP_ALIGN.CENTER)

# 四个价值卡片
values = [
    ("可落地", "从试点到真实业务流程\n数字员工直接参与日常工作\n开箱即用的场景能力包"),
    ("可治理", "统一权限管理与审计\n数据安全与合规保障\n清晰的成本可控策略"),
    ("可复制", "优秀经验沉淀为组织资产\n能力模板快速复制推广\n新业务场景快速上线"),
    ("可扩展", "灵活的模型与算力管理\n混合部署支持多种环境\n按需扩展的数字员工团队"),
]

val_width = Inches(2.7)
val_height = Inches(2.5)

for i, (title, desc) in enumerate(values):
    x = Inches(0.8 + i * 3.05)
    y = Inches(4.2)
    
    card = add_shape(slide, x, y, val_width, val_height, fill_color=DARK_CARD)
    
    add_text_box(slide, x + Inches(0.25), y + Inches(0.2), val_width - Inches(0.5), Inches(0.4),
                 title, font_size=22, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_decorative_line(slide, x + Inches(0.8), y + Inches(0.65), Inches(1.1), ACCENT_CYAN, Pt(2))
    add_text_box(slide, x + Inches(0.25), y + Inches(0.85), val_width - Inches(0.5), Inches(1.5),
                 desc, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

print("价值主张页完成")

# ========== 第8页：应用场景 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
             "SCENARIOS", font_size=12, color=ACCENT_CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.8),
             "多场景覆盖，深入业务一线", font_size=36, color=WHITE, bold=True)
add_decorative_line(slide, Inches(0.8), Inches(1.55), Inches(2), ACCENT_BLUE)

scenarios = [
    ("📝", "办公与运营", ["通知起草与发送", "会议纪要整理", "周报/日报/月报生成", "文档归纳与知识查询", "跨部门事务跟进"]),
    ("📊", "数据与分析", ["数据报表生成", "多源数据汇总分析", "KPI监控与预警", "经营分析报告", "趋势分析与预测"]),
    ("📦", "供应链与采购", ["供应商信息管理", "采购流程辅助", "合同摘要与比对", "库存分析与建议", "价格趋势追踪"]),
    ("👥", "人力与行政", ["招聘JD生成", "员工问答与培训", "制度文档管理", "考勤与假期处理", "员工入职引导"]),
    ("🏭", "制造与质量", ["工艺文档管理", "质量异常分析", "检测报告生成", "设备维护计划", "标准操作指引"]),
]

sc_width = Inches(2.3)
sc_height = Inches(4.5)

for i, (icon, title, items) in enumerate(scenarios):
    x = Inches(0.5 + i * 2.5)
    y = Inches(2.2)
    
    card = add_shape(slide, x, y, sc_width, sc_height, fill_color=DARK_CARD)
    
    # 图标
    add_icon_circle(slide, x + Inches(0.7), y + Inches(0.25), Inches(0.7), RGBColor(0x00, 0x1F, 0x45))
    add_text_box(slide, x + Inches(0.7), y + Inches(0.3), Inches(0.7), Inches(0.6),
                 icon, font_size=22, color=WHITE, alignment=PP_ALIGN.CENTER)
    
    add_text_box(slide, x + Inches(0.15), y + Inches(1.1), sc_width - Inches(0.3), Inches(0.4),
                 title, font_size=16, color=ACCENT_CYAN, bold=True, alignment=PP_ALIGN.CENTER)
    
    add_decorative_line(slide, x + Inches(0.5), y + Inches(1.55), Inches(1.3), ACCENT_BLUE, Pt(1))
    
    for j, item in enumerate(items):
        add_text_box(slide, x + Inches(0.2), y + Inches(1.75 + j * 0.45), sc_width - Inches(0.4), Inches(0.4),
                     f"• {item}", font_size=11, color=LIGHT_GRAY)

print("场景页完成")

# ========== 第9页：目标客户 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
             "TARGET", font_size=12, color=ACCENT_CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.8),
             "iWorker 最适合谁", font_size=36, color=WHITE, bold=True)
add_decorative_line(slide, Inches(0.8), Inches(1.55), Inches(2), ACCENT_BLUE)

# 行业卡片
industries = ["制造业", "连锁零售", "医疗健康", "金融保险", "政企集团", "高合规企业"]
ind_card_w = Inches(1.8)
ind_card_h = Inches(0.7)

for i, ind in enumerate(industries):
    x = Inches(0.8 + i * 2.0)
    y = Inches(2.2)
    card = add_shape(slide, x, y, ind_card_w, ind_card_h, fill_color=ACCENT_BLUE)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.1), ind_card_w - Inches(0.2), Inches(0.5),
                 ind, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# 企业特征
add_text_box(slide, Inches(0.8), Inches(3.3), Inches(11), Inches(0.4),
             "最适合的企业特征", font_size=18, color=ACCENT_CYAN, bold=True)

features = [
    "希望AI真正进入业务流程，而不仅是员工自行试用",
    "关注数据安全、权限控制和审计要求",
    "拥有大量重复性、规则性、知识性工作",
    "希望先试点、后复制，逐步构建企业级AI能力",
    "希望把AI建成长期基础设施，而不是短期热点项目",
]

for i, feat in enumerate(features):
    add_text_box(slide, Inches(1.0), Inches(3.85 + i * 0.4), Inches(10), Inches(0.4),
                 f"✓  {feat}", font_size=14, color=LIGHT_GRAY)

# 决策者角色
add_text_box(slide, Inches(0.8), Inches(5.9), Inches(11), Inches(0.4),
             "关键决策角色", font_size=18, color=ACCENT_CYAN, bold=True)

roles = ["董事长/总经理", "CIO/CTO", "数字化负责人", "运营管理负责人", "财务管理负责人", "HR负责人"]
role_w = Inches(1.7)

for i, role in enumerate(roles):
    x = Inches(0.8 + i * 2.0)
    y = Inches(6.4)
    card = add_shape(slide, x, y, role_w, Inches(0.6), fill_color=DARK_CARD, line_color=ACCENT_BLUE)
    add_text_box(slide, x + Inches(0.05), y + Inches(0.08), role_w - Inches(0.1), Inches(0.45),
                 role, font_size=13, color=WHITE, alignment=PP_ALIGN.CENTER)

print("目标客户页完成")

# ========== 第10页：核心优势对比 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
             "COMPETITIVE ADVANTAGE", font_size=12, color=ACCENT_CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.8),
             "iWorker vs 传统AI工具", font_size=36, color=WHITE, bold=True)
add_decorative_line(slide, Inches(0.8), Inches(1.55), Inches(2), ACCENT_BLUE)

# 对比表格
headers = ["对比维度", "传统AI工具", "iWorker"]
rows = [
    ["定位", "个人效率工具", "企业组织能力系统"],
    ["使用方式", "各自试用，零散使用", "统一部署，组织级管理"],
    ["能力沉淀", "依赖个人经验", "沉淀到系统，可复制"],
    ["协作能力", "单点对话，无协作", "数字员工协作网络"],
    ["流程支持", "仅问答生成", "端到端流程闭环"],
    ["安全治理", "缺乏管控", "统一审计与权限管理"],
    ["扩展性", "难以规模化", "按需扩展，持续增长"],
]

# 表头
table_x = Inches(0.8)
table_y = Inches(2.2)
col_widths = [Inches(2.5), Inches(4.5), Inches(4.5)]
row_height = Inches(0.6)

# 绘制表头
for j, (header, cw) in enumerate(zip(headers, col_widths)):
    cx = table_x + sum(col_widths[:j])
    cell = add_shape(slide, cx, table_y, cw, row_height, fill_color=ACCENT_BLUE)
    add_text_box(slide, cx + Inches(0.15), table_y + Inches(0.05), cw - Inches(0.3), row_height - Inches(0.1),
                 header, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# 绘制数据行
for i, row in enumerate(rows):
    y = table_y + (i + 1) * row_height
    bg = DARK_CARD if i % 2 == 0 else RGBColor(0x0E, 0x12, 0x2E)
    
    for j, (cell_text, cw) in enumerate(zip(row, col_widths)):
        cx = table_x + sum(col_widths[:j])
        cell_color = bg
        text_color = WHITE if j == 0 else (MID_GRAY if j == 1 else ACCENT_CYAN)
        
        cell = add_shape(slide, cx, y, cw, row_height, fill_color=cell_color)
        add_text_box(slide, cx + Inches(0.15), y + Inches(0.05), cw - Inches(0.3), row_height - Inches(0.1),
                     cell_text, font_size=13, color=text_color, 
                     bold=(j == 0 or j == 2), alignment=PP_ALIGN.CENTER)

print("对比页完成")

# ========== 第11页：各方价值 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
             "STAKEHOLDER VALUE", font_size=12, color=ACCENT_CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.8),
             "为每个决策角色创造独特价值", font_size=36, color=WHITE, bold=True)
add_decorative_line(slide, Inches(0.8), Inches(1.55), Inches(2), ACCENT_BLUE)

stakeholders = [
    ("👔", "管理层", [
        "AI投入从工具采购变为组织能力建设",
        "数字化成果沉淀为长期资产",
        "建设面向未来的人机协作底座",
    ]),
    ("📊", "业务负责人", [
        "更快复制优秀岗位能力",
        "更稳定推动流程标准化",
        "更容易扩大人均产出",
    ]),
    ("💻", "IT负责人", [
        "统一的技术接入与治理方式",
        "清晰的安全边界与权限体系",
        "可控的模型与算力成本",
    ]),
    ("👤", "一线员工", [
        "更低的AI使用门槛",
        "更自然的任务入口",
        "更少重复劳动",
    ]),
]

stk_width = Inches(2.8)
stk_height = Inches(4.0)

for i, (icon, title, items) in enumerate(stakeholders):
    x = Inches(0.5 + i * 3.1)
    y = Inches(2.2)
    
    card = add_shape(slide, x, y, stk_width, stk_height, fill_color=DARK_CARD)
    
    add_icon_circle(slide, x + Inches(0.9), y + Inches(0.25), Inches(0.8), RGBColor(0x00, 0x1F, 0x45))
    add_text_box(slide, x + Inches(0.9), y + Inches(0.3), Inches(0.8), Inches(0.7),
                 icon, font_size=26, color=WHITE, alignment=PP_ALIGN.CENTER)
    
    add_text_box(slide, x + Inches(0.2), y + Inches(1.15), stk_width - Inches(0.4), Inches(0.4),
                 title, font_size=20, color=ACCENT_CYAN, bold=True, alignment=PP_ALIGN.CENTER)
    
    add_decorative_line(slide, x + Inches(0.7), y + Inches(1.6), Inches(1.4), ACCENT_BLUE, Pt(2))
    
    for j, item in enumerate(items):
        add_text_box(slide, x + Inches(0.25), y + Inches(1.85 + j * 0.6), stk_width - Inches(0.5), Inches(0.55),
                     f"▸ {item}", font_size=12, color=LIGHT_GRAY)

print("各方价值页完成")

# ========== 第12页：部署与安全 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
             "DEPLOYMENT & SECURITY", font_size=12, color=ACCENT_CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.8),
             "企业级部署，安全可控", font_size=36, color=WHITE, bold=True)
add_decorative_line(slide, Inches(0.8), Inches(1.55), Inches(2), ACCENT_BLUE)

# 部署模式
deploy_card = add_shape(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.5), fill_color=DARK_CARD)
add_text_box(slide, Inches(1.2), Inches(2.5), Inches(4.8), Inches(0.5),
             "🏗️  灵活部署模式", font_size=20, color=ACCENT_CYAN, bold=True)

deploy_items = [
    ("☁️  云端SaaS", "快速开通，按需使用"),
    ("🏢  专属环境", "独立资源，数据隔离"),
    ("🔒  私有化部署", "完全内网，自主可控"),
    ("🔀  混合模式", "灵活组合，兼顾安全与效率"),
]

for i, (mode, desc) in enumerate(deploy_items):
    add_text_box(slide, Inches(1.3), Inches(3.3 + i * 0.75), Inches(2.2), Inches(0.35),
                 mode, font_size=15, color=WHITE, bold=True)
    add_text_box(slide, Inches(3.5), Inches(3.3 + i * 0.75), Inches(2.5), Inches(0.35),
                 desc, font_size=13, color=LIGHT_GRAY)

# 安全保障
security_card = add_shape(slide, Inches(6.8), Inches(2.2), Inches(5.7), Inches(4.5), fill_color=DARK_CARD)
add_text_box(slide, Inches(7.2), Inches(2.5), Inches(5.0), Inches(0.5),
             "🛡️  安全与合规保障", font_size=20, color=ACCENT_CYAN, bold=True)

security_items = [
    "统一身份认证与权限管理",
    "数据传输与存储加密",
    "全链路操作审计与追溯",
    "敏感数据脱敏处理",
    "支持等保与行业合规要求",
    "模型调用的安全隔离与监控",
]

for i, item in enumerate(security_items):
    add_text_box(slide, Inches(7.3), Inches(3.3 + i * 0.55), Inches(5.0), Inches(0.4),
                 f"✓  {item}", font_size=13, color=LIGHT_GRAY)

print("部署安全页完成")

# ========== 第13页：商业模式 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
             "BUSINESS MODEL", font_size=12, color=ACCENT_CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.8),
             "可持续的商业模式", font_size=36, color=WHITE, bold=True)
add_decorative_line(slide, Inches(0.8), Inches(1.55), Inches(2), ACCENT_BLUE)

# 三大收入模式
biz_models = [
    ("💎", "平台订阅", "按组织/用户数订阅\n持续收入，可预测增长", "核心收入来源\n高续费率保障"),
    ("🧩", "场景能力包", "按行业/场景售卖能力包\n数字员工模板+流程方案", "高附加值\n边际成本低"),
    ("⚙️", "定制与实施", "大型客户定制化需求\n专属部署与深度集成", "高客单价\n标杆项目效应"),
]

bm_width = Inches(3.6)
bm_height = Inches(4.2)

for i, (icon, title, desc, value) in enumerate(biz_models):
    x = Inches(0.6 + i * 4.1)
    y = Inches(2.2)
    
    card = add_shape(slide, x, y, bm_width, bm_height, fill_color=DARK_CARD)
    
    add_icon_circle(slide, x + Inches(1.3), y + Inches(0.3), Inches(0.8), RGBColor(0x00, 0x1F, 0x45))
    add_text_box(slide, x + Inches(1.3), y + Inches(0.35), Inches(0.8), Inches(0.7),
                 icon, font_size=26, color=WHITE, alignment=PP_ALIGN.CENTER)
    
    add_text_box(slide, x + Inches(0.3), y + Inches(1.25), bm_width - Inches(0.6), Inches(0.4),
                 title, font_size=20, color=ACCENT_CYAN, bold=True, alignment=PP_ALIGN.CENTER)
    
    add_decorative_line(slide, x + Inches(1.0), y + Inches(1.7), Inches(1.6), ACCENT_BLUE, Pt(2))
    
    add_text_box(slide, x + Inches(0.3), y + Inches(1.9), bm_width - Inches(0.6), Inches(1.0),
                 desc, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    
    # 价值标注
    add_shape(slide, x + Inches(0.3), y + Inches(3.1), bm_width - Inches(0.6), Inches(0.8),
              fill_color=RGBColor(0x00, 0x14, 0x33))
    add_text_box(slide, x + Inches(0.4), y + Inches(3.15), bm_width - Inches(0.8), Inches(0.7),
                 value, font_size=12, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

print("商业模式页完成")

# ========== 第14页：市场愿景 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
             "VISION", font_size=12, color=ACCENT_CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.8),
             "从企业AI工具，到人机协作基础设施", font_size=36, color=WHITE, bold=True)
add_decorative_line(slide, Inches(0.8), Inches(1.55), Inches(2), ACCENT_BLUE)

# 发展路径
phases = [
    ("Phase 1", "单点突破", "打造标杆场景\n快速验证价值", "0-6个月"),
    ("Phase 2", "场景扩展", "覆盖核心业务线\n形成行业方案", "6-18个月"),
    ("Phase 3", "平台生态", "构建能力市场\n开放协作网络", "18-36个月"),
    ("Phase 4", "行业基础设施", "成为企业AI底座\n人机协作操作系统", "36个月+"),
]

# 连接线
add_shape(slide, Inches(1.5), Inches(3.35), Inches(10.3), Pt(3), fill_color=ACCENT_BLUE)

phase_w = Inches(2.5)
phase_h = Inches(3.5)

for i, (phase, title, desc, time) in enumerate(phases):
    x = Inches(0.8 + i * 3.1)
    y = Inches(2.3)
    
    # 时间节点圆
    add_icon_circle(slide, x + Inches(0.95), y + Inches(0.8), Inches(0.5), ACCENT_BLUE)
    add_text_box(slide, x + Inches(0.95), y + Inches(0.83), Inches(0.5), Inches(0.45),
                 str(i+1), font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    
    # 时间标注
    add_text_box(slide, x + Inches(0.3), y + Inches(1.4), phase_w - Inches(0.6), Inches(0.3),
                 time, font_size=11, color=ORANGE, alignment=PP_ALIGN.CENTER)
    
    # 内容卡片
    card = add_shape(slide, x, y + Inches(1.7), phase_w, Inches(2.3), fill_color=DARK_CARD)
    add_text_box(slide, x + Inches(0.2), y + Inches(1.85), phase_w - Inches(0.4), Inches(0.3),
                 phase, font_size=11, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), y + Inches(2.2), phase_w - Inches(0.4), Inches(0.4),
                 title, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_decorative_line(slide, x + Inches(0.6), y + Inches(2.65), Inches(1.3), ACCENT_CYAN, Pt(1))
    add_text_box(slide, x + Inches(0.2), y + Inches(2.85), phase_w - Inches(0.4), Inches(1.0),
                 desc, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

print("愿景页完成")

# ========== 第15页：结尾CTA ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

# 装饰光晕
add_shape(slide, Inches(5), Inches(0), Inches(8), Inches(7.5), 
          fill_color=RGBColor(0x00, 0x12, 0x35), shape_type=MSO_SHAPE.OVAL)

# 主标语
add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1.5),
             "让AI不再只是工具\n而是企业的核心能力",
             font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_decorative_line(slide, Inches(5.5), Inches(3.3), Inches(2.3), ACCENT_CYAN, Pt(3))

# 副标语
add_text_box(slide, Inches(2.5), Inches(3.8), Inches(8), Inches(1.0),
             "iWorker — 帮助每一家企业构建属于自己的AI组织能力系统",
             font_size=22, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)

# CTA按钮
cta = add_shape(slide, Inches(4.5), Inches(5.2), Inches(4.3), Inches(0.8), fill_color=ACCENT_BLUE)
add_text_box(slide, Inches(4.5), Inches(5.25), Inches(4.3), Inches(0.7),
             "联系我们  ·  开启AI组织能力建设",
             font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# 底部
add_text_box(slide, Inches(2), Inches(6.5), Inches(9), Inches(0.5),
             "iWorker  |  AI Native 企业组织能力系统  |  让能力沉淀在系统中",
             font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

print("结尾页完成")

# ========== 保存文件 ==========
output_path = r"D:\workprj\aicoder\docs\iWorker_宣传PPT.pptx"
prs.save(output_path)
print(f"\nPPT已保存至: {output_path}")
print(f"共 {len(prs.slides)} 页幻灯片")
