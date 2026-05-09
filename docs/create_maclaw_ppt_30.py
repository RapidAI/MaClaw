# -*- coding: utf-8 -*-
"""
MaClaw（码卡龙）产品介绍 PPT 生成脚本 — 30 页完整版
面向用户，突出功能特点、核心优势、竞品对比
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ========== 配色 ==========
PRIMARY      = RGBColor(0x0A, 0x0E, 0x27)
ACCENT_BLUE  = RGBColor(0x00, 0x7B, 0xFF)
ACCENT_CYAN  = RGBColor(0x00, 0xD4, 0xFF)
ACCENT_PURPLE= RGBColor(0x8B, 0x5C, 0xF6)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xB0, 0xB8, 0xCC)
MID_GRAY     = RGBColor(0x6B, 0x7B, 0x94)
DARK_CARD    = RGBColor(0x12, 0x18, 0x38)
DARKER_CARD  = RGBColor(0x0E, 0x13, 0x30)
ORANGE       = RGBColor(0xFF, 0x6B, 0x35)
GREEN        = RGBColor(0x00, 0xC8, 0x53)
RED          = RGBColor(0xFF, 0x44, 0x44)
YELLOW       = RGBColor(0xFF, 0xD6, 0x00)
PINK         = RGBColor(0xFF, 0x6B, 0xA0)

# ========== 工具函数 ==========
def bg(slide, color=PRIMARY):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def rect(slide, l, t, w, h, fill=None, line=None, st=MSO_SHAPE.ROUNDED_RECTANGLE):
    s = slide.shapes.add_shape(st, l, t, w, h); s.shadow.inherit = False
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else: s.fill.background()
    if line: s.line.color.rgb = line; s.line.width = Pt(1)
    else: s.line.fill.background()
    return s

def tb(slide, l, t, w, h, text, sz=18, c=WHITE, b=False, a=PP_ALIGN.LEFT, fn='Microsoft YaHei'):
    bx = slide.shapes.add_textbox(l, t, w, h); tf = bx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz); p.font.color.rgb = c
    p.font.bold = b; p.font.name = fn; p.alignment = a; return bx

def mp(slide, l, t, w, h, lines, sz=14, c=WHITE, a=PP_ALIGN.LEFT, fn='Microsoft YaHei'):
    bx = slide.shapes.add_textbox(l, t, w, h); tf = bx.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(line, tuple):
            p.text = line[0]; p.font.size = Pt(line[1] if len(line)>1 else sz)
            p.font.color.rgb = line[2] if len(line)>2 else c
            p.font.bold = line[3] if len(line)>3 else False
        else:
            p.text = line; p.font.size = Pt(sz); p.font.color.rgb = c
        p.font.name = fn; p.alignment = a; p.space_after = Pt(sz * 0.3)
    return bx

def line(slide, l, t, w, c=ACCENT_BLUE, h=Pt(3)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background(); s.shadow.inherit = False
    return s

def card(slide, x, y, w, h, title, items, icon="", tc=ACCENT_CYAN, ac=ACCENT_BLUE):
    rect(slide, x, y, w, h, fill=DARK_CARD, line=RGBColor(0x1E,0x2A,0x50))
    if icon: tb(slide, x+Inches(0.2), y+Inches(0.12), w-Inches(0.4), Inches(0.35), icon, sz=18, c=ac)
    tb(slide, x+Inches(0.2), y+Inches(0.45), w-Inches(0.4), Inches(0.35), title, sz=15, c=tc, b=True)
    line(slide, x+Inches(0.2), y+Inches(0.82), Inches(0.6), ac, Pt(2))
    mp(slide, x+Inches(0.2), y+Inches(0.95), w-Inches(0.4), h-Inches(1.1),
       [f"•  {it}" for it in items], sz=10, c=LIGHT_GRAY)

def section_header(slide, label, title, subtitle=""):
    tb(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.4), label, sz=11, c=ACCENT_CYAN, b=True)
    tb(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.7), title, sz=34, c=WHITE, b=True)
    line(slide, Inches(0.8), Inches(1.5), Inches(2.2), ACCENT_BLUE)
    if subtitle:
        tb(slide, Inches(0.8), Inches(1.7), Inches(11), Inches(0.4), subtitle, sz=14, c=LIGHT_GRAY)

def new_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); return s

print("工具函数准备完成")

# ================================================================
# P1: 封面
# ================================================================
s = new_slide()
bg(s, RGBColor(0x06,0x09,0x1C))
line(s, Inches(0), Inches(0), Inches(13.333), ACCENT_BLUE, Pt(4))
tb(s, Inches(1), Inches(1.2), Inches(11), Inches(0.7), "MaClaw  码卡龙", sz=58, c=WHITE, b=True, a=PP_ALIGN.CENTER)
line(s, Inches(4.8), Inches(2.15), Inches(3.7), ACCENT_CYAN, Pt(3))
tb(s, Inches(1.5), Inches(2.6), Inches(10), Inches(0.8), "通用可自进化智能体平台", sz=38, c=ACCENT_CYAN, a=PP_ALIGN.CENTER)
tb(s, Inches(2), Inches(3.7), Inches(9), Inches(0.8), "你的个人数智工作伙伴 — 不只是聊天，而是替你干活", sz=22, c=LIGHT_GRAY, a=PP_ALIGN.CENTER)
tags = ["长期记忆", "知识库", "自进化", "SSH远程", "GUI自动化", "19种工作流", "MIS数据"]
tw = Inches(1.55); th = Inches(0.45); total = len(tags)*tw.inches + (len(tags)-1)*0.18
sx = Inches((13.333-total)/2)
for i, tag in enumerate(tags):
    x = sx + Inches(i*(tw.inches+0.18))
    rect(s, x, Inches(5.0), tw, th, line=ACCENT_BLUE)
    tb(s, x, Inches(5.02), tw, th, tag, sz=12, c=ACCENT_CYAN, a=PP_ALIGN.CENTER)
tb(s, Inches(0), Inches(6.6), Inches(13.333), Inches(0.5), "Wails + Go + React  |  桌面 · IM · TUI 多形态  |  开箱即用  |  数据本地化", sz=12, c=MID_GRAY, a=PP_ALIGN.CENTER)
print("P1 封面 完成")

# ================================================================
# P2: 目录
# ================================================================
s = new_slide()
section_header(s, "CONTENTS", "目录")
toc_items = [
    ("01", "产品定位与核心理念", "P3-P5"),
    ("02", "六大核心能力总览", "P6-P8"),
    ("03", "内中外脑 · 三层记忆架构", "P9-P12"),
    ("04", "中外脑 · 知识库系统", "P13-P15"),
    ("05", "企业级结构化数据 (MIS)", "P16-P18"),
    ("06", "自进化与经验学习", "P19-P21"),
    ("07", "SSH / GUI 执行能力", "P22-P24"),
    ("08", "多通道协作与智能调度", "P25-P26"),
    ("09", "ag-UI 驱动结构化数据", "P27"),
    ("10", "技术架构与安全隐私", "P28"),
    ("11", "竞品对比 · 应用场景 · 快速开始", "P29-P31"),
]
for i, (num, title, pages) in enumerate(toc_items):
    y = Inches(2.0) + i * Inches(0.5)
    c = ACCENT_BLUE if i < 5 else ACCENT_PURPLE
    tb(s, Inches(1.2), y, Inches(0.6), Inches(0.4), num, sz=16, c=c, b=True)
    tb(s, Inches(2.0), y, Inches(6), Inches(0.4), title, sz=15, c=WHITE)
    tb(s, Inches(9.5), y, Inches(2), Inches(0.4), pages, sz=12, c=MID_GRAY, a=PP_ALIGN.RIGHT)
    if i < len(toc_items)-1:
        line(s, Inches(1.2), y+Inches(0.42), Inches(10), RGBColor(0x1E,0x2A,0x50), Pt(1))
print("P2 目录 完成")

# ================================================================
# P3: 产品定位
# ================================================================
s = new_slide()
section_header(s, "POSITIONING", "MaClaw 能做什么",
               "MaClaw 不只是聊天机器人，而是能理解意图、记住偏好、自主规划并执行复杂任务的智能体")
# 左侧大标语
rect(s, Inches(0.8), Inches(2.3), Inches(5.5), Inches(4.5), fill=DARK_CARD, line=RGBColor(0x1E,0x2A,0x50))
mp(s, Inches(1.3), Inches(2.6), Inches(4.5), Inches(3.8), [
    ("你说想法，它出成果", 30, ACCENT_CYAN, True),
    ("", 8, WHITE, False),
    ("内置 19 种结构化工作流模板，覆盖商业、研究、合规、学术、内容、技术全场景。", 14, LIGHT_GRAY, False),
    ("", 6, WHITE, False),
    ("每个工作流遵循「需求确认 → 方案设计 → 分步执行」的质量闭环，每个阶段产出文档等待审阅。", 14, LIGHT_GRAY, False),
    ("", 6, WHITE, False),
    ("你只需要说出想法，MaClaw 从需求梳理到成果交付全程陪你走完。", 14, ACCENT_BLUE, True),
    ("", 6, WHITE, False),
    ("> 不只是聊天，而是替你干活", 16, YELLOW, True),
])
# 右侧领域矩阵
domains = [
    ("商业与战略", "商业计划书 · 竞品分析 · 项目提案\n创新方案 · 招投标文件", ACCENT_BLUE),
    ("研究与分析", "文献综述 · 研究报告 · 实验设计\n专利分析", ACCENT_PURPLE),
    ("合规与尽调", "合同审查 · 尽职调查 · 合规审计", GREEN),
    ("学术写作", "基金申请书 · 论文写作", ORANGE),
    ("内容创作", "PPT 设计 · 活动策划", ACCENT_CYAN),
    ("产品与技术", "产品设计(PRD) · 软件测试 · 开发", YELLOW),
]
for i, (title, desc, color) in enumerate(domains):
    col = i % 2; row = i // 2
    x = Inches(6.8) + col*Inches(3.0); y = Inches(2.3) + row*Inches(1.5)
    rect(s, x, y, Inches(2.8), Inches(1.3), fill=DARK_CARD, line=RGBColor(0x1E,0x2A,0x50))
    tb(s, x+Inches(0.15), y+Inches(0.08), Inches(2.5), Inches(0.3), title, sz=13, c=color, b=True)
    line(s, x+Inches(0.15), y+Inches(0.38), Inches(0.5), color, Pt(2))
    tb(s, x+Inches(0.15), y+Inches(0.48), Inches(2.5), Inches(0.7), desc, sz=9, c=LIGHT_GRAY)
print("P3 产品定位 完成")

# ================================================================
# P4: 痛点——用户为什么需要 MaClaw
# ================================================================
s = new_slide()
section_header(s, "THE PROBLEM", "AI 应用的现实困境",
               "为什么现有 AI 工具无法真正帮到企业和个人？")
pains = [
    ("01", "用过就忘", "每次对话从零开始\n无法记住用户偏好和历史\n重复交代背景浪费时间", ACCENT_BLUE),
    ("02", "只会聊天", "只能生成文字建议\n无法操作电脑和服务器\n无法执行实际工作", ORANGE),
    ("03", "能力固化", "功能出厂即固定\n不会从错误中学习\n用 100 次和用 1 次没区别", RED),
    ("04", "数据外流", "企业数据传到云端\n无法本地部署\n合规审计困难", YELLOW),
    ("05", "单点工具", "需要在不同工具间切换\n无法串联完整工作流\n成果难以沉淀复用", MID_GRAY),
]
cw = Inches(2.2); ch = Inches(4.3); gap = Inches(0.22)
sx = Inches((13.333 - 5*cw.inches - 4*gap.inches)/2)
for i, (num, title, desc, color) in enumerate(pains):
    x = sx + i*(cw+gap); y = Inches(2.2)
    rect(s, x, y, cw, ch, fill=DARK_CARD, line=RGBColor(0x1E,0x2A,0x50))
    tb(s, x+Inches(0.2), y+Inches(0.25), Inches(0.8), Inches(0.5), num, sz=30, c=color, b=True)
    line(s, x+Inches(0.2), y+Inches(0.85), Inches(0.8), ACCENT_CYAN, Pt(2))
    tb(s, x+Inches(0.2), y+Inches(1.0), cw-Inches(0.4), Inches(0.4), title, sz=18, c=WHITE, b=True)
    tb(s, x+Inches(0.2), y+Inches(1.6), cw-Inches(0.4), Inches(2.5), desc, sz=12, c=LIGHT_GRAY)
print("P4 痛点 完成")

# ================================================================
# P5: 理念——MaClaw 如何解决
# ================================================================
s = new_slide()
section_header(s, "OUR SOLUTION", "MaClaw 的解决之道",
               "不是又一个聊天机器人，而是一个能记、能学、能做的智能体")
# 左右对比
rect(s, Inches(0.8), Inches(2.3), Inches(5.5), Inches(4.5), fill=DARK_CARD, line=RGBColor(0x1E,0x2A,0x50))
tb(s, Inches(1.2), Inches(2.5), Inches(4.8), Inches(0.4), "❌  传统 AI 助手", sz=20, c=ORANGE, b=True)
old = ["无状态，每次从零开始", "只能生成文字，无法执行", "能力固定，不会进化",
       "数据必须上传云端", "单一入口，无法远程", "无知识积累机制"]
for i, it in enumerate(old):
    tb(s, Inches(1.3), 3.3*Inches(1)+i*Inches(0.55), Inches(4.5), Inches(0.4),
       f"✗  {it}", sz=13, c=MID_GRAY)

rect(s, Inches(6.8), Inches(2.3), Inches(5.7), Inches(4.5), fill=DARK_CARD, line=ACCENT_BLUE)
tb(s, Inches(7.2), Inches(2.5), Inches(5.0), Inches(0.4), "✅  MaClaw 码卡龙", sz=20, c=GREEN, b=True)
new = ["内中外脑三层记忆，永不遗忘", "SSH/GUI 全域执行，真正替你干活",
       "Combee 经验蒸馏，越用越聪明", "数据本地化，支持私有部署",
       "桌面/IM/TUI 多通道，随时随地", "知识库蒸馏 · 结构化数据管理"]
for i, it in enumerate(new):
    tb(s, Inches(7.3), 3.3*Inches(1)+i*Inches(0.55), Inches(5.0), Inches(0.4),
       f"✓  {it}", sz=13, c=ACCENT_CYAN)
print("P5 理念 完成")

# ================================================================
# P6: 六大核心能力总览
# ================================================================
s = new_slide()
section_header(s, "CORE CAPABILITIES", "六大核心能力，构建完整智能体")
caps = [
    ("🧠", "长期记忆系统", ACCENT_BLUE,
     ["内中外脑三层架构","Ebbinghaus 遗忘曲线 + 自适应衰减","语义图谱 · 时序树 · 话题聚类","Claude 四类型分类法"]),
    ("📚", "中外脑知识库", ACCENT_PURPLE,
     ["支持 PDF/Word/Excel/MD/URL","文档结构化蒸馏，非简单切块","知识卡片 + 实体关系抽取","目录批量导入 · URL 安全抓取"]),
    ("🧬", "自我进化能力", GREEN,
     ["能力缺口自动检测与补全","Combee 并行经验聚合","技能自修复 · Nudge 推荐","工具路由学习与自适应"]),
    ("🖥️", "SSH/GUI 远程操控", ORANGE,
     ["SSH 多会话 · 后台任务 · SFTP","浏览器自动化 (CDP)","桌面 GUI 自动化","流程录制回放 · 定时触发"]),
    ("🔗", "多通道协作", ACCENT_CYAN,
     ["桌面/IM/TUI/REST 四形态","飞书 · 微信 · QQ · 钉钉","Swarm 多智能体协同","AgentNet P2P 知识网络"]),
    ("🏗️", "企业级结构化数据", YELLOW,
     ["多租户 · 多工作域","Schema 渐进演化 · 字段版本","API Key 权限策略 · 审计","外部连接器 · 数据导入导出"]),
]
cw = Inches(3.8); ch = Inches(2.0)
for i, (icon, title, color, items) in enumerate(caps):
    col = i%3; row = i//3
    x = Inches(0.8)+col*(cw+Inches(0.25)); y = Inches(2.0)+row*(ch+Inches(0.2))
    card(s, x, y, cw, ch, title, items, icon, tc=color, ac=color)
print("P6 核心能力 完成")

# ================================================================
# P7: 能力全景图（能力间关系）
# ================================================================
s = new_slide()
section_header(s, "CAPABILITY MAP", "能力全景——从感知到执行的全链路")
# 5 阶段流程
stages = [
    ("感知", "意图理解\n三层UIC", ACCENT_CYAN),
    ("记忆", "内中外脑\n长期记忆", ACCENT_BLUE),
    ("规划", "工作流引擎\n19种模板", ACCENT_PURPLE),
    ("执行", "SSH/GUI\n工具调用", GREEN),
    ("进化", "经验蒸馏\n技能沉淀", ORANGE),
]
sw = Inches(2.1); sh = Inches(1.2)
total = len(stages)*sw.inches + (len(stages)-1)*0.3
sx = Inches((13.333-total)/2)
for i, (name, desc, color) in enumerate(stages):
    x = sx + Inches(i*(sw.inches+0.3)); y = Inches(2.2)
    rect(s, x, y, sw, sh, fill=DARK_CARD, line=color)
    tb(s, x, y+Inches(0.05), sw, Inches(0.35), name, sz=16, c=color, b=True, a=PP_ALIGN.CENTER)
    line(s, x+Inches(0.3), y+Inches(0.42), sw-Inches(0.6), color, Pt(2))
    tb(s, x+Inches(0.1), y+Inches(0.5), sw-Inches(0.2), Inches(0.6), desc, sz=11, c=LIGHT_GRAY, a=PP_ALIGN.CENTER)
    if i < len(stages)-1:
        tb(s, x+sw, y+Inches(0.3), Inches(0.3), Inches(0.4), "→", sz=18, c=ACCENT_BLUE, a=PP_ALIGN.CENTER)

# 下半部分：详细说明
details = [
    ("感知层", "关键词规则(<1ms) → BM25语义检索(<5ms) → LLM确认(10-30s)\n渐进式意图识别，简单任务零延迟，复杂任务不遗漏", ACCENT_CYAN),
    ("记忆层", "对话记忆 → 工作记忆 → 长期记忆 → 知识库\nEbbinghaus衰减 + 语义图谱 + 时序树 三维度存储", ACCENT_BLUE),
    ("规划层", "结构化工作流 · 任务分解 · TDD验收\n需求→设计→编码→测试→集成→验收 全流程质量闭环", ACCENT_PURPLE),
    ("执行层", "197+动态工具 · SSH远程 · 浏览器CDP · 桌面GUI · 文件操作\n工具按上下文按需激活，不污染简单任务", GREEN),
    ("进化层", "Combee经验聚合 · 能力缺口检测 · 技能自修复 · Nudge推荐\n保守学习原则：单次仅提示，重复证据才改权重", ORANGE),
]
for i, (title, desc, color) in enumerate(details):
    y = Inches(3.8) + i*Inches(0.72)
    tb(s, Inches(1.0), y, Inches(1.2), Inches(0.35), title, sz=12, c=color, b=True)
    tb(s, Inches(2.3), y, Inches(10.0), Inches(0.6), desc, sz=10, c=LIGHT_GRAY)
    if i < len(details)-1:
        line(s, Inches(1.0), y+Inches(0.62), Inches(11.3), RGBColor(0x1E,0x2A,0x50), Pt(1))
print("P7 能力全景 完成")

# ================================================================
# P8: 核心差异化——一句话说清楚
# ================================================================
s = new_slide()
section_header(s, "DIFFERENTIATION", "MaClaw 的核心差异化")
# 中间大标语
rect(s, Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.8), fill=DARK_CARD, line=ACCENT_BLUE)
mp(s, Inches(2.0), Inches(2.15), Inches(9.3), Inches(1.5), [
    ("能记  ·  能学  ·  能做", 34, ACCENT_CYAN, True),
    ("", 6, WHITE, False),
    ("记忆是基础，进化是灵魂，执行是价值", 18, LIGHT_GRAY, False),
])
# 三大差异化卡片
diffs = [
    ("🧠", "三层记忆架构", ACCENT_BLUE, [
        "内脑：7种记忆类型 · Ebbinghaus衰减",
        "中脑：智能调度 · 5级渐进检索",
        "外脑：知识库蒸馏 · Card/Fact/Node",
        "语义图谱多跳推理 · 时序记忆树",
        "Claude四类型 · 冲突检测 · 语义去重",
    ], "其他产品：无状态 或 简单对话记忆"),
    ("🧬", "自进化引擎", GREEN, [
        "Combee并行经验聚合",
        "保守学习：单次提示，重复改权重",
        "保留具体路径/命令/错误号等细节",
        "能力缺口检测 → 自动搜索安装技能",
        "技能自修复 · Nudge成功序列封装",
    ], "其他产品：能力固定，不会学习"),
    ("🖥️", "全域执行能力", ORANGE, [
        "SSH：多会话 · 后台任务 · 文件传输",
        "浏览器：CDP · 流程录制回放",
        "桌面GUI：Accessibility + YOLO视觉",
        "197+动态工具按需激活",
        "桌面/IM/TUI多通道远程操控",
    ], "其他产品：只能聊天或调用API"),
]
dw = Inches(3.7); dh = Inches(3.0)
for i, (icon, title, color, items, vs) in enumerate(diffs):
    x = Inches(0.8)+i*(dw+Inches(0.25)); y = Inches(4.1)
    rect(s, x, y, dw, dh, fill=DARK_CARD, line=color)
    tb(s, x+Inches(0.2), y+Inches(0.12), dw-Inches(0.4), Inches(0.35),
       f"{icon}  {title}", sz=16, c=color, b=True)
    line(s, x+Inches(0.2), y+Inches(0.5), Inches(0.8), color, Pt(2))
    mp(s, x+Inches(0.2), y+Inches(0.6), dw-Inches(0.4), Inches(1.6),
       [f"•  {it}" for it in items], sz=10, c=LIGHT_GRAY)
    # 对比标注
    rect(s, x+Inches(0.1), y+dh-Inches(0.45), dw-Inches(0.2), Inches(0.35), fill=DARKER_CARD)
    tb(s, x+Inches(0.15), y+dh-Inches(0.42), dw-Inches(0.3), Inches(0.3),
       vs, sz=8, c=MID_GRAY, a=PP_ALIGN.CENTER)
print("P8 差异化 完成")

# ================================================================
# P9: 内中外脑——三层记忆架构总览（重点）
# ================================================================
s = new_slide()
section_header(s, "MEMORY ARCHITECTURE", "内中外脑——三层记忆架构",
               "从临时对话到永久知识，分层存储与智能调度")
brains = [
    ("内  脑", "corelib/memory", ACCENT_BLUE,
     ["用户长期偏好与身份 (self_identity)", "项目知识与决策记录 (project_knowledge)",
      "会话摘要与任务检查点 (conversation_summary)", "行为规则与反馈修正 (preference/instruction)",
      "高价值结论与反馈 (user_fact)", "Claude 四类型分类法"],
     "持久化 · Ebbinghaus衰减\n向量索引+BM25+语义图谱"),
    ("中  脑", "智能调度层", ACCENT_PURPLE,
     ["识别当前话题与用户意图", "决定什么时候存、存到哪",
      "分层召回策略 L1-L5", "查询复杂度自适应 (simple/hybrid/complex)",
      "冲突检测与语义去重", "RecallGating 防止过度召回"],
     "上下文感知 · 实时调度\n时序记忆树 · 5级渐进检索"),
    ("外  脑", "corelib/knowledge", GREEN,
     ["大型 PDF/Word/Excel/网页", "文档结构化节点 (DocumentNode)",
      "知识卡片 (KnowledgeCard)", "实体关系事实 (KnowledgeFact)",
      "目录批量导入 · URL 安全抓取", "分层召回: 卡片→事实→原文"],
     "海量存储 · SQLite FTS5\n结构化优先，避免原始RAG"),
]
bw = Inches(3.7); bh = Inches(4.3)
for i, (title, pkg, color, features, note) in enumerate(brains):
    x = Inches(0.7)+i*(bw+Inches(0.2)); y = Inches(2.0)
    rect(s, x, y, bw, bh, fill=DARK_CARD, line=color)
    tb(s, x+Inches(0.2), y+Inches(0.12), bw-Inches(0.4), Inches(0.35), title, sz=24, c=color, b=True, a=PP_ALIGN.CENTER)
    tb(s, x+Inches(0.2), y+Inches(0.5), bw-Inches(0.4), Inches(0.25), pkg, sz=10, c=MID_GRAY, a=PP_ALIGN.CENTER)
    line(s, x+Inches(0.4), y+Inches(0.78), bw-Inches(0.8), color, Pt(2))
    mp(s, x+Inches(0.2), y+Inches(0.9), bw-Inches(0.4), Inches(2.5),
       [f"✦  {f}" for f in features], sz=10, c=LIGHT_GRAY)
    rect(s, x+Inches(0.1), y+bh-Inches(0.7), bw-Inches(0.2), Inches(0.6), fill=DARKER_CARD)
    tb(s, x+Inches(0.15), y+bh-Inches(0.65), bw-Inches(0.3), Inches(0.5), note, sz=9, c=MID_GRAY, a=PP_ALIGN.CENTER)
    if i < 2:
        tb(s, x+bw, y+bh/2-Inches(0.15), Inches(0.2), Inches(0.3), "⟷", sz=16, c=color, a=PP_ALIGN.CENTER)
# 底部亮点
rect(s, Inches(0.7), Inches(6.6), Inches(11.9), Inches(0.6), fill=DARKER_CARD, line=RGBColor(0x1E,0x2A,0x50))
tb(s, Inches(1.0), Inches(6.65), Inches(11.3), Inches(0.5),
   "🧠 语义图谱多跳推理  |  📊 BM25+向量混合检索  |  🔄 Dream Cycle 自修复  |  📌 钉选保护  |  🔒 Owner隔离",
   sz=11, c=ACCENT_CYAN, a=PP_ALIGN.CENTER)
print("P9 内中外脑 完成")

# ================================================================
# P10: 内脑——7种记忆类型详解
# ================================================================
s = new_slide()
section_header(s, "INNER BRAIN — MEMORY TYPES", "内脑——7种记忆类型",
               "每种记忆有独立的存储策略、衰减规则和保护等级")
mem_types = [
    ("🛡️", "自我身份", "self_identity", "受保护，永不驱逐", ACCENT_BLUE),
    ("👤", "用户事实", "user_fact", "用户角色、目标、关键知识", ACCENT_CYAN),
    ("⚙️", "偏好指令", "preference / instruction", "行为规则与反馈修正", GREEN),
    ("📁", "项目知识", "project_knowledge", "决策、期限、架构、约束", ACCENT_PURPLE),
    ("💬", "会话摘要", "conversation_summary", "对话压缩与上下文保留", ORANGE),
    ("📋", "任务产物", "task_artifact", "工作流阶段输出与结论", YELLOW),
    ("🧬", "经验沉淀", "experience", "工具路由 · 技能模式 · 错误修复", PINK),
]
for i, (icon, name, code, desc, color) in enumerate(mem_types):
    col = i%2 if i < 6 else 0
    row = i//2 if i < 6 else 3
    x = Inches(0.8)+col*Inches(6.1); y = Inches(2.0)+row*Inches(1.25)
    rect(s, x, y, Inches(5.8), Inches(1.1), fill=DARK_CARD, line=RGBColor(0x1E,0x2A,0x50))
    tb(s, x+Inches(0.15), y+Inches(0.08), Inches(0.4), Inches(0.35), icon, sz=16, c=color)
    tb(s, x+Inches(0.55), y+Inches(0.08), Inches(2.0), Inches(0.35), name, sz=14, c=color, b=True)
    tb(s, x+Inches(2.6), y+Inches(0.08), Inches(3.0), Inches(0.35), code, sz=9, c=MID_GRAY)
    line(s, x+Inches(0.15), y+Inches(0.45), Inches(5.5), RGBColor(0x1E,0x2A,0x50), Pt(1))
    tb(s, x+Inches(0.15), y+Inches(0.55), Inches(5.5), Inches(0.45), desc, sz=11, c=LIGHT_GRAY)
print("P10 记忆类型 完成")

# ================================================================
# P11: 内脑——记忆管线与遗忘曲线
# ================================================================
s = new_slide()
section_header(s, "MEMORY PIPELINE", "记忆管线——从存储到检索的完整流程",
               "Ebbinghaus 遗忘曲线 + Dream Cycle 自修复 + 5级渐进检索")
# 左：存储管线
rect(s, Inches(0.8), Inches(2.1), Inches(5.5), Inches(5.0), fill=DARK_CARD, line=RGBColor(0x1E,0x2A,0x50))
tb(s, Inches(1.1), Inches(2.2), Inches(5.0), Inches(0.35), "存储管线", sz=16, c=ACCENT_BLUE, b=True)
line(s, Inches(1.1), Inches(2.55), Inches(1.0), ACCENT_BLUE, Pt(2))
pipeline_steps = [
    ("Step 1: 体验摄入", "ExperienceIngest", "对话消息 / 工具调用 / 工作流输出", ACCENT_CYAN),
    ("Step 2: 压缩提炼", "ConversationCompress", "LLM 生成摘要，保留关键信息", ACCENT_BLUE),
    ("Step 3: 提升入库", "PromoteToLongTerm", "高价值信息写入长期记忆", GREEN),
    ("Step 4: 反思整合", "ReflectAndConsolidate", "去重、冲突检测、图谱更新", ACCENT_PURPLE),
    ("Step 5: 衰减维护", "DreamCycle", "Ebbinghaus 衰减 + 过期清理", ORANGE),
]
for i, (step, func, desc, color) in enumerate(pipeline_steps):
    y = Inches(2.8)+i*Inches(0.82)
    rect(s, Inches(1.1), y, Inches(4.8), Inches(0.7), fill=DARKER_CARD, line=RGBColor(0x1E,0x2A,0x50))
    tb(s, Inches(1.3), y+Inches(0.02), Inches(2.5), Inches(0.25), step, sz=11, c=color, b=True)
    tb(s, Inches(3.8), y+Inches(0.02), Inches(2.0), Inches(0.25), func, sz=8, c=MID_GRAY, a=PP_ALIGN.RIGHT)
    tb(s, Inches(1.3), y+Inches(0.3), Inches(4.5), Inches(0.35), desc, sz=10, c=LIGHT_GRAY)
    if i < len(pipeline_steps)-1:
        tb(s, Inches(3.2), y+Inches(0.7), Inches(0.5), Inches(0.15), "↓", sz=10, c=MID_GRAY, a=PP_ALIGN.CENTER)

# 右：检索策略
rect(s, Inches(6.8), Inches(2.1), Inches(5.7), Inches(5.0), fill=DARK_CARD, line=RGBColor(0x1E,0x2A,0x50))
tb(s, Inches(7.1), Inches(2.2), Inches(5.0), Inches(0.35), "5 级渐进检索", sz=16, c=ACCENT_PURPLE, b=True)
line(s, Inches(7.1), Inches(2.55), Inches(1.0), ACCENT_PURPLE, Pt(2))
levels = [
    ("L1", "关键词精确匹配", "BM25 全文检索", "<10ms", ACCENT_BLUE),
    ("L2", "语义向量相似", "Embedding 向量召回", "<50ms", ACCENT_CYAN),
    ("L3", "时序记忆树", "Session → Day → Week → Profile", "<100ms", GREEN),
    ("L4", "语义图谱多跳", "Subject→Predicate→Object 推理", "<500ms", ORANGE),
    ("L5", "LLM 综合判断", "复杂查询最终裁决", "2-5s", ACCENT_PURPLE),
]
for i, (lv, name, desc, latency, color) in enumerate(levels):
    y = Inches(2.8)+i*Inches(0.82)
    rect(s, Inches(7.1), y, Inches(5.1), Inches(0.7), fill=DARKER_CARD, line=RGBColor(0x1E,0x2A,0x50))
    rect(s, Inches(7.1), y, Inches(0.6), Inches(0.7), fill=color, st=MSO_SHAPE.RECTANGLE)
    tb(s, Inches(7.1), y+Inches(0.12), Inches(0.6), Inches(0.35), lv, sz=14, c=WHITE, b=True, a=PP_ALIGN.CENTER)
    tb(s, Inches(7.9), y+Inches(0.02), Inches(2.0), Inches(0.25), name, sz=11, c=WHITE, b=True)
    tb(s, Inches(9.9), y+Inches(0.02), Inches(1.0), Inches(0.25), latency, sz=9, c=YELLOW)
    tb(s, Inches(7.9), y+Inches(0.3), Inches(4.2), Inches(0.35), desc, sz=9, c=MID_GRAY)
print("P11 记忆管线 完成")

# ================================================================
# P12: 内脑——语义图谱与冲突检测
# ================================================================
s = new_slide()
section_header(s, "SEMANTIC GRAPH", "语义图谱与冲突检测",
               "Subject-Predicate-Object 三元组 + 极性竞争 + 精确去重")
# 左：语义图谱
rect(s, Inches(0.8), Inches(2.0), Inches(5.5), Inches(5.0), fill=DARK_CARD, line=RGBColor(0x1E,0x2A,0x50))
tb(s, Inches(1.1), Inches(2.1), Inches(5.0), Inches(0.35), "语义图谱 (SemanticGraph)", sz=16, c=ACCENT_BLUE, b=True)
line(s, Inches(1.1), Inches(2.45), Inches(1.2), ACCENT_BLUE, Pt(2))
graph_features = [
    ("三元组结构", "Subject → Predicate → Object\n例: 「项目A」→「使用框架」→「React」", ACCENT_CYAN),
    ("实体索引", "双向索引：从实体查关系，从关系查实体\n支持 O(1) 实体查找 + O(k) 关系遍历", GREEN),
    ("多跳推理", "用户问「项目A的技术栈？」\n→ 图谱遍历发现 React/TypeScript/Wails", ACCENT_PURPLE),
    ("极性竞争", "矛盾事实自动标记\n「偏好中文」vs「用英文回复」→ 标记冲突供 LLM 裁决", ORANGE),
    ("增量更新", "新记忆入库时自动更新图谱\nADD/DELETE/UPDATE/MERGE 四操作", ACCENT_BLUE),
]
for i, (title, desc, color) in enumerate(graph_features):
    y = Inches(2.7)+i*Inches(0.82)
    tb(s, Inches(1.1), y, Inches(5.0), Inches(0.25), f"▸ {title}", sz=12, c=color, b=True)
    tb(s, Inches(1.3), y+Inches(0.25), Inches(4.8), Inches(0.5), desc, sz=9, c=LIGHT_GRAY)

# 右：冲突检测与去重
rect(s, Inches(6.8), Inches(2.0), Inches(5.7), Inches(5.0), fill=DARK_CARD, line=RGBColor(0x1E,0x2A,0x50))
tb(s, Inches(7.1), Inches(2.1), Inches(5.0), Inches(0.35), "冲突检测与语义去重", sz=16, c=ACCENT_PURPLE, b=True)
line(s, Inches(7.1), Inches(2.45), Inches(1.2), ACCENT_PURPLE, Pt(2))
dedup_features = [
    ("LLM 精确去重", "新记忆入库前，LLM 对比已有记忆\n「用户喜欢Python」vs「偏好Python编程」→ 去重", ACCENT_CYAN),
    ("矛盾事实标记", "检测到冲突时标记 PolarityConflict\n不自动删除，留给 LLM 裁决", RED),
    ("四操作更新", "ADD: 全新记忆\nDELETE: 确认过时\nUPDATE: 信息更新\nMERGE: 合并重叠", GREEN),
    ("来源溯源", "每条记忆记录 SourceTrace\n对话ID / 工具调用ID / 工作流阶段\n可追溯「为什么记住这个」", ORANGE),
    ("钉选保护", "Pin 机制防止高价值记忆被驱逐\nself_identity 类型默认钉选\n用户可手动钉选/解钉", YELLOW),
]
for i, (title, desc, color) in enumerate(dedup_features):
    y = Inches(2.7)+i*Inches(0.82)
    tb(s, Inches(7.1), y, Inches(5.0), Inches(0.25), f"▸ {title}", sz=12, c=color, b=True)
    tb(s, Inches(7.3), y+Inches(0.25), Inches(4.8), Inches(0.5), desc, sz=9, c=LIGHT_GRAY)
print("P12 语义图谱 完成")

# ================================================================
# P13: 知识库——来源与导入
# ================================================================
s = new_slide()
section_header(s, "KNOWLEDGE BASE — SOURCES", "中外脑知识库——来源与导入",
               "支持多种文档格式，结构化解析非简单切块")
# 来源卡片
sources = [
    ("📄", "PDF", "学术论文\n技术文档\n扫描件(OCR)", ACCENT_BLUE),
    ("📝", "DOCX", "合同\n报告\n方案书", ACCENT_PURPLE),
    ("📊", "XLSX", "数据表\n清单\n财务数据", GREEN),
    ("📋", "Markdown", "设计文档\nREADME\n笔记", ORANGE),
    ("🌐", "URL", "网页\n博客\n在线文档", ACCENT_CYAN),
    ("📁", "目录批量", "整站导入\n知识库\n文档集", YELLOW),
    ("🎤", "语音", "ASR转写\n会议录音\n声纹识别", PINK),
]
sw2 = Inches(1.55); sh2 = Inches(1.15)
for i, (icon, name, desc, color) in enumerate(sources):
    x = Inches(0.8)+i*(sw2+Inches(0.1)); y = Inches(2.0)
    rect(s, x, y, sw2, sh2, fill=DARK_CARD, line=RGBColor(0x1E,0x2A,0x50))
    tb(s, x, y+Inches(0.05), sw2, Inches(0.3), f"{icon} {name}", sz=12, c=color, b=True, a=PP_ALIGN.CENTER)
    line(s, x+Inches(0.3), y+Inches(0.38), sw2-Inches(0.6), color, Pt(1))
    tb(s, x+Inches(0.1), y+Inches(0.45), sw2-Inches(0.2), Inches(0.6), desc, sz=9, c=MID_GRAY, a=PP_ALIGN.CENTER)

# 导入流程
tb(s, Inches(0.8), Inches(3.5), Inches(12), Inches(0.4),
   "━━━━━━━━━━━━━━━━━━━  导入流水线  ━━━━━━━━━━━━━━━━━━━", sz=11, c=ACCENT_PURPLE, a=PP_ALIGN.CENTER)
pipeline = [
    ("上传文件", "拖拽 / URL / 目录", ACCENT_CYAN),
    ("格式解析", "结构化节点提取", ACCENT_BLUE),
    ("分段切分", "按章节/段落切分", GREEN),
    ("内容蒸馏", "Card + Fact 生成", ORANGE),
    ("入库索引", "FTS5 + 向量 + 图谱", ACCENT_PURPLE),
]
pw = Inches(2.0); ph = Inches(1.4)
total_p = len(pipeline)*pw.inches + (len(pipeline)-1)*0.2
sp = Inches((13.333-total_p)/2)
for i, (name, desc, color) in enumerate(pipeline):
    x = sp+Inches(i*(pw.inches+0.2)); y = Inches(4.1)
    rect(s, x, y, pw, ph, fill=DARK_CARD, line=color)
    tb(s, x, y+Inches(0.15), pw, Inches(0.3), name, sz=13, c=color, b=True, a=PP_ALIGN.CENTER)
    line(s, x+Inches(0.3), y+Inches(0.5), pw-Inches(0.6), color, Pt(2))
    tb(s, x+Inches(0.15), y+Inches(0.6), pw-Inches(0.3), Inches(0.6), desc, sz=10, c=LIGHT_GRAY, a=PP_ALIGN.CENTER)
    if i < len(pipeline)-1:
        tb(s, x+pw, y+Inches(0.4), Inches(0.2), Inches(0.3), "→", sz=14, c=MID_GRAY, a=PP_ALIGN.CENTER)

# 底部关键点
rect(s, Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.2), fill=DARK_CARD, line=RGBColor(0x1E,0x2A,0x50))
mp(s, Inches(1.1), Inches(5.9), Inches(11.1), Inches(1.0), [
    ("关键差异化：", 13, ACCENT_CYAN, True),
    ("•  不是简单的文档切块 (Chunk & Embed)，而是结构化解析保留文档层级 (DocumentNode)", 11, LIGHT_GRAY, False),
    ("•  蒸馏产出三层知识：KnowledgeCard(卡片) + KnowledgeFact(事实) + DocumentNode(原文节点)", 11, LIGHT_GRAY, False),
    ("•  分层召回策略：先查卡片摘要 → 再查事实三元组 → 最后定位原文片段，精度远超原始 RAG", 11, LIGHT_GRAY, False),
])
print("P13 知识库来源 完成")

# ================================================================
# P14: 知识库——三层产出物
# ================================================================
s = new_slide()
section_header(s, "KNOWLEDGE BASE — OUTPUTS", "知识库三层产出物",
               "KnowledgeCard + KnowledgeFact + DocumentNode，分层存储与分级召回")
outputs = [
    ("KnowledgeCard", "知识卡片", ACCENT_BLUE, [
        "5种类型: Summary / Concept / Procedure / Reference / FAQ",
        "质量评分 (Quality Score: 0-100)",
        "卡片级 Pin 与版本追踪",
        "支持 Batch 批量蒸馏",
        "自动标签与主题分类",
        "召回时优先返回卡片摘要",
    ], "第一层召回", "≈ 摘要卡片"),
    ("KnowledgeFact", "事实三元组", GREEN, [
        "Subject-Predicate-Object 结构化",
        "实体关系图谱自动构建",
        "与记忆系统语义图谱联动",
        "冲突检测与精确去重",
        "来源溯源 (SourceLink)",
        "多跳推理的基础单元",
    ], "第二层召回", "≈ 结构化事实"),
    ("DocumentNode", "文档结构树", ORANGE, [
        "章节 → 段落 → 句子层级",
        "SourceTimelineEvent 时间线",
        "原文片段精确定位",
        "增量更新与版本追踪",
        "全文搜索 BM25 索引",
        "最终回退到原文引用",
    ], "第三层召回", "≈ 原文片段"),
]
ow = Inches(3.7); oh = Inches(4.4)
for i, (cls, cn, color, features, recall, example) in enumerate(outputs):
    x = Inches(0.8)+i*(ow+Inches(0.2)); y = Inches(2.0)
    rect(s, x, y, ow, oh, fill=DARK_CARD, line=color)
    tb(s, x+Inches(0.2), y+Inches(0.15), ow-Inches(0.4), Inches(0.35), cn, sz=20, c=color, b=True)
    tb(s, x+Inches(0.2), y+Inches(0.5), ow-Inches(0.4), Inches(0.25), cls, sz=9, c=MID_GRAY)
    line(s, x+Inches(0.2), y+Inches(0.78), ow-Inches(0.4), color, Pt(2))
    mp(s, x+Inches(0.2), y+Inches(0.9), ow-Inches(0.4), Inches(2.5),
       [f"•  {f}" for f in features], sz=10, c=LIGHT_GRAY)
    # 召回标签
    rect(s, x+Inches(0.2), y+oh-Inches(0.9), ow-Inches(0.4), Inches(0.7), fill=DARKER_CARD)
    tb(s, x+Inches(0.2), y+oh-Inches(0.85), ow-Inches(0.4), Inches(0.3),
       recall, sz=11, c=color, b=True, a=PP_ALIGN.CENTER)
    tb(s, x+Inches(0.2), y+oh-Inches(0.5), ow-Inches(0.4), Inches(0.3),
       example, sz=9, c=MID_GRAY, a=PP_ALIGN.CENTER)
print("P14 三层产出 完成")

# ================================================================
# P15: 知识库——检索与使用场景
# ================================================================
s = new_slide()
section_header(s, "KNOWLEDGE BASE — USAGE", "知识库检索与应用场景",
               "分层召回策略确保精度，结合工作流实现端到端知识应用")
# 左：检索流程
rect(s, Inches(0.8), Inches(2.0), Inches(6.0), Inches(5.0), fill=DARK_CARD, line=RGBColor(0x1E,0x2A,0x50))
tb(s, Inches(1.1), Inches(2.1), Inches(5.5), Inches(0.35), "分层召回策略", sz=16, c=ACCENT_BLUE, b=True)
line(s, Inches(1.1), Inches(2.45), Inches(1.0), ACCENT_BLUE, Pt(2))
recall_steps = [
    ("用户提问 / 工作流引用", "触发知识库检索", ACCENT_CYAN),
    ("L1: 卡片摘要匹配", "BM25 + 语义向量快速过滤\n找到最相关的 KnowledgeCard", ACCENT_BLUE),
    ("L2: 事实三元组匹配", "从卡片中提取实体\n遍历 KnowledgeFact 图谱", GREEN),
    ("L3: 原文节点定位", "精确定位到 DocumentNode 段落\n获取完整上下文", ORANGE),
    ("结果融合与排序", "合并三层结果\n重排序后注入 LLM 上下文", ACCENT_PURPLE),
]
for i, (step, desc, color) in enumerate(recall_steps):
    y = Inches(2.7)+i*Inches(0.82)
    rect(s, Inches(1.1), y, Inches(5.4), Inches(0.7), fill=DARKER_CARD, line=RGBColor(0x1E,0x2A,0x50))
    tb(s, Inches(1.3), y+Inches(0.02), Inches(5.0), Inches(0.25), step, sz=11, c=color, b=True)
    tb(s, Inches(1.3), y+Inches(0.3), Inches(5.0), Inches(0.35), desc, sz=9, c=LIGHT_GRAY)
    if i < len(recall_steps)-1:
        tb(s, Inches(3.5), y+Inches(0.7), Inches(0.5), Inches(0.15), "↓", sz=10, c=MID_GRAY, a=PP_ALIGN.CENTER)

# 右：应用场景
rect(s, Inches(7.2), Inches(2.0), Inches(5.3), Inches(5.0), fill=DARK_CARD, line=RGBColor(0x1E,0x2A,0x50))
tb(s, Inches(7.5), Inches(2.1), Inches(4.8), Inches(0.35), "典型应用场景", sz=16, c=ACCENT_PURPLE, b=True)
line(s, Inches(7.5), Inches(2.45), Inches(1.0), ACCENT_PURPLE, Pt(2))
scenarios = [
    ("📚 企业知识库", "导入公司文档、制度、FAQ\n员工自然语言查询，精确返回答案", ACCENT_BLUE),
    ("⚖️ 合同审查", "上传合同PDF → 结构化解析\n自动提取关键条款、风险点、时间节点", GREEN),
    ("📊 数据分析", "导入 Excel → 表格结构识别\n自动生成分析报告和可视化建议", ORANGE),
    ("🔬 学术研究", "导入论文PDF → 蒸馏知识卡片\n文献综述自动关联相关事实", ACCENT_PURPLE),
    ("📖 技术文档", "导入项目文档 → 建立知识图谱\n开发团队快速查询架构决策", ACCENT_CYAN),
]
for i, (title, desc, color) in enumerate(scenarios):
    y = Inches(2.7)+i*Inches(0.82)
    tb(s, Inches(7.5), y, Inches(4.8), Inches(0.25), title, sz=12, c=color, b=True)
    tb(s, Inches(7.5), y+Inches(0.28), Inches(4.8), Inches(0.45), desc, sz=9, c=LIGHT_GRAY)
    if i < len(scenarios)-1:
        line(s, Inches(7.5), y+Inches(0.7), Inches(4.8), RGBColor(0x1E,0x2A,0x50), Pt(1))
print("P15 知识库使用 完成")

# ================================================================
# P16: MIS——结构化数据模型
# ================================================================
s = new_slide()
section_header(s, "STRUCTURED DATA — MODEL", "MIS 结构化数据模型",
               "Schema-less 渐进演化，16 种字段类型，版本全追踪")
# 核心概念卡片
mis_concepts = [
    ("🏢", "多租户隔离", ACCENT_BLUE, [
        "Tenant 组织边界隔离",
        "Domain 业务域: sales/hr/finance",
        "Dataset 数据集: schema-less",
        "Owner 级访问控制",
    ]),
    ("📋", "Record 数据模型", GREEN, [
        "灵活 JSON data + 固定元数据",
        "Schema 不强制，渐进演化",
        "Tag 标签系统 · 自动归一化",
        "record_ref 跨数据集关联",
    ]),
    ("📝", "FieldDefinition", ACCENT_PURPLE, [
        "16 种字段类型",
        "text/number/time/enum/file_ref",
        "record_ref/person_ref/org_ref",
        "字段级索引 · 非强制约束",
    ]),
    ("📊", "Revision 版本", ORANGE, [
        "RecordRevision 修改追溯",
        "谁修改 · 什么时候 · 改了什么",
        "完整变更历史回溯",
        "支持回滚到任意版本",
    ]),
]
cw2 = Inches(2.75); ch2 = Inches(2.2)
for i, (icon, title, color, items) in enumerate(mis_concepts):
    x = Inches(0.8)+i*(cw2+Inches(0.2)); y = Inches(2.0)
    card(s, x, y, cw2, ch2, title, items, icon, tc=color, ac=color)

# 下半：数据示例
rect(s, Inches(0.8), Inches(4.5), Inches(11.7), Inches(2.5), fill=DARK_CARD, line=RGBColor(0x1E,0x2A,0x50))
tb(s, Inches(1.1), Inches(4.6), Inches(5), Inches(0.35), "数据示例：Sales Orders", sz=14, c=ACCENT_BLUE, b=True)
line(s, Inches(1.1), Inches(4.95), Inches(1.0), ACCENT_BLUE, Pt(2))
example_fields = [
    ("id", "text", "ORD-2024-001", "主键"),
    ("title", "text", "Q4 设备采购订单", "标题"),
    ("customer", "person_ref", "→ HR-001 张三", "人员引用"),
    ("supplier", "org_ref", "→ SUP-001 华为", "组织引用"),
    ("amount", "number", "158000.00", "金额"),
    ("contract", "file_ref", "→ /contracts/ORD-001.pdf", "文件引用"),
    ("status", "enum", "pending_approval", "状态枚举"),
    ("tags", "text[]", "[\"紧急\",\"Q4\"]", "标签数组"),
]
# 表头
th_x = Inches(1.1); th_y = Inches(5.15)
for j, (header, w) in enumerate([("字段名", Inches(2.0)), ("类型", Inches(1.5)), ("示例值", Inches(3.0)), ("说明", Inches(1.5))]):
    rect(s, th_x, th_y, w, Inches(0.35), fill=RGBColor(0x0C,0x10,0x2E), st=MSO_SHAPE.RECTANGLE)
    tb(s, th_x+Inches(0.1), th_y+Inches(0.03), w-Inches(0.2), Inches(0.3), header, sz=10, c=MID_GRAY, b=True)
    th_x += w
# 数据行
for i, (name, ftype, example, note) in enumerate(example_fields):
    th_x = Inches(1.1); y = Inches(5.5)+i*Inches(0.35)
    for j, (val, w) in enumerate([(name, Inches(2.0)), (ftype, Inches(1.5)), (example, Inches(3.0)), (note, Inches(1.5))]):
        c2 = ACCENT_CYAN if j == 2 else (WHITE if j == 0 else MID_GRAY)
        tb(s, th_x+Inches(0.1), y, w-Inches(0.2), Inches(0.3), val, sz=9, c=c2)
        th_x += w
print("P16 MIS模型 完成")

# ================================================================
# P17: MIS——权限审计与质量管控
# ================================================================
s = new_slide()
section_header(s, "STRUCTURED DATA — GOVERNANCE", "权限、审计与质量管控",
               "企业级数据治理，从访问控制到质量检查全链路覆盖")
# 三大治理卡片
gov_items = [
    ("🔐", "权限策略", ACCENT_BLUE, [
        "API Key Policy 策略控制",
        "allowed_domains: 限制可访问域",
        "allowed_datasets: 限制可访问数据集",
        "allow_raw_data: 原始数据操作审批",
        "ConnectorContractBinding 连接器合约",
        "Owner 级数据隔离 (多租户)",
    ], "每个 API Key 绑定独立的访问策略\n越权操作自动拦截并审计记录"),
    ("📋", "操作审计", GREEN, [
        "DataEventLog 全量操作日志",
        "record.create / update / delete",
        "记录: who/when/what/before/after",
        "审计日志 CSV 导出",
        "字段级变更追踪",
        "时间线回放 (RecordTimeline)",
    ], "任何数据变更都有据可查\n支持合规审计与事后追溯"),
    ("✅", "数据质量", ORANGE, [
        "QualityCheck 质量检查规则",
        "schema_validation: Schema 合规",
        "unique_duplicates: 唯一性校验",
        "relationship_integrity: 引用完整性",
        "敏感字段标记与脱敏",
        "字段类型自动推断",
    ], "导入/更新时自动触发质量检查\n问题数据进入收件箱等待处理"),
]
gw = Inches(3.7); gh = Inches(4.2)
for i, (icon, title, color, items, note) in enumerate(gov_items):
    x = Inches(0.8)+i*(gw+Inches(0.2)); y = Inches(2.0)
    rect(s, x, y, gw, gh, fill=DARK_CARD, line=color)
    tb(s, x+Inches(0.2), y+Inches(0.12), gw-Inches(0.4), Inches(0.35),
       f"{icon}  {title}", sz=16, c=color, b=True)
    line(s, x+Inches(0.2), y+Inches(0.5), Inches(0.8), color, Pt(2))
    mp(s, x+Inches(0.2), y+Inches(0.6), gw-Inches(0.4), Inches(2.2),
       [f"•  {it}" for it in items], sz=10, c=LIGHT_GRAY)
    rect(s, x+Inches(0.1), y+gh-Inches(0.9), gw-Inches(0.2), Inches(0.8), fill=DARKER_CARD)
    tb(s, x+Inches(0.2), y+gh-Inches(0.85), gw-Inches(0.4), Inches(0.7), note, sz=9, c=MID_GRAY, a=PP_ALIGN.CENTER)
print("P17 权限审计 完成")

# ================================================================
# P18: MIS——外部连接器与数据导入导出
# ================================================================
s = new_slide()
section_header(s, "STRUCTURED DATA — INTEGRATION", "外部连接器与数据管道",
               "连接外部 CRM/ERP/HR 系统，数据导入导出全流程")
# 左：外部连接器
rect(s, Inches(0.8), Inches(2.0), Inches(5.7), Inches(5.0), fill=DARK_CARD, line=RGBColor(0x1E,0x2A,0x50))
tb(s, Inches(1.1), Inches(2.1), Inches(5.2), Inches(0.35), "外部连接器 (Connector)", sz=16, c=ACCENT_BLUE, b=True)
line(s, Inches(1.1), Inches(2.45), Inches(1.2), ACCENT_BLUE, Pt(2))
conn_features = [
    ("ExternalConnector 定义", "配置外部系统的连接信息\nREST API / 数据库 / 文件系统", ACCENT_CYAN),
    ("连接器健康监控", "ConnectorHealth 定期检查\n自动告警与恢复机制", GREEN),
    ("合约绑定", "ConnectorContractBinding\n控制连接器可访问的数据集", ORANGE),
    ("数据同步管道", "支持增量/全量同步\nCursor 分页 · 批量提交", ACCENT_PURPLE),
    ("安全接入", "token_ref 密钥引用\n不存储明文凭据 · 安全代理", RED),
]
for i, (title, desc, color) in enumerate(conn_features):
    y = Inches(2.7)+i*Inches(0.82)
    tb(s, Inches(1.1), y, Inches(5.2), Inches(0.25), f"▸ {title}", sz=12, c=color, b=True)
    tb(s, Inches(1.3), y+Inches(0.25), Inches(5.0), Inches(0.5), desc, sz=9, c=LIGHT_GRAY)

# 右：导入导出
rect(s, Inches(6.8), Inches(2.0), Inches(5.7), Inches(5.0), fill=DARK_CARD, line=RGBColor(0x1E,0x2A,0x50))
tb(s, Inches(7.1), Inches(2.1), Inches(5.2), Inches(0.35), "数据导入导出", sz=16, c=ACCENT_PURPLE, b=True)
line(s, Inches(7.1), Inches(2.45), Inches(1.2), ACCENT_PURPLE, Pt(2))
io_features = [
    ("CSV 批量导入", "start_csv_import_job\n第一行为字段头，支持元数据列", ACCENT_CYAN),
    ("JSONL 流式导入", "start_jsonl_import_job\n每行一个 JSON 对象，灵活结构", GREEN),
    ("CSV/JSONL 导出", "start_csv_export_job / start_jsonl_export_job\n异步任务，支持过滤和排序", ORANGE),
    ("备份与恢复", "create_backup → download_backup → restore_backup\n校验和验证，支持增量备份", ACCENT_PURPLE),
    ("操作计划", "create_operation_plan → review → apply\n高风险操作需管理员审阅确认", RED),
]
for i, (title, desc, color) in enumerate(io_features):
    y = Inches(2.7)+i*Inches(0.82)
    tb(s, Inches(7.1), y, Inches(5.2), Inches(0.25), f"▸ {title}", sz=12, c=color, b=True)
    tb(s, Inches(7.3), y+Inches(0.25), Inches(5.0), Inches(0.5), desc, sz=9, c=LIGHT_GRAY)
print("P18 连接器 完成")

# ================================================================
# P19: 自进化——经验学习机制
# ================================================================
s = new_slide()
section_header(s, "SELF-EVOLUTION — EXPERIENCE", "自进化引擎——经验学习机制",
               "Combee 并行经验聚合：从对话、工作流、工具调用中提取可复用经验")
# 中间标语
rect(s, Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.2), fill=DARK_CARD, line=ACCENT_BLUE)
mp(s, Inches(2.0), Inches(2.1), Inches(9.3), Inches(1.0), [
    ("保守学习原则：单次事件仅生成提示，重复证据才改变权重", 18, ACCENT_CYAN, True),
    ("保护少数高价值信号，不被通用摘要淹没", 13, LIGHT_GRAY, False),
])
# 经验来源 + 处理流程
evo_sources = [
    ("💬", "对话经验", "从用户对话中提取\n成功的操作序列", ACCENT_BLUE),
    ("📋", "工作流经验", "从 19 种工作流的\n执行路径中提炼", GREEN),
    ("🤖", "Swarm协作", "多智能体协同中\n发现最佳实践", ACCENT_PURPLE),
    ("🌐", "A2A 群聊", "跨智能体交流中\n获取外部经验", ORANGE),
]
sw3 = Inches(2.5); sh3 = Inches(1.4)
total_es = len(evo_sources)*sw3.inches + (len(evo_sources)-1)*0.2
ses = Inches((13.333-total_es)/2)
for i, (icon, title, desc, color) in enumerate(evo_sources):
    x = ses+Inches(i*(sw3.inches+0.2)); y = Inches(3.5)
    rect(s, x, y, sw3, sh3, fill=DARK_CARD, line=color)
    tb(s, x, y+Inches(0.1), sw3, Inches(0.3), f"{icon}  {title}", sz=13, c=color, b=True, a=PP_ALIGN.CENTER)
    tb(s, x+Inches(0.15), y+Inches(0.45), sw3-Inches(0.3), Inches(0.8), desc, sz=10, c=LIGHT_GRAY, a=PP_ALIGN.CENTER)
    if i < len(evo_sources)-1:
        tb(s, x+sw3, y+Inches(0.4), Inches(0.2), Inches(0.3), "→", sz=12, c=MID_GRAY, a=PP_ALIGN.CENTER)

# 处理流程
rect(s, Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.8), fill=DARK_CARD, line=RGBColor(0x1E,0x2A,0x50))
tb(s, Inches(1.1), Inches(5.3), Inches(5), Inches(0.35), "ExperienceDistiller 处理流程", sz=14, c=ACCENT_BLUE, b=True)
line(s, Inches(1.1), Inches(5.65), Inches(1.0), ACCENT_BLUE, Pt(2))
process_steps = [
    ("1. 收集", "从各来源获取原始事件流", ACCENT_CYAN),
    ("2. 过滤", "去除低价值/噪声事件", MID_GRAY),
    ("3. 提取", "保留具体路径、命令、错误号", GREEN),
    ("4. 聚合", "Combee 合并同类经验", ACCENT_PURPLE),
    ("5. 评估", "保守原则：验证后再写入", ORANGE),
    ("6. 存储", "经验沉淀为永久记忆", ACCENT_BLUE),
]
for i, (step, desc, color) in enumerate(process_steps):
    x = Inches(1.0)+i*Inches(1.9)
    tb(s, x, Inches(5.85), Inches(1.7), Inches(0.3), step, sz=11, c=color, b=True, a=PP_ALIGN.CENTER)
    tb(s, x, Inches(6.15), Inches(1.7), Inches(0.6), desc, sz=9, c=LIGHT_GRAY, a=PP_ALIGN.CENTER)
print("P19 经验学习 完成")

# ================================================================
# P20: 自进化——技能系统与 Nudge
# ================================================================
s = new_slide()
section_header(s, "SELF-EVOLUTION — SKILLS", "技能系统与 Nudge 推荐",
               "从临时脚本到永久技能，能力缺口自动补全")
# 左：技能生命周期
rect(s, Inches(0.8), Inches(2.0), Inches(5.7), Inches(5.0), fill=DARK_CARD, line=RGBColor(0x1E,0x2A,0x50))
tb(s, Inches(1.1), Inches(2.1), Inches(5.2), Inches(0.35), "技能生命周期", sz=16, c=ACCENT_BLUE, b=True)
line(s, Inches(1.1), Inches(2.45), Inches(1.0), ACCENT_BLUE, Pt(2))
skill_lifecycle = [
    ("craft_tool", "临时脚本", "遇到无法完成的任务\n自动生成一次性脚本", ACCENT_CYAN, "单次使用"),
    ("验证 & 注册", "技能验证", "验证脚本有效性\n注册到 SkillHub 本地", GREEN, "可复用"),
    ("SkillHub 发布", "技能发布", "上传到 SkillHub 市场\n供其他用户安装使用", ACCENT_PURPLE, "共享"),
    ("Nudge 推荐", "主动推荐", "检测到重复操作序列\n建议封装为技能", ORANGE, "自动化"),
]
for i, (step, name, desc, color, stage) in enumerate(skill_lifecycle):
    y = Inches(2.7)+i*Inches(1.0)
    rect(s, Inches(1.1), y, Inches(5.2), Inches(0.85), fill=DARKER_CARD, line=RGBColor(0x1E,0x2A,0x50))
    tb(s, Inches(1.3), y+Inches(0.02), Inches(1.5), Inches(0.25), step, sz=10, c=MID_GRAY)
    tb(s, Inches(1.3), y+Inches(0.25), Inches(2.5), Inches(0.25), name, sz=12, c=color, b=True)
    tb(s, Inches(3.8), y+Inches(0.02), Inches(1.5), Inches(0.25), stage, sz=9, c=YELLOW, a=PP_ALIGN.RIGHT)
    tb(s, Inches(1.3), y+Inches(0.5), Inches(4.8), Inches(0.3), desc, sz=9, c=LIGHT_GRAY)
    if i < len(skill_lifecycle)-1:
        tb(s, Inches(3.5), y+Inches(0.85), Inches(0.5), Inches(0.15), "↓", sz=10, c=MID_GRAY, a=PP_ALIGN.CENTER)

# 右：Nudge 与能力缺口
rect(s, Inches(6.8), Inches(2.0), Inches(5.7), Inches(2.3), fill=DARK_CARD, line=RGBColor(0x1E,0x2A,0x50))
tb(s, Inches(7.1), Inches(2.1), Inches(5.2), Inches(0.35), "💡  Nudge 推荐系统", sz=16, c=ORANGE, b=True)
line(s, Inches(7.1), Inches(2.45), Inches(1.0), ORANGE, Pt(2))
mp(s, Inches(7.1), Inches(2.6), Inches(5.2), Inches(1.5), [
    "•  检测到用户重复执行相同操作序列时自动触发",
    "•  建议封装为可复用技能 (Skill)",
    "•  用户确认后自动生成技能定义",
    "•  技能可在后续对话中直接调用",
], sz=10, c=LIGHT_GRAY)

rect(s, Inches(6.8), Inches(4.5), Inches(5.7), Inches(2.5), fill=DARK_CARD, line=RGBColor(0x1E,0x2A,0x50))
tb(s, Inches(7.1), Inches(4.6), Inches(5.2), Inches(0.35), "🔍  能力缺口检测", sz=16, c=GREEN, b=True)
line(s, Inches(7.1), Inches(4.95), Inches(1.0), GREEN, Pt(2))
mp(s, Inches(7.1), Inches(5.1), Inches(5.2), Inches(1.7), [
    "•  遇到无法完成的任务时，自动搜索 SkillHub 市场",
    "•  找到匹配技能 → 一键安装 → 立即可用",
    "•  搜索无果 → 记录能力缺口 → 等待后续补全",
    "•  缺口记录成为产品迭代的输入",
], sz=10, c=LIGHT_GRAY)
print("P20 技能系统 完成")

# ================================================================
# P21: 自进化——工具路由进化
# ================================================================
s = new_slide()
section_header(s, "SELF-EVOLUTION — ROUTING", "工具路由进化——越用越精准",
               "记录每次工具调用的执行轨迹，自适应调整路由策略")
# 上方：路由流程
rect(s, Inches(0.8), Inches(2.0), Inches(11.7), Inches(2.2), fill=DARK_CARD, line=RGBColor(0x1E,0x2A,0x50))
tb(s, Inches(1.1), Inches(2.1), Inches(5), Inches(0.35), "工具路由决策流程", sz=16, c=ACCENT_BLUE, b=True)
line(s, Inches(1.1), Inches(2.45), Inches(1.0), ACCENT_BLUE, Pt(2))
route_steps = [
    ("用户请求", ACCENT_CYAN),
    ("意图分类", ACCENT_BLUE),
    ("工具候选集", GREEN),
    ("上下文匹配", ACCENT_PURPLE),
    ("工具执行", ORANGE),
    ("结果验证", YELLOW),
    ("轨迹记录", ACCENT_BLUE),
]
rw = Inches(1.5); rh = Inches(1.0)
total_r = len(route_steps)*rw.inches + (len(route_steps)-1)*0.1
sr = Inches((13.333-total_r)/2)
for i, (step, color) in enumerate(route_steps):
    x = sr+Inches(i*(rw.inches+0.1)); y = Inches(2.8)
    rect(s, x, y, rw, rh, fill=DARKER_CARD, line=color)
    tb(s, x, y+Inches(0.25), rw, Inches(0.3), step, sz=11, c=color, b=True, a=PP_ALIGN.CENTER)
    if i < len(route_steps)-1:
        tb(s, x+rw, y+Inches(0.25), Inches(0.1), Inches(0.3), "→", sz=12, c=MID_GRAY, a=PP_ALIGN.CENTER)

# 下方：轨迹数据与进化策略
rect(s, Inches(0.8), Inches(4.5), Inches(5.7), Inches(2.5), fill=DARK_CARD, line=RGBColor(0x1E,0x2A,0x50))
tb(s, Inches(1.1), Inches(4.6), Inches(5.2), Inches(0.35), "📊  执行轨迹数据", sz=16, c=ACCENT_PURPLE, b=True)
line(s, Inches(1.1), Inches(4.95), Inches(1.0), ACCENT_PURPLE, Pt(2))
mp(s, Inches(1.1), Inches(5.1), Inches(5.2), Inches(1.7), [
    "•  task_type: 任务类型分类",
    "•  tool_sequence: 工具调用序列",
    "•  error_class: 错误类型与错误号",
    "•  retry_count: 重试次数",
    "•  recovery_tool: 恢复使用的替代工具",
    "•  final_outcome: 最终执行结果",
], sz=10, c=LIGHT_GRAY)

rect(s, Inches(6.8), Inches(4.5), Inches(5.7), Inches(2.5), fill=DARK_CARD, line=RGBColor(0x1E,0x2A,0x50))
tb(s, Inches(7.1), Inches(4.6), Inches(5.2), Inches(0.35), "🧬  进化策略", sz=16, c=GREEN, b=True)
line(s, Inches(7.1), Inches(4.95), Inches(1.0), GREEN, Pt(2))
mp(s, Inches(7.1), Inches(5.1), Inches(5.2), Inches(1.7), [
    "•  高失败率工具自动降权",
    "•  新工具渐进激活，避免突然变化",
    "•  DistillRoutingHints 聚合最近轨迹",
    "•  聚合结果注入 LLM 上下文作为路由提示",
    "•  按上下文关键词按需激活工具",
    "•  不污染简单任务，保持低延迟",
], sz=10, c=LIGHT_GRAY)
print("P21 工具路由 完成")

# ================================================================
# P22: SSH远程服务器管理
# ================================================================
s = new_slide()
section_header(s, "EXECUTION — SSH", "SSH 远程服务器管理",
               "多会话并发，密码/密钥/Agent 多种认证，支持后台任务与文件传输")
features_ssh = [
    ("🔌", "连接管理", ACCENT_BLUE, [
        "多会话并发连接",
        "密码 / 密钥 / Agent 认证",
        "自动 sudo token 管理",
        "Shell 响应性检测",
        "连接保活与超时处理",
    ]),
    ("⚡", "命令执行", GREEN, [
        "同步执行 · 实时输出",
        "后台任务 (exec_background)",
        "长命令自动转后台",
        "任务进度查询 (check_task)",
        "连续失败自动清理",
    ]),
    ("📁", "文件传输", ACCENT_PURPLE, [
        "SFTP 上传 / 下载",
        "目录递归传输",
        "大文件分块传输",
        "断点续传支持",
        "传输进度追踪",
    ]),
    ("🔒", "安全审计", ORANGE, [
        "操作全程审计日志",
        "命令白名单 / 黑名单",
        "敏感操作二次确认",
        "会话超时自动断开",
        "IM 通道远程操控",
    ]),
]
fw = Inches(2.8); fh = Inches(2.2)
for i, (icon, title, color, items) in enumerate(features_ssh):
    x = Inches(0.8)+i*(fw+Inches(0.2)); y = Inches(2.0)
    card(s, x, y, fw, fh, title, items, icon, tc=color, ac=color)

# 典型场景
rect(s, Inches(0.8), Inches(4.5), Inches(11.7), Inches(2.5), fill=DARK_CARD, line=RGBColor(0x1E,0x2A,0x50))
tb(s, Inches(1.1), Inches(4.6), Inches(5), Inches(0.35), "典型使用场景", sz=14, c=ACCENT_BLUE, b=True)
line(s, Inches(1.1), Inches(4.95), Inches(1.0), ACCENT_BLUE, Pt(2))
ssh_scenarios = [
    ("运维巡检", "\"检查生产服务器状态\" → SSH连接 → 执行 top/df/free\n→ 分析输出 → 生成巡检报告", ACCENT_CYAN),
    ("日志排查", "\"查看 nginx 错误日志\" → SSH连接 → tail/grep\n→ 过滤关键错误 → 总结问题原因", GREEN),
    ("部署发布", "\"部署最新版本到测试环境\" → SSH连接 → git pull\n→ 构建部署 → 健康检查 → 汇报结果", ORANGE),
    ("安全审计", "\"检查服务器安全配置\" → SSH连接 → 扫描配置\n→ 对比安全基线 → 生成审计报告", RED),
]
for i, (title, desc, color) in enumerate(ssh_scenarios):
    y = Inches(5.15)+i*Inches(0.45)
    tb(s, Inches(1.3), y, Inches(1.5), Inches(0.3), title, sz=10, c=color, b=True)
    tb(s, Inches(2.8), y, Inches(9.5), Inches(0.35), desc, sz=9, c=LIGHT_GRAY)
print("P22 SSH 完成")

# ================================================================
# P23: 浏览器自动化 (CDP)
# ================================================================
s = new_slide()
section_header(s, "EXECUTION — BROWSER", "浏览器自动化",
               "基于 Chrome DevTools Protocol，支持复杂网页交互与数据提取")
# 功能矩阵
browser_features = [
    ("🌐", "页面导航", ACCENT_BLUE, ["URL 导航", "前进/后退/刷新", "多标签页管理", "页面加载等待"]),
    ("🖱️", "元素交互", GREEN, ["点击 (CSS Selector/Ref)", "文本输入", "下拉框选择", "文件上传", "滚动页面"]),
    ("👁️", "内容提取", ACCENT_PURPLE, ["文本提取 (Selector/整页)", "HTML 获取", "OCR 文字识别", "结构化数据提取"]),
    ("📸", "截图与录制", ORANGE, ["全页/区域截图", "流程录制与回放", "参数化变量替换", "后台异步执行"]),
]
bw2 = Inches(2.8); bh2 = Inches(1.8)
for i, (icon, title, color, items) in enumerate(browser_features):
    x = Inches(0.8)+i*(bw2+Inches(0.2)); y = Inches(2.0)
    card(s, x, y, bw2, bh2, title, items, icon, tc=color, ac=color)

# OCR 集成说明
rect(s, Inches(0.8), Inches(4.1), Inches(5.7), Inches(3.0), fill=DARK_CARD, line=RGBColor(0x1E,0x2A,0x50))
tb(s, Inches(1.1), Inches(4.2), Inches(5.2), Inches(0.35), "🔍  OCR 集成", sz=14, c=ACCENT_BLUE, b=True)
line(s, Inches(1.1), Inches(4.55), Inches(1.0), ACCENT_BLUE, Pt(2))
mp(s, Inches(1.1), Inches(4.7), Inches(5.2), Inches(2.2), [
    ("双引擎 OCR 方案：", 12, ACCENT_CYAN, True),
    ("1. RapidOCR (本地)", 11, LIGHT_GRAY, False),
    ("   离线运行，零延迟，保护隐私", 10, MID_GRAY, False),
    ("2. LLM 视觉理解 (云端)", 11, LIGHT_GRAY, False),
    ("   处理复杂图表、手写体、多语言", 10, MID_GRAY, False),
    ("", 6, WHITE, False),
    ("自动选择策略：", 12, ACCENT_CYAN, True),
    ("•  简单文字识别 → RapidOCR (快)", 10, LIGHT_GRAY, False),
    ("•  复杂内容理解 → LLM (准)", 10, LIGHT_GRAY, False),
], sz=10, c=LIGHT_GRAY)

# 典型应用
rect(s, Inches(6.8), Inches(4.1), Inches(5.7), Inches(3.0), fill=DARK_CARD, line=RGBColor(0x1E,0x2A,0x50))
tb(s, Inches(7.1), Inches(4.2), Inches(5.2), Inches(0.35), "📋  典型应用场景", sz=14, c=ACCENT_PURPLE, b=True)
line(s, Inches(7.1), Inches(4.55), Inches(1.0), ACCENT_PURPLE, Pt(2))
browser_cases = [
    ("竞品监控", "定时访问竞品网站 → 截图对比 → 变更提醒", ACCENT_CYAN),
    ("数据采集", "登录系统 → 导航到报表页 → 提取数据 → 生成报告", GREEN),
    ("自动化测试", "执行测试步骤 → 验证页面状态 → 生成测试报告", ORANGE),
    ("表单填写", "读取 Excel 数据 → 自动填写网页表单 → 提交", ACCENT_PURPLE),
    ("价格跟踪", "定时访问电商页面 → OCR 提取价格 → 价格趋势图", YELLOW),
]
for i, (title, desc, color) in enumerate(browser_cases):
    y = Inches(4.8)+i*Inches(0.42)
    tb(s, Inches(7.1), y, Inches(1.3), Inches(0.3), f"▸ {title}", sz=10, c=color, b=True)
    tb(s, Inches(8.4), y, Inches(3.8), Inches(0.3), desc, sz=9, c=LIGHT_GRAY)
print("P23 浏览器 完成")

# ================================================================
# P24: 桌面 GUI 自动化
# ================================================================
s = new_slide()
section_header(s, "EXECUTION — DESKTOP GUI", "桌面 GUI 自动化",
               "Accessibility 元素树 + YOLO 视觉检测，双引擎覆盖所有桌面应用")
# 双引擎对比
rect(s, Inches(0.8), Inches(2.0), Inches(5.7), Inches(3.0), fill=DARK_CARD, line=ACCENT_BLUE)
tb(s, Inches(1.1), Inches(2.1), Inches(5.2), Inches(0.35), "🌳  Accessibility 元素树", sz=16, c=ACCENT_BLUE, b=True)
line(s, Inches(1.1), Inches(2.45), Inches(1.0), ACCENT_BLUE, Pt(2))
mp(s, Inches(1.1), Inches(2.6), Inches(5.2), Inches(2.2), [
    "•  调用系统 Accessibility API 获取元素树",
    "•  支持 Windows / macOS / Linux",
    "•  精确定位：按名称、角色、ID 查找元素",
    "•  可靠性高，不受 UI 外观变化影响",
    "•  限制：部分应用不支持 Accessibility",
    "•  适用：标准桌面应用、系统设置",
], sz=11, c=LIGHT_GRAY)

rect(s, Inches(6.8), Inches(2.0), Inches(5.7), Inches(3.0), fill=DARK_CARD, line=ORANGE)
tb(s, Inches(7.1), Inches(2.1), Inches(5.2), Inches(0.35), "👁️  YOLO 视觉检测 (OmniParser V2)", sz=16, c=ORANGE, b=True)
line(s, Inches(7.1), Inches(2.45), Inches(1.0), ORANGE, Pt(2))
mp(s, Inches(7.1), Inches(2.6), Inches(5.2), Inches(2.2), [
    "•  基于 YOLO 模型检测屏幕上的可交互元素",
    "•  不依赖 Accessibility API",
    "•  覆盖所有应用：游戏、虚拟机、远程桌面",
    "•  坐标级精确定位与点击",
    "•  支持 OmniParser V2 增强识别",
    "•  适用：任何可见的 GUI 界面",
], sz=11, c=LIGHT_GRAY)

# 下半：功能列表
gui_features = [
    ("🖱️", "鼠标操作", ACCENT_BLUE, ["精确点击 (坐标/元素)", "双击/右键/拖拽", "滚轮滚动", "坐标偏移微调"]),
    ("⌨️", "键盘操作", GREEN, ["文本输入", "快捷键组合", "特殊键支持", "输入法兼容"]),
    ("📸", "截图分析", ACCENT_PURPLE, ["全屏/区域截图", "OCR 文字识别", "颜色检测", "图像对比"]),
    ("🔄", "流程录制", ORANGE, ["操作序列录制", "参数化变量替换", "回放执行", "定时触发"]),
]
gw2 = Inches(2.8); gh2 = Inches(2.0)
for i, (icon, title, color, items) in enumerate(gui_features):
    x = Inches(0.8)+i*(gw2+Inches(0.2)); y = Inches(5.3)
    card(s, x, y, gw2, gh2, title, items, icon, tc=color, ac=color)
print("P24 GUI 完成")

# ================================================================
# P25: 多通道协作与 Swarm
# ================================================================
s = new_slide()
section_header(s, "MULTI-CHANNEL & SWARM", "多通道协作与 Swarm 多智能体",
               "桌面/IM/TUI 统一体验，Swarm 协同解决复杂任务")
# 上方：四种交互形态
forms = [
    ("🖥️", "桌面 AI 面板", ACCENT_BLUE, "Wails 桌面应用",
     ["本地 GUI 完整体验", "文件操作 · 屏幕截图 · 进程管理", "即时响应，零网络延迟", "功能最完整的入口"]),
    ("💬", "IM 消息通道", GREEN, "飞书 / 微信 / QQ / 钉钉 / 企微",
     ["通过 Hub 消息代理", "随时随地远程操控", "语音消息自动 ASR 转文字", "支持文件/截图收发"]),
    ("⌨️", "TUI 终端", ACCENT_PURPLE, "Bubble Tea 终端界面",
     ["开发者友好", "SSH 服务器环境可用", "REST API (MaClawSrv)", "轻量级部署"]),
    ("🌐", "REST API", ORANGE, "MaClawSrv 服务端",
     ["HTTP API 接口", "第三方系统集成", "自动化管道接入", "Webhook 回调支持"]),
]
fw2 = Inches(2.8); fh2 = Inches(2.3)
for i, (icon, title, color, sub, items) in enumerate(forms):
    x = Inches(0.8)+i*(fw2+Inches(0.2)); y = Inches(2.0)
    rect(s, x, y, fw2, fh2, fill=DARK_CARD, line=RGBColor(0x1E,0x2A,0x50))
    tb(s, x+Inches(0.15), y+Inches(0.08), fw2-Inches(0.3), Inches(0.3), f"{icon}  {title}", sz=13, c=color, b=True)
    line(s, x+Inches(0.15), y+Inches(0.38), fw2-Inches(0.3), color, Pt(1))
    mp(s, x+Inches(0.15), y+Inches(0.48), fw2-Inches(0.3), Inches(1.3),
       [f"•  {it}" for it in items], sz=9, c=LIGHT_GRAY)
    rect(s, x+Inches(0.05), y+fh2-Inches(0.35), fw2-Inches(0.1), Inches(0.28), fill=DARKER_CARD)
    tb(s, x+Inches(0.1), y+fh2-Inches(0.33), fw2-Inches(0.2), Inches(0.24), sub, sz=8, c=MID_GRAY, a=PP_ALIGN.CENTER)

# 下方：Swarm + 远程拓扑
rect(s, Inches(0.8), Inches(4.6), Inches(5.7), Inches(2.5), fill=DARK_CARD, line=RGBColor(0x1E,0x2A,0x50))
tb(s, Inches(1.1), Inches(4.7), Inches(5.2), Inches(0.35), "🐝  Swarm 多智能体协同", sz=14, c=ACCENT_PURPLE, b=True)
line(s, Inches(1.1), Inches(5.05), Inches(1.0), ACCENT_PURPLE, Pt(2))
mp(s, Inches(1.1), Inches(5.2), Inches(5.2), Inches(1.7), [
    "•  多个 Agent 实例并行工作",
    "•  任务自动分解与分配",
    "•  Agent 间消息通信与状态同步",
    "•  协同经验回流到记忆系统",
    "•  支持跨 Agent 知识共享",
], sz=10, c=LIGHT_GRAY)

rect(s, Inches(6.8), Inches(4.6), Inches(5.7), Inches(2.5), fill=DARK_CARD, line=RGBColor(0x1E,0x2A,0x50))
tb(s, Inches(7.1), Inches(4.7), Inches(5.2), Inches(0.35), "🌐  远程控制拓扑", sz=14, c=ACCENT_CYAN, b=True)
line(s, Inches(7.1), Inches(5.05), Inches(1.0), ACCENT_CYAN, Pt(2))
topo = "IM发消息 → Hub路由 → MaClaw桌面端执行 → 结果回传 → IM收到回复"
tb(s, Inches(7.1), Inches(5.3), Inches(5.2), Inches(0.4), topo, sz=10, c=LIGHT_GRAY)
mp(s, Inches(7.1), Inches(5.8), Inches(5.2), Inches(1.0), [
    "•  支持文件/截图/语音消息收发",
    "•  群聊 @提及 触发 MaClaw 响应",
    "•  多人共享同一 MaClaw 实例",
    "•  操作审计与权限隔离",
], sz=10, c=MID_GRAY)
print("P25 多通道 完成")

# ================================================================
# P26: 智能调度——UIC + Steering
# ================================================================
s = new_slide()
section_header(s, "INTELLIGENT ROUTING", "智能调度——意图理解与规则引擎",
               "三层意图分类 + Steering 规则系统，精准理解需求高效匹配工具")
# 上：三层UIC
rect(s, Inches(0.8), Inches(2.0), Inches(11.7), Inches(2.3), fill=DARK_CARD, line=RGBColor(0x1E,0x2A,0x50))
tb(s, Inches(1.1), Inches(2.1), Inches(5), Inches(0.35), "三层意图分类 (UIC)", sz=16, c=ACCENT_BLUE, b=True)
line(s, Inches(1.1), Inches(2.45), Inches(1.0), ACCENT_BLUE, Pt(2))
uic = [
    ("Layer 1", "关键词规则", "<1ms", "精确匹配高频模式\n如「写代码」「发邮件」", ACCENT_BLUE),
    ("Layer 2", "BM25 语义", "<5ms", "模糊匹配与工作流模板\n自然语言意图识别", GREEN),
    ("Layer 3", "LLM 确认", "10-30s", "复杂意图理解与澄清\n多轮对话确认需求", ORANGE),
]
for i, (name, method, latency, desc, color) in enumerate(uic):
    x = Inches(1.0)+i*Inches(3.8); y = Inches(2.8)
    rect(s, x, y, Inches(3.5), Inches(1.3), fill=DARKER_CARD, line=color)
    tb(s, x+Inches(0.15), y+Inches(0.05), Inches(3.2), Inches(0.3),
       f"{name}:  {method}", sz=12, c=color, b=True)
    tb(s, x+Inches(0.15), y+Inches(0.35), Inches(1.0), Inches(0.25), f"⏱ {latency}", sz=10, c=YELLOW)
    tb(s, x+Inches(1.2), y+Inches(0.35), Inches(2.1), Inches(0.25), "渐进升级", sz=9, c=MID_GRAY)
    tb(s, x+Inches(0.15), y+Inches(0.65), Inches(3.2), Inches(0.55), desc, sz=10, c=LIGHT_GRAY)
    if i < 2:
        tb(s, x+Inches(3.5), y+Inches(0.35), Inches(0.3), Inches(0.3), "→", sz=14, c=MID_GRAY, a=PP_ALIGN.CENTER)

# 下：Steering 规则
rect(s, Inches(0.8), Inches(4.6), Inches(11.7), Inches(2.5), fill=DARK_CARD, line=RGBColor(0x1E,0x2A,0x50))
tb(s, Inches(1.1), Inches(4.7), Inches(5), Inches(0.35), "⚙️  Steering 规则系统", sz=16, c=ACCENT_PURPLE, b=True)
line(s, Inches(1.1), Inches(5.05), Inches(1.0), ACCENT_PURPLE, Pt(2))
steering_features = [
    ("声明式规则", "Markdown 文件声明行为规则\n无需改代码，非程序员可配置", ACCENT_CYAN),
    ("四种注入模式", "始终 / 文件匹配 / 关键词匹配 / 手动引用\n灵活控制规则生效时机", GREEN),
    ("两级作用域", "用户级 (~/.maclaw/steering/)\n项目级 (.maclaw/steering/)", ORANGE),
    ("热加载", "修改后 30 秒内自动生效\n无需重启应用", ACCENT_PURPLE),
]
for i, (title, desc, color) in enumerate(steering_features):
    col = i%2; row = i//2
    x = Inches(1.0)+col*Inches(5.8); y = Inches(5.3)+row*Inches(0.8)
    tb(s, x, y, Inches(5.5), Inches(0.25), f"▸ {title}", sz=11, c=color, b=True)
    tb(s, x+Inches(0.2), y+Inches(0.25), Inches(5.3), Inches(0.45), desc, sz=9, c=LIGHT_GRAY)
print("P26 智能调度 完成")

# ================================================================
# P27: ag-UI 协议集成 —— 用对话驱动结构化数据管理
# ================================================================
s = new_slide()
section_header(s, "AG-UI PROTOCOL", "ag-UI — 用对话驱动结构化数据管理",
               "Agent 自然语言交互 → 结构化数据 CRUD，替代传统 MIS/ERP/CRM 前端")

# ===== 顶部标语条 =====
rect(s, Inches(0.8), Inches(2.05), Inches(11.7), Inches(0.65), fill=RGBColor(0x0E,0x1A,0x40), line=ACCENT_BLUE)
mp(s, Inches(1.1), Inches(2.1), Inches(11.2), Inches(0.55), [
    ("一句话录入订单、查询库存、审批报销 —— Agent 理解意图，MIS 结构化存储，ag-UI 实时呈现", 14, ACCENT_CYAN, True),
])

# ===== 左上：传统 MIS vs MaClaw+AG-UI 对比 =====
rect(s, Inches(0.8), Inches(2.9), Inches(5.6), Inches(2.1), fill=DARK_CARD, line=RGBColor(0x1E,0x2A,0x50))
tb(s, Inches(1.1), Inches(3.0), Inches(5), Inches(0.3), "📊 传统 MIS 前端  vs  MaClaw + ag-UI", sz=14, c=ACCENT_CYAN, b=True)
line(s, Inches(1.1), Inches(3.3), Inches(1.0), ACCENT_CYAN, Pt(2))

compare_items = [
    ("❌  传统方式", "表单填写 · 多级菜单 · 培训上岗\n固定报表 · 开发周期长 · 改动需发版", RED),
    ("✅  MaClaw + ag-UI", "自然语言指令 → Agent 自动映射 MIS 操作\n生成式 UI 动态渲染表单 · 零培训上手", GREEN),
]
for i, (title, desc, color) in enumerate(compare_items):
    y = Inches(3.5) + i * Inches(0.72)
    tb(s, Inches(1.1), y, Inches(5.2), Inches(0.25), title, sz=12, c=color, b=True)
    tb(s, Inches(1.3), y + Inches(0.25), Inches(4.8), Inches(0.45), desc, sz=9, c=LIGHT_GRAY)

# ===== 右上：ag-UI 驱动结构化数据的核心交互流程 =====
rect(s, Inches(6.7), Inches(2.9), Inches(5.8), Inches(2.1), fill=DARK_CARD, line=RGBColor(0x1E,0x2A,0x50))
tb(s, Inches(7.0), Inches(3.0), Inches(5), Inches(0.3), "🔄 对话 → 结构化的完整链路", sz=14, c=ORANGE, b=True)
line(s, Inches(7.0), Inches(3.3), Inches(1.0), ORANGE, Pt(2))

flow_steps = [
    ("🗣️ 用户意图", "\"帮我录入一笔Q4采购订单\n供应商华为，金额15.8万\"", ACCENT_CYAN),
    ("🤖 Agent 解析", "NLU 提取实体 → 匹配 business_action\nresolve_intent → sales.order_upsert", ACCENT_BLUE),
    ("📦 MIS 存储", "Schema 校验 → Record 创建\n关联 person_ref / org_ref / file_ref", GREEN),
    ("🖥️ ag-UI 呈现", "生成式 UI 渲染订单卡片\n共享状态实时同步前端", ORANGE),
]
for i, (step, desc, color) in enumerate(flow_steps):
    x = Inches(7.0) + i * Inches(1.35)
    tb(s, x, Inches(3.5), Inches(1.3), Inches(0.25), step, sz=10, c=color, b=True, a=PP_ALIGN.CENTER)
    tb(s, x, Inches(3.8), Inches(1.3), Inches(0.65), desc, sz=8, c=LIGHT_GRAY, a=PP_ALIGN.CENTER)
    if i < 3:
        tb(s, x + Inches(1.25), Inches(3.6), Inches(0.2), Inches(0.2), "→", sz=12, c=MID_GRAY, a=PP_ALIGN.CENTER)

# ===== 左下：ag-UI 六大能力支撑结构化数据 =====
rect(s, Inches(0.8), Inches(5.2), Inches(5.6), Inches(1.95), fill=DARK_CARD, line=RGBColor(0x1E,0x2A,0x50))
tb(s, Inches(1.1), Inches(5.3), Inches(5), Inches(0.3), "⚡ ag-UI 能力 × 结构化数据场景", sz=14, c=ACCENT_BLUE, b=True)
line(s, Inches(1.1), Inches(5.6), Inches(1.0), ACCENT_BLUE, Pt(2))

agui_data_items = [
    ("生成式 UI", "动态渲染表单/卡片/列表/图表\n替代传统 MIS 固定页面", ORANGE),
    ("共享状态", "前端与 MIS 数据双向同步\n实时反映 Record 变更", ACCENT_PURPLE),
    ("前端工具调用", "Agent 调用前端组件完成复杂输入\n日期选择/文件上传/关联选择", ACCENT_CYAN),
]
for i, (title, desc, color) in enumerate(agui_data_items):
    y = Inches(5.75) + i * Inches(0.48)
    tb(s, Inches(1.1), y, Inches(5.2), Inches(0.2), f"● {title}", sz=10, c=color, b=True)
    tb(s, Inches(2.5), y, Inches(3.7), Inches(0.4), desc, sz=9, c=LIGHT_GRAY)

# ===== 右下：MaClaw 结构化数据产品定位 =====
rect(s, Inches(6.7), Inches(5.2), Inches(5.8), Inches(1.95), fill=DARK_CARD, line=ACCENT_BLUE)
tb(s, Inches(7.0), Inches(5.3), Inches(5), Inches(0.3), "🎯 产品核心定位", sz=14, c=WHITE, b=True)
line(s, Inches(7.0), Inches(5.6), Inches(1.0), ACCENT_BLUE, Pt(2))

mp(s, Inches(7.2), Inches(5.75), Inches(5.2), Inches(1.2), [
    ("MIS = 结构化数据存储引擎（后端能力）", 12, YELLOW, True),
    ("ag-UI = 自然语言驱动的前端界面（交互层）", 12, ACCENT_CYAN, True),
    ("", 6, WHITE, False),
    ("MaClaw = Agent 时代的企业信息系统", 16, WHITE, True),
    ("无需 ERP/CRM 前端开发，对话即系统", 11, LIGHT_GRAY, False),
])

print("P27 ag-UI 结构化数据定位 完成")

# ================================================================
# P28: 技术架构全景与安全隐私 (原P27)
# ================================================================
s = new_slide()
section_header(s, "ARCHITECTURE & SECURITY", "技术架构与安全隐私",
               "五层架构 + 数据本地化 + 企业级安全保障")
# 架构层
arch_layers = [
    ("交互层", ACCENT_CYAN, "Wails Desktop · IM 通道 (飞书/微信/QQ) · Bubble Tea TUI · REST API"),
    ("Agent 服务层", ACCENT_BLUE, "IMMessageHandler · 三层UIC · 19种工作流引擎 · 工具路由 · Steering 规则 · AgentInstance"),
    ("核心能力层", ACCENT_PURPLE, "memory.Store · knowledge.Store · structureddata · ExperienceDistiller · SkillHub"),
    ("执行层", GREEN, "SSH 远程管理 · 浏览器CDP · 桌面GUI (Accessibility+YOLO) · 文件系统 · 定时任务"),
    ("存储层", ORANGE, "SQLite WAL · FTS5 · PostgreSQL (Phase2) · 向量索引 · 语义图谱 · ~/.maclaw/"),
]
for i, (name, color, items) in enumerate(arch_layers):
    y = Inches(2.0)+i*Inches(0.75)
    rect(s, Inches(0.8), y, Inches(1.4), Inches(0.6), fill=color, st=MSO_SHAPE.ROUNDED_RECTANGLE)
    tb(s, Inches(0.8), y+Inches(0.1), Inches(1.4), Inches(0.35), name, sz=12, c=WHITE, b=True, a=PP_ALIGN.CENTER)
    rect(s, Inches(2.4), y, Inches(10.1), Inches(0.6), fill=DARK_CARD, line=RGBColor(0x1E,0x2A,0x50), st=MSO_SHAPE.ROUNDED_RECTANGLE)
    tb(s, Inches(2.6), y+Inches(0.1), Inches(9.7), Inches(0.4), items, sz=10, c=LIGHT_GRAY)

# 安全隐私
rect(s, Inches(0.8), Inches(5.9), Inches(11.7), Inches(1.2), fill=DARK_CARD, line=RGBColor(0x1E,0x2A,0x50))
security = [
    ("🔒", "数据本地化", "所有数据存储在本地 ~/.maclaw/\n不上传任何用户数据到云端"),
    ("🛡️", "API Key 安全", "LLM API 调用走官方接口\nToken 不经第三方"),
    ("🔐", "权限隔离", "多租户 Owner 级隔离\nAPI Key 策略控制"),
    ("📋", "审计追踪", "全量操作日志\n变更历史可追溯"),
]
for i, (icon, title, desc) in enumerate(security):
    x = Inches(1.0)+i*Inches(2.9)
    tb(s, x, Inches(6.0), Inches(2.7), Inches(0.3), f"{icon}  {title}", sz=12, c=ACCENT_CYAN, b=True)
    tb(s, x, Inches(6.3), Inches(2.7), Inches(0.6), desc, sz=9, c=LIGHT_GRAY)

# 底部技术栈
rect(s, Inches(0.8), Inches(7.15), Inches(11.7), Inches(0.3), fill=DARKER_CARD, st=MSO_SHAPE.ROUNDED_RECTANGLE)
tb(s, Inches(0.8), Inches(7.17), Inches(11.7), Inches(0.25),
   "Go 1.25  |  Wails v2  |  React 18  |  TypeScript  |  SQLite (WAL)  |  WebSocket  |  Bubble Tea  |  WASM (wazero)",
   sz=9, c=MID_GRAY, a=PP_ALIGN.CENTER)
print("P28 架构安全 完成")

# ================================================================
# P29: 竞品对比（重点）(原P28)
# ================================================================
s = new_slide()
section_header(s, "COMPETITIVE COMPARISON", "与新一代智能体的全面对比",
               "MaClaw vs Hermes vs Macaron AI vs 传统 Agent 方案")
# 对比表格
headers = ["能力维度", "MaClaw", "Hermes", "Macaron AI", "传统 Agent"]
col_w = [Inches(2.0), Inches(2.8), Inches(2.4), Inches(2.4), Inches(2.4)]
row_h = Inches(0.5)
# 表头
cx = Inches(0.4); cy = Inches(2.0)
for j, (h, w) in enumerate(zip(headers, col_w)):
    fill = RGBColor(0x0C,0x10,0x2E) if j > 0 else DARK_CARD
    rect(s, cx, cy, w, row_h, fill=fill, line=RGBColor(0x1E,0x2A,0x50), st=MSO_SHAPE.RECTANGLE)
    tb(s, cx, cy+Inches(0.08), w, row_h-Inches(0.1), h, sz=12,
       c=ACCENT_BLUE if j==1 else LIGHT_GRAY, b=True, a=PP_ALIGN.CENTER)
    cx += w
rows = [
    ["长期记忆", "✅ 三层架构 · 7类型 · Ebbinghaus", "⚠️ 短期上下文为主", "⚠️ 会话级记忆", "❌ 无状态"],
    ["知识库", "✅ 12种来源 · Card/Fact/Node", "⚠️ RAG 基础检索", "⚠️ 简单文档上传", "❌ 无知识库"],
    ["自进化", "✅ Combee蒸馏 · 缺口检测", "❌ 无自进化", "❌ 无自进化", "❌ 无进化"],
    ["SSH远程", "✅ 多会话 · 后台任务 · SFTP", "❌ 无SSH能力", "❌ 无SSH", "❌ 无SSH"],
    ["GUI自动化", "✅ Accessibility + YOLO", "❌ 无GUI能力", "❌ 无GUI", "❌ 无GUI"],
    ["浏览器", "✅ CDP · 录制回放 · OCR", "✅ Web搜索为主", "⚠️ 基础浏览", "⚠️ API调用"],
    ["结构化数据", "✅ MIS内置 · 审计 · 质量管控", "❌ 无MIS", "❌ 无MIS", "❌ 需外部系统"],
    ["工作流", "✅ 19种模板 · 质量闭环", "⚠️ 单轮对话", "⚠️ 基础对话", "⚠️ 简单Chain"],
    ["部署形态", "✅ 桌面/IM/TUI/REST", "☁️ 云端Web", "☁️ 云端Web", "☁️ API服务"],
    ["数据安全", "✅ 本地化 · 开源自托管", "⚠️ 云端处理", "⚠️ 云端处理", "⚠️ 云端处理"],
]
colors_map = [MID_GRAY, GREEN, YELLOW, YELLOW, MID_GRAY]
for i, row in enumerate(rows):
    cx = Inches(0.4); y = cy + (i+1)*row_h
    bg_c = DARK_CARD if i%2==0 else DARKER_CARD
    for j, (cell, w) in enumerate(zip(row, col_w)):
        fc = bg_c if j!=1 else RGBColor(0x0E,0x1A,0x40)
        rect(s, cx, y, w, row_h, fill=fc, line=RGBColor(0x1E,0x2A,0x50), st=MSO_SHAPE.RECTANGLE)
        cc2 = GREEN if cell.startswith("✅") else (YELLOW if cell.startswith("⚠️") else (RED if cell.startswith("❌") else WHITE))
        tb(s, cx+Inches(0.08), y+Inches(0.05), w-Inches(0.16), row_h-Inches(0.1),
           cell, sz=9, c=cc2, b=(j==0), a=PP_ALIGN.CENTER if j>0 else PP_ALIGN.LEFT)
        cx += w
print("P29 竞品对比 完成")

# ================================================================
# P30: 应用场景与快速开始 (原P29)
# ================================================================
s = new_slide()
section_header(s, "SCENARIOS & QUICK START", "应用场景与快速上手",
               "6大行业场景覆盖，4步即可开始使用")
# 上方：场景
scenarios = [
    ("📝", "商业与战略", ACCENT_BLUE, "商业计划书 · 竞品分析 · 项目提案 · 招投标"),
    ("📊", "研发与技术", GREEN, "软件开发全流程 · PRD · 技术文档 · 代码审查"),
    ("⚖️", "合规与法务", ORANGE, "合同审查 · 尽职调查 · 合规审计 · 风险评估"),
    ("🔬", "研究与学术", ACCENT_PURPLE, "文献综述 · 研究报告 · 实验设计 · 专利分析"),
    ("🖥️", "IT 运维", YELLOW, "SSH远程管理 · 日志分析 · 自动化部署 · 安全巡检"),
    ("📊", "数据分析", ACCENT_CYAN, "多源数据汇总 · 经营分析 · KPI监控 · 趋势预测"),
]
scw = Inches(1.85); sch = Inches(1.2)
for i, (icon, title, color, desc) in enumerate(scenarios):
    x = Inches(0.8)+i*(scw+Inches(0.1)); y = Inches(2.0)
    rect(s, x, y, scw, sch, fill=DARK_CARD, line=color)
    tb(s, x, y+Inches(0.08), scw, Inches(0.3), f"{icon}  {title}", sz=11, c=color, b=True, a=PP_ALIGN.CENTER)
    line(s, x+Inches(0.3), y+Inches(0.4), scw-Inches(0.6), color, Pt(1))
    tb(s, x+Inches(0.1), y+Inches(0.5), scw-Inches(0.2), Inches(0.6), desc, sz=8, c=MID_GRAY, a=PP_ALIGN.CENTER)

# 下方：快速开始 4 步
steps = [
    ("01", "安装", "下载 MaClaw 安装包\n一键安装桌面应用", ACCENT_BLUE),
    ("02", "配置 API", "配置 LLM API Key\n支持 OpenAI / Claude / DeepSeek", GREEN),
    ("03", "开始对话", "打开桌面或连接 IM\n自然语言描述需求", ACCENT_PURPLE),
    ("04", "获得成果", "工作流自动执行\n阶段产出等你确认", ORANGE),
]
stw = Inches(2.8); sth = Inches(3.0)
for i, (num, title, desc, color) in enumerate(steps):
    x = Inches(0.8)+i*(stw+Inches(0.2)); y = Inches(3.5)
    rect(s, x, y, stw, sth, fill=DARK_CARD, line=color)
    # 序号圆
    rect(s, x+stw/2-Inches(0.35), y+Inches(0.25), Inches(0.7), Inches(0.7),
         fill=color, st=MSO_SHAPE.OVAL)
    tb(s, x+stw/2-Inches(0.35), y+Inches(0.33), Inches(0.7), Inches(0.5),
       num, sz=24, c=WHITE, b=True, a=PP_ALIGN.CENTER)
    tb(s, x+Inches(0.2), y+Inches(1.1), stw-Inches(0.4), Inches(0.4),
       title, sz=20, c=color, b=True, a=PP_ALIGN.CENTER)
    line(s, x+Inches(0.5), y+Inches(1.55), stw-Inches(1.0), color, Pt(2))
    tb(s, x+Inches(0.3), y+Inches(1.75), stw-Inches(0.6), Inches(1.0),
       desc, sz=13, c=LIGHT_GRAY, a=PP_ALIGN.CENTER)
    if i < 3:
        tb(s, x+stw, y+sth/2-Inches(0.15), Inches(0.2), Inches(0.3), "→", sz=16, c=ACCENT_BLUE, a=PP_ALIGN.CENTER)
print("P30 场景开始 完成")

# ================================================================
# P31: 结尾——核心理念 (原P30)
# ================================================================
s = new_slide()
bg(s, RGBColor(0x06,0x09,0x1C))
line(s, Inches(0), Inches(0), Inches(13.333), ACCENT_BLUE, Pt(4))

# 中心区域
tb(s, Inches(1), Inches(0.8), Inches(11), Inches(0.6),
   "MaClaw  码卡龙", sz=52, c=WHITE, b=True, a=PP_ALIGN.CENTER)
line(s, Inches(5.0), Inches(1.6), Inches(3.3), ACCENT_CYAN, Pt(3))

tb(s, Inches(1.5), Inches(2.0), Inches(10), Inches(0.8),
   "不只是聊天，而是替你干活", sz=36, c=ACCENT_CYAN, a=PP_ALIGN.CENTER)

# 三大理念
ideas = [
    ("🧠", "能记", "内中外脑三层记忆\n永不遗忘你的偏好与知识", ACCENT_BLUE),
    ("🧬", "能学", "Combee 经验蒸馏\n越用越聪明，能力持续进化", GREEN),
    ("🖥️", "能做", "SSH/GUI 全域执行\n真正替你完成实际工作", ORANGE),
]
iw = Inches(3.3); ih = Inches(1.6)
for i, (icon, title, desc, color) in enumerate(ideas):
    x = Inches((13.333-3*iw.inches-2*0.3)/2)+Inches(i*(iw.inches+0.3)); y = Inches(3.2)
    rect(s, x, y, iw, ih, fill=DARK_CARD, line=color)
    tb(s, x, y+Inches(0.1), iw, Inches(0.5), f"{icon}  {title}", sz=26, c=color, b=True, a=PP_ALIGN.CENTER)
    line(s, x+Inches(0.5), y+Inches(0.6), iw-Inches(1.0), color, Pt(2))
    tb(s, x+Inches(0.2), y+Inches(0.7), iw-Inches(0.4), Inches(0.7), desc, sz=12, c=LIGHT_GRAY, a=PP_ALIGN.CENTER)

# 特性标签
tags = ["开源", "自托管", "数据本地化", "多LLM支持", "零配置上手", "企业级安全"]
tw2 = Inches(1.3)
total_t = len(tags)*tw2.inches + (len(tags)-1)*0.15
stx = Inches((13.333-total_t)/2)
for i, tag in enumerate(tags):
    x = stx+Inches(i*(tw2.inches+0.15))
    rect(s, x, Inches(5.2), Inches(tw2), Inches(0.4), line=ACCENT_BLUE)
    tb(s, x, Inches(5.22), Inches(tw2), Inches(0.36), tag, sz=11, c=ACCENT_CYAN, a=PP_ALIGN.CENTER)

# 关键数字
nums = [
    ("19", "种工作流模板"),
    ("7", "种记忆类型"),
    ("197+", "动态工具"),
    ("6", "大行业场景"),
    ("4", "种交互形态"),
]
nw = Inches(2.0); nh = Inches(0.8)
for i, (num, label) in enumerate(nums):
    x = Inches((13.333-5*nw.inches-4*0.15)/2)+Inches(i*(nw.inches+0.15)); y = Inches(5.9)
    rect(s, x, y, Inches(nw), nh, fill=DARK_CARD, line=RGBColor(0x1E,0x2A,0x50))
    tb(s, x, y+Inches(0.05), Inches(0.8), Inches(0.3), num, sz=22, c=ACCENT_CYAN, b=True, a=PP_ALIGN.CENTER)
    tb(s, x+Inches(0.8), y+Inches(0.1), Inches(1.1), Inches(0.25), label, sz=10, c=LIGHT_GRAY)

# 底部
tb(s, Inches(0), Inches(6.9), Inches(13.333), Inches(0.4),
   "github.com/RapidAI/CodeClaw", sz=14, c=MID_GRAY, a=PP_ALIGN.CENTER)
print("P31 结尾 完成")

# ========== 保存 ==========
output_dir = os.path.dirname(os.path.abspath(__file__))
output_path = os.path.join(os.path.dirname(output_dir), "MaClaw_产品介绍_CN_30p.pptx")
prs.save(output_path)
print(f"\n✅ PPT 已保存: {output_path}")
print(f"共 {len(prs.slides)} 页")
