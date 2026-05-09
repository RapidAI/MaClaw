# -*- coding: utf-8 -*-
"""
MaClaw（码卡龙）产品介绍 PPT 生成脚本
面向用户，突出功能特点、核心优势、竞品对比
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 创建演示文稿 (16:9)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ========== 配色方案 ==========
PRIMARY = RGBColor(0x0A, 0x0E, 0x27)
ACCENT_BLUE = RGBColor(0x00, 0x7B, 0xFF)
ACCENT_CYAN = RGBColor(0x00, 0xD4, 0xFF)
ACCENT_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xB8, 0xCC)
MID_GRAY = RGBColor(0x6B, 0x7B, 0x94)
DARK_CARD = RGBColor(0x12, 0x18, 0x38)
DARKER_CARD = RGBColor(0x0E, 0x13, 0x30)
ORANGE = RGBColor(0xFF, 0x6B, 0x35)
GREEN = RGBColor(0x00, 0xC8, 0x53)
RED = RGBColor(0xFF, 0x44, 0x44)
YELLOW = RGBColor(0xFF, 0xD6, 0x00)


def set_slide_bg(slide, color=PRIMARY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None,
              shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.shadow.inherit = False
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE,
                 bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multi_para(slide, left, top, width, height, lines, font_size=14, color=WHITE,
                   alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei', line_spacing=1.3):
    """添加多段落文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if isinstance(line, tuple):
            p.text = line[0]
            p.font.size = Pt(line[1]) if len(line) > 1 else Pt(font_size)
            p.font.color.rgb = line[2] if len(line) > 2 else color
            p.font.bold = line[3] if len(line) > 3 else False
        else:
            p.text = line
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(font_size * 0.3)
    return txBox


def add_decorative_line(slide, left, top, width, color=ACCENT_BLUE, height=Pt(3)):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    line.shadow.inherit = False
    return line


def add_icon_circle(slide, left, top, size, color=ACCENT_BLUE):
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    circle.shadow.inherit = False
    return circle


def add_card(slide, x, y, w, h, title, desc_lines, icon_text="",
             title_color=ACCENT_CYAN, accent=ACCENT_BLUE):
    """统一卡片组件"""
    card = add_shape(slide, x, y, w, h, fill_color=DARK_CARD, line_color=RGBColor(0x1E, 0x2A, 0x50))
    if icon_text:
        add_text_box(slide, x + Inches(0.2), y + Inches(0.15), w - Inches(0.4), Inches(0.4),
                     icon_text, font_size=20, color=accent)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.55), w - Inches(0.4), Inches(0.4),
                 title, font_size=16, color=title_color, bold=True)
    add_decorative_line(slide, x + Inches(0.2), y + Inches(0.95), Inches(0.6), accent, Pt(2))
    add_multi_para(slide, x + Inches(0.2), y + Inches(1.1), w - Inches(0.4), h - Inches(1.3),
                   desc_lines, font_size=11, color=LIGHT_GRAY)
    return card


# ================================================================
# 第1页：封面
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

# 装饰性背景元素
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5),
          fill_color=RGBColor(0x06, 0x09, 0x1C))

# 顶部装饰线
add_decorative_line(slide, Inches(0), Inches(0), Inches(13.333), ACCENT_BLUE, Pt(4))

# 中央区域
add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.6),
             "MaClaw  码卡龙", font_size=56, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_decorative_line(slide, Inches(5), Inches(2.4), Inches(3.3), ACCENT_CYAN, Pt(3))

add_text_box(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(1.0),
             "通用可自进化智能体平台", font_size=36, color=ACCENT_CYAN,
             bold=False, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(3.9), Inches(9), Inches(1.0),
             "你的个人数智工作伙伴 — 不只是聊天，而是替你干活",
             font_size=22, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# 底部特性标签
tags = ["长期记忆", "知识库", "自进化", "SSH远程", "GUI自动化", "19种工作流"]
tag_w = Inches(1.7)
tag_h = Inches(0.5)
total_w = len(tags) * tag_w.inches + (len(tags) - 1) * 0.2
start_x = Inches((13.333 - total_w) / 2)

for i, tag in enumerate(tags):
    x = start_x + Inches(i * (tag_w.inches + 0.2))
    add_shape(slide, x, Inches(5.3), tag_w, tag_h,
              fill_color=None, line_color=ACCENT_BLUE, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text_box(slide, x, Inches(5.32), tag_w, tag_h,
                 tag, font_size=13, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)

# 底部
add_text_box(slide, Inches(0), Inches(6.8), Inches(13.333), Inches(0.5),
             "Wails + Go + React  |  桌面 · IM · TUI 多形态  |  开箱即用",
             font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

print("封面完成")

# ================================================================
# 第2页：产品定位——不只是聊天机器人
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
             "POSITIONING", font_size=12, color=ACCENT_CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.8),
             "MaClaw 能做什么", font_size=36, color=WHITE, bold=True)
add_decorative_line(slide, Inches(0.8), Inches(1.55), Inches(2), ACCENT_BLUE)

# 左侧大标语
big_quote = add_shape(slide, Inches(0.8), Inches(2.1), Inches(5.5), Inches(4.8),
                       fill_color=DARK_CARD, line_color=RGBColor(0x1E, 0x2A, 0x50))
add_multi_para(slide, Inches(1.3), Inches(2.4), Inches(4.5), Inches(4.2),
    [
        ("你说想法，它出成果", 28, ACCENT_CYAN, True),
        ("", 10, WHITE, False),
        ("MaClaw 内置 19 种结构化工作流模板，覆盖商业与战略、研究与分析、合规与尽调、学术写作、内容创作、产品与技术全场景。", 14, LIGHT_GRAY, False),
        ("", 8, WHITE, False),
        ("每个工作流遵循「需求确认 → 方案设计 → 分步执行」的质量闭环，每个阶段产出文档后等待你审阅、修改、确认。", 14, LIGHT_GRAY, False),
        ("", 8, WHITE, False),
        ("你只需要说出想法，MaClaw 从需求梳理到成果交付全程陪你走完。", 14, ACCENT_BLUE, True),
    ])

# 右侧能力矩阵
domains = [
    ("商业与战略", "商业计划书 · 竞品分析 · 项目提案\n创新方案 · 招投标文件", ACCENT_BLUE),
    ("研究与分析", "文献综述 · 研究报告 · 实验设计\n专利分析", ACCENT_PURPLE),
    ("合规与尽调", "合同审查 · 尽职调查 · 合规审计", GREEN),
    ("学术写作", "基金申请书 · 论文写作", ORANGE),
    ("内容创作", "PPT 设计 · 活动策划", ACCENT_CYAN),
    ("产品与技术", "产品设计(PRD) · 软件测试 · 软件开发", YELLOW),
]

for i, (title, desc, color) in enumerate(domains):
    row = i // 2
    col = i % 2
    x = Inches(6.8) + col * Inches(3.0)
    y = Inches(2.1) + row * Inches(1.6)
    card = add_shape(slide, x, y, Inches(2.8), Inches(1.4), fill_color=DARK_CARD,
                     line_color=RGBColor(0x1E, 0x2A, 0x50))
    add_text_box(slide, x + Inches(0.15), y + Inches(0.1), Inches(2.5), Inches(0.35),
                 title, font_size=14, color=color, bold=True)
    add_decorative_line(slide, x + Inches(0.15), y + Inches(0.45), Inches(0.5), color, Pt(2))
    add_text_box(slide, x + Inches(0.15), y + Inches(0.55), Inches(2.5), Inches(0.8),
                 desc, font_size=10, color=LIGHT_GRAY)

print("产品定位页完成")

# ================================================================
# 第3页：核心能力总览
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
             "CORE CAPABILITIES", font_size=12, color=ACCENT_CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.8),
             "六大核心能力，构建完整智能体", font_size=36, color=WHITE, bold=True)
add_decorative_line(slide, Inches(0.8), Inches(1.55), Inches(2), ACCENT_BLUE)

capabilities = [
    ("🧠", "长期记忆系统", [
        "内中外脑三层记忆架构",
        "Ebbinghaus 遗忘曲线 + 自适应衰减",
        "语义图谱 · 实体索引 · 话题聚类",
        "时序记忆树 · 渐进式检索",
        "Claude 四类型分类法",
    ], ACCENT_BLUE),
    ("📚", "中外脑知识库", [
        "支持 PDF/Word/Excel/Markdown/URL",
        "文档结构化解析，非简单切块",
        "知识卡片蒸馏 + 实体关系抽取",
        "目录批量导入 · URL 安全抓取",
        "分层召回：卡片 → 事实 → 原文节点",
    ], ACCENT_PURPLE),
    ("🧬", "自我进化能力", [
        "能力缺口自动检测与补全",
        "经验蒸馏 → 可复用技能沉淀",
        "技能自修复 · Nudge 技能推荐",
        "工具路由学习与自适应",
        "Combee 并行经验聚合",
    ], GREEN),
    ("🖥️", "SSH/GUI 远程操控", [
        "SSH 多会话 · 密码/密钥/Agent 认证",
        "后台任务 · 文件传输 · sudo 管理",
        "浏览器自动化 (CDP)",
        "桌面 GUI 自动化 (Accessibility+YOLO)",
        "流程录制与回放 · 定时触发",
    ], ORANGE),
    ("🔗", "多通道协作", [
        "桌面 AI 面板 / IM 消息 / TUI 终端",
        "飞书 · 微信 · QQ · 钉钉 · 企微",
        "Hub 消息代理 · A2A 群聊协作",
        "Swarm 多智能体协同",
        "AgentNet P2P 知识网络",
    ], ACCENT_CYAN),
    ("🏗️", "企业级结构化数据", [
        "租户隔离 · 多工作域数据管理",
        "Schema 渐进演化 · 字段版本追踪",
        "API Key 权限策略 · 访问审计",
        "外部连接器 · 数据导入导出",
        "SQLite 本地 / PostgreSQL 团队部署",
    ], YELLOW),
]

card_w = Inches(3.8)
card_h = Inches(2.2)
gap_x = Inches(0.3)
gap_y = Inches(0.3)
start_x = Inches(0.8)
start_y = Inches(2.0)

for i, (icon, title, lines, color) in enumerate(capabilities):
    col = i % 3
    row = i // 3
    x = start_x + col * (card_w + gap_x)
    y = start_y + row * (card_h + gap_y)

    card = add_shape(slide, x, y, card_w, card_h, fill_color=DARK_CARD,
                     line_color=RGBColor(0x1E, 0x2A, 0x50))
    add_text_box(slide, x + Inches(0.15), y + Inches(0.1), card_w - Inches(0.3), Inches(0.4),
                 f"{icon}  {title}", font_size=16, color=color, bold=True)
    add_decorative_line(slide, x + Inches(0.15), y + Inches(0.5), Inches(0.8), color, Pt(2))
    add_multi_para(slide, x + Inches(0.15), y + Inches(0.6), card_w - Inches(0.3), card_h - Inches(0.7),
                   [f"• {l}" for l in lines], font_size=10, color=LIGHT_GRAY)

print("核心能力页完成")

# ================================================================
# 第4页：内中外脑——三层记忆架构（重点页）
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
             "MEMORY ARCHITECTURE", font_size=12, color=ACCENT_CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.8),
             "内中外脑——它记得你的一切", font_size=36, color=WHITE, bold=True)
add_decorative_line(slide, Inches(0.8), Inches(1.55), Inches(2), ACCENT_BLUE)

# 三层架构图
brains = [
    ("内  脑", "corelib/memory", ACCENT_BLUE,
     ["用户长期偏好与身份", "项目知识与决策记录", "会话摘要与任务检查点",
      "高价值结论与反馈", "Claude 四类型: user/feedback/project/reference"],
     "持久化 · Ebbinghaus 衰减\n向量索引 + BM25 + 语义图谱"),
    ("中  脑", "智能调度层", ACCENT_PURPLE,
     ["识别当前话题与用户意图", "决定什么时候存、存到哪", "分层召回策略 (L1-L5)",
      "查询复杂度评估 (simple/hybrid/complex)", "冲突检测与语义去重"],
     "上下文感知 · 实时调度\n时序记忆树 · 5级渐进检索"),
    ("外  脑", "corelib/knowledge", GREEN,
     ["大型 PDF/Word/Excel/网页", "文档结构化节点 (DocumentNode)", "知识卡片 (KnowledgeCard)",
      "实体关系事实 (KnowledgeFact)", "目录批量导入 · URL 安全抓取"],
     "海量存储 · SQLite FTS5\n结构化优先，避免原始 RAG"),
]

brain_w = Inches(3.7)
brain_h = Inches(4.2)
gap = Inches(0.25)

for i, (title, pkg, color, features, note) in enumerate(brains):
    x = Inches(0.7) + i * (brain_w + gap)
    y = Inches(2.0)

    # 卡片
    card = add_shape(slide, x, y, brain_w, brain_h, fill_color=DARK_CARD,
                     line_color=color)

    # 标题
    add_text_box(slide, x + Inches(0.2), y + Inches(0.15), brain_w - Inches(0.4), Inches(0.4),
                 title, font_size=24, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.55), brain_w - Inches(0.4), Inches(0.3),
                 pkg, font_size=10, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
    add_decorative_line(slide, x + Inches(0.5), y + Inches(0.85), brain_w - Inches(1.0), color, Pt(2))

    # 特性列表
    add_multi_para(slide, x + Inches(0.2), y + Inches(1.0), brain_w - Inches(0.4), Inches(2.2),
                   [f"✦  {f}" for f in features], font_size=11, color=LIGHT_GRAY)

    # 底部说明
    add_shape(slide, x + Inches(0.15), y + brain_h - Inches(0.65), brain_w - Inches(0.3), Inches(0.55),
              fill_color=DARKER_CARD)
    add_text_box(slide, x + Inches(0.2), y + brain_h - Inches(0.6), brain_w - Inches(0.4), Inches(0.5),
                 note, font_size=9, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    # 连接箭头（在前两个之间）
    if i < 2:
        arrow_x = x + brain_w
        add_text_box(slide, arrow_x, y + brain_h / 2 - Inches(0.2), gap, Inches(0.4),
                     "⟷", font_size=18, color=color, alignment=PP_ALIGN.CENTER)

# 底部亮点条
bar = add_shape(slide, Inches(0.7), Inches(6.5), Inches(11.9), Inches(0.7),
                fill_color=DARKER_CARD, line_color=RGBColor(0x1E, 0x2A, 0x50))
highlights = "🧠 语义图谱多跳推理  |  📊 BM25+向量混合检索  |  🔄 Dream Cycle 自修复  |  📌 钉选保护机制  |  🔒 多租户 Owner 隔离"
add_text_box(slide, Inches(1.0), Inches(6.55), Inches(11.3), Inches(0.6),
             highlights, font_size=12, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)

print("内中外脑页完成")

# ================================================================
# 第5页：长期记忆系统深度解析
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
             "LONG-TERM MEMORY", font_size=12, color=ACCENT_CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.8),
             "长期记忆——从存储到智能调度", font_size=36, color=WHITE, bold=True)
add_decorative_line(slide, Inches(0.8), Inches(1.55), Inches(2), ACCENT_BLUE)

# 左侧：记忆类型分类
left_card = add_shape(slide, Inches(0.8), Inches(2.0), Inches(5.8), Inches(5.0),
                       fill_color=DARK_CARD, line_color=RGBColor(0x1E, 0x2A, 0x50))
add_text_box(slide, Inches(1.1), Inches(2.1), Inches(5.2), Inches(0.4),
             "记忆分类与层级", font_size=18, color=ACCENT_BLUE, bold=True)
add_decorative_line(slide, Inches(1.1), Inches(2.5), Inches(1.2), ACCENT_BLUE, Pt(2))

mem_types = [
    ("自我身份 (self_identity)", "受保护，永不驱逐", "🛡️"),
    ("用户事实 (user_fact)", "用户角色、目标、知识", "👤"),
    ("偏好指令 (preference/instruction)", "行为规则与反馈修正", "⚙️"),
    ("项目知识 (project_knowledge)", "决策、期限、架构", "📁"),
    ("会话摘要 (conversation_summary)", "对话压缩与上下文保留", "💬"),
    ("任务产物 (task_artifact)", "工作流阶段输出", "📋"),
    ("经验沉淀 (experience)", "工具路由 · 技能模式", "🧬"),
]

for i, (name, desc, icon) in enumerate(mem_types):
    y = Inches(2.75) + i * Inches(0.55)
    add_text_box(slide, Inches(1.1), y, Inches(0.35), Inches(0.35),
                 icon, font_size=14, color=ACCENT_CYAN)
    add_text_box(slide, Inches(1.5), y, Inches(2.5), Inches(0.3),
                 name, font_size=11, color=WHITE, bold=True)
    add_text_box(slide, Inches(4.0), y, Inches(2.5), Inches(0.3),
                 desc, font_size=10, color=MID_GRAY)

# 右侧：技术特性
right_card = add_shape(slide, Inches(7.0), Inches(2.0), Inches(5.5), Inches(5.0),
                        fill_color=DARK_CARD, line_color=RGBColor(0x1E, 0x2A, 0x50))
add_text_box(slide, Inches(7.3), Inches(2.1), Inches(5.0), Inches(0.4),
             "核心技术特性", font_size=18, color=ACCENT_PURPLE, bold=True)
add_decorative_line(slide, Inches(7.3), Inches(2.5), Inches(1.2), ACCENT_PURPLE, Pt(2))

tech_features = [
    ("语义图谱 (SemanticGraph)", "Subject-Predicate-Object 三元组\n实体索引 · 多跳推理 · 极性竞争检测", ACCENT_BLUE),
    ("时序记忆树 (TemporalTree)", "Segment → Session → Day → Week → Profile\n5 级层次化压缩与检索", GREEN),
    ("遗忘曲线 + Dream Cycle", "Ebbinghaus 强度衰减 (半衰期 9.6 天)\n后台自修复：过期检测 · 关系发现 · Hash 补填", ORANGE),
    ("语义去重 & 冲突检测", "LLM 精确去重 · 矛盾事实自动标记\n四操作更新: ADD/DELETE/UPDATE/MERGE", ACCENT_CYAN),
    ("查询复杂度自适应", "Simple(L1-L3) / Hybrid(L1-L4) / Complex(L1-L5)\nRecallGating 防止过度召回", YELLOW),
]

for i, (title, desc, color) in enumerate(tech_features):
    y = Inches(2.75) + i * Inches(0.85)
    add_text_box(slide, Inches(7.3), y, Inches(4.8), Inches(0.3),
                 f"▸ {title}", font_size=12, color=color, bold=True)
    add_text_box(slide, Inches(7.5), y + Inches(0.25), Inches(4.6), Inches(0.5),
                 desc, font_size=9, color=LIGHT_GRAY)

print("长期记忆深度页完成")

# ================================================================
# 第6页：知识库——外脑能力
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
             "KNOWLEDGE BASE", font_size=12, color=ACCENT_CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.8),
             "中外脑知识库——让 AI 真正读懂你的资料", font_size=36, color=WHITE, bold=True)
add_decorative_line(slide, Inches(0.8), Inches(1.55), Inches(2), ACCENT_BLUE)

# 上方：来源类型 + 处理流水线
sources = [
    ("📄", "PDF", "学术论文\n技术文档"),
    ("📝", "DOCX", "合同\n报告"),
    ("📊", "XLSX", "数据表\n清单"),
    ("📋", "Markdown", "设计文档\nREADME"),
    ("🌐", "URL", "网页\n博客"),
    ("📁", "目录批量", "整站导入\n知识库"),
    ("🎤", "语音", "ASR 转写\n声纹"),
]

src_w = Inches(1.55)
src_h = Inches(1.0)
for i, (icon, name, desc) in enumerate(sources):
    x = Inches(0.8) + i * (src_w + Inches(0.1))
    card = add_shape(slide, x, Inches(2.0), src_w, src_h,
                     fill_color=DARK_CARD, line_color=RGBColor(0x1E, 0x2A, 0x50))
    add_text_box(slide, x, Inches(2.0), src_w, Inches(0.35),
                 f"{icon} {name}", font_size=12, color=ACCENT_BLUE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(2.4), src_w, Inches(0.55),
                 desc, font_size=9, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# 箭头
add_text_box(slide, Inches(0.8), Inches(3.15), Inches(11.7), Inches(0.35),
             "━━━━━━━━━━━━━━━━━━━  Distill Pipeline  ━━━━━━━━━━━━━━━━━━━",
             font_size=11, color=ACCENT_PURPLE, alignment=PP_ALIGN.CENTER)

# 下半区：三列产出物
outputs = [
    ("KnowledgeCard", "知识卡片", ACCENT_BLUE, [
        "5 种类型: Summary/Concept/Procedure/Reference/FAQ",
        "质量评分 (Quality Score: 0-100)",
        "卡片级别的 Pin 与版本追踪",
        "支持 Batch 批量蒸馏",
        "自动标签与主题分类",
    ]),
    ("KnowledgeFact", "事实三元组", GREEN, [
        "Subject-Predicate-Object 结构",
        "实体关系图谱自动构建",
        "与记忆系统语义图谱联动",
        "冲突检测与精确去重",
        "来源溯源 (SourceLink)",
    ]),
    ("DocumentNode", "文档结构树", ORANGE, [
        "章节 → 段落 → 句子层级",
        "SourceTimelineEvent 时间线",
        "原文片段精确定位",
        "增量更新与版本追踪",
        "分层召回: 卡片→事实→原文",
    ]),
]

out_w = Inches(3.7)
out_h = Inches(3.1)
for i, (cls_name, cn_name, color, features) in enumerate(outputs):
    x = Inches(0.8) + i * (out_w + Inches(0.25))
    y = Inches(3.6)
    card = add_shape(slide, x, y, out_w, out_h, fill_color=DARK_CARD,
                     line_color=color)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.15), out_w - Inches(0.4), Inches(0.35),
                 f"{cn_name}", font_size=16, color=color, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.45), out_w - Inches(0.4), Inches(0.25),
                 cls_name, font_size=9, color=MID_GRAY)
    add_decorative_line(slide, x + Inches(0.2), y + Inches(0.7), Inches(0.8), color, Pt(2))
    add_multi_para(slide, x + Inches(0.2), y + Inches(0.85), out_w - Inches(0.4), out_h - Inches(1.0),
                   [f"•  {f}" for f in features], font_size=10, color=LIGHT_GRAY)

print("知识库页完成")

# ================================================================
# 第7页：企业级结构化数据管理 (MIS)
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
             "STRUCTURED DATA (MIS)", font_size=12, color=ACCENT_CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.8),
             "结构化数据管理——MIS 能力内置", font_size=36, color=WHITE, bold=True)
add_decorative_line(slide, Inches(0.8), Inches(1.55), Inches(2), ACCENT_BLUE)

# 6 个核心概念卡片
mis_items = [
    ("🏢", "多租户 & 工作域", ACCENT_BLUE,
     ["租户隔离 (Tenant) 组织边界",
      "业务域 (Domain): sales/hr/finance/legal",
      "数据集 (Dataset): schema-less 渐进演化",
      "字段定义 (FieldDefinition) 非强制",
      "16 种字段类型含 record_ref 关联"]),
    ("📋", "Record & 版本", GREEN,
     ["灵活 JSON data + 固定元数据",
      "Record Revision 修改追溯",
      "谁修改 · 什么时候 · 改了什么",
      "FTS5 全文搜索 · 字段级索引",
      "标签系统 · 自动归一化"]),
    ("🔐", "权限与审计", ORANGE,
     ["API Key Policy 策略控制",
      "ConnectorContractBinding 绑定",
      "DataEventLog 全量操作审计",
      "外部连接器安全接入",
      "ConnectorHealth 健康监控"]),
    ("🔗", "外部连接器", ACCENT_PURPLE,
     ["ExternalConnector 定义",
      "支持 REST/数据库/文件系统",
      "数据导入导出管道",
      "合约绑定与访问策略",
      "健康检查与自动恢复"]),
    ("📊", "数据质量", YELLOW,
     ["质量评分与检查规则",
      "关系引用完整性校验",
      "敏感字段标记与脱敏",
      "数据变更事件日志",
      "字段类型自动推断"]),
    ("💾", "存储引擎", ACCENT_CYAN,
     ["SQLite WAL 单机/本地",
      "字段级索引 (text/number/time)",
      "PostgreSQL 团队部署 (Phase 2)",
      "租户隔离 Schema 级",
      "FTS5 全文检索引擎"]),
]

mis_w = Inches(3.7)
mis_h = Inches(2.2)
for i, (icon, title, color, lines) in enumerate(mis_items):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + col * (mis_w + Inches(0.3))
    y = Inches(2.0) + row * (mis_h + Inches(0.3))
    add_card(slide, x, y, mis_w, mis_h, title, [f"• {l}" for l in lines],
             icon, title_color=color, accent=color)

print("结构化数据页完成")

# ================================================================
# 第8页：自进化——经验积累与能力成长
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
             "SELF-EVOLUTION", font_size=12, color=ACCENT_CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.8),
             "自进化——越用越聪明的智能体", font_size=36, color=WHITE, bold=True)
add_decorative_line(slide, Inches(0.8), Inches(1.55), Inches(2), ACCENT_BLUE)

# 中间大标语
quote = add_shape(slide, Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.2),
                  fill_color=DARK_CARD, line_color=RGBColor(0x1E, 0x2A, 0x50))
add_text_box(slide, Inches(1.8), Inches(2.1), Inches(9.7), Inches(1.0),
             "\"MaClaw 不是被动执行命令的工具，而是会学习、会反思、会自我改进的工作伙伴\"",
             font_size=20, color=ACCENT_CYAN, bold=True, alignment=PP_ALIGN.CENTER)

# 四个进化机制卡片
evo_items = [
    ("🔍", "能力缺口检测", ACCENT_BLUE,
     ["遇到无法完成的任务时，自动搜索 SkillHub",
      "找到匹配技能并安装，无需人工干预",
      "搜索无果时记录缺口，待后续补全",
      "craft_tool 临时脚本 → 永久技能转化"]),
    ("🧪", "经验蒸馏 (Combee)", ACCENT_PURPLE,
     ["从对话/工作流/Swarm/A2A中提取经验",
      "保留具体路径、命令、错误号等细节",
      "保守学习：单次事件仅生成提示，重复证据才改变权重",
      "保护少数高价值信号，不被通用摘要淹没"]),
    ("🔧", "技能自修复", GREEN,
     ["技能执行失败后 LLM 自动分析错误原因",
      "修补技能定义（步骤、参数、路径等）",
      "Nudge 系统建议封装成功操作序列",
      "修复结果持久化，避免重复出错"]),
    ("🛤️", "工具路由进化", ORANGE,
     ["记录工具调用的完整执行轨迹",
      "task_type · tool_sequence · error_class",
      "retry_count · recovery_tool · final_outcome",
      "高失败率工具自动降权 · 新工具渐进激活"]),
]

evo_w = Inches(5.7)
evo_h = Inches(1.8)
for i, (icon, title, color, lines) in enumerate(evo_items):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + col * (evo_w + Inches(0.3))
    y = Inches(3.5) + row * (evo_h + Inches(0.2))
    add_card(slide, x, y, evo_w, evo_h, title, [f"• {l}" for l in lines],
             icon, title_color=color, accent=color)

print("自进化页完成")

# ================================================================
# 第9页：执行能力——SSH远程 + GUI自动化
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
             "EXECUTION", font_size=12, color=ACCENT_CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.8),
             "执行能力——不仅能想，还能做", font_size=36, color=WHITE, bold=True)
add_decorative_line(slide, Inches(0.8), Inches(1.55), Inches(2), ACCENT_BLUE)

# 三大执行能力卡片
exec_items = [
    ("🖥️", "SSH 远程服务器管理", ACCENT_BLUE, [
        "多会话并发，密码/密钥/Agent 多种认证",
        "同步执行 · 后台任务 · SFTP 文件传输",
        "自动 sudo token 管理 · Shell 响应性检测",
        "连续失败自动清理 · 操作全程审计",
        "支持通过 IM 通道远程操控服务器",
    ], "connect → exec → upload/download → check_task → close"),
    ("🌐", "浏览器自动化 (CDP)", GREEN, [
        "基于 Chrome DevTools Protocol",
        "页面导航 · 元素点击 · 文本输入 · 截图",
        "内容提取 · OCR 集成 (RapidOCR + LLM)",
        "流程录制与回放 · 定时触发",
        "参数化变量 · 后台异步执行",
    ], "导航 → 交互 → 提取 → 验证 → 输出"),
    ("🖥️", "桌面 GUI 自动化", ORANGE, [
        "Accessibility 元素树 (Win/Mac/Linux)",
        "YOLO 视觉检测 (OmniParser V2)",
        "不依赖 Accessibility API 的视觉方案",
        "鼠标键盘操作 · 坐标精确点击",
        "流程录制回放 · 参数化覆盖",
    ], "截图 → 检测元素 → 定位 → 操作 → 验证"),
]

exec_w = Inches(3.7)
exec_h = Inches(4.6)
for i, (icon, title, color, lines, pipeline) in enumerate(exec_items):
    x = Inches(0.8) + i * (exec_w + Inches(0.3))
    y = Inches(2.0)
    card = add_shape(slide, x, y, exec_w, exec_h, fill_color=DARK_CARD,
                     line_color=color)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.15), exec_w - Inches(0.4), Inches(0.4),
                 f"{icon}  {title}", font_size=16, color=color, bold=True)
    add_decorative_line(slide, x + Inches(0.2), y + Inches(0.55), Inches(1.0), color, Pt(2))
    add_multi_para(slide, x + Inches(0.2), y + Inches(0.7), exec_w - Inches(0.4), Inches(2.5),
                   [f"•  {l}" for l in lines], font_size=11, color=LIGHT_GRAY)
    # 底部 pipeline
    add_shape(slide, x + Inches(0.15), y + exec_h - Inches(0.7), exec_w - Inches(0.3), Inches(0.55),
              fill_color=DARKER_CARD)
    add_text_box(slide, x + Inches(0.2), y + exec_h - Inches(0.65), exec_w - Inches(0.4), Inches(0.5),
                 pipeline, font_size=9, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

print("执行能力页完成")

# ================================================================
# 第10页：竞品对比——vs Hermes / Macaron / 传统方案
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
             "COMPETITIVE COMPARISON", font_size=12, color=ACCENT_CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.8),
             "与新一代智能体的对比", font_size=36, color=WHITE, bold=True)
add_decorative_line(slide, Inches(0.8), Inches(1.55), Inches(2), ACCENT_BLUE)

# 对比表格
headers = ["能力维度", "MaClaw", "Hermes", "Macaron AI", "传统 Agent"]
col_widths = [Inches(2.0), Inches(2.8), Inches(2.4), Inches(2.4), Inches(2.4)]
col_colors = [MID_GRAY, ACCENT_BLUE, LIGHT_GRAY, LIGHT_GRAY, MID_GRAY]

# 表头
table_x = Inches(0.4)
table_y = Inches(2.0)
row_h = Inches(0.55)

# 绘制表头
cur_x = table_x
for j, (header, cw, cc) in enumerate(zip(headers, col_widths, col_colors)):
    cell = add_shape(slide, cur_x, table_y, cw, row_h,
                     fill_color=DARK_CARD if j > 0 else RGBColor(0x0C, 0x10, 0x2E),
                     line_color=RGBColor(0x1E, 0x2A, 0x50),
                     shape_type=MSO_SHAPE.RECTANGLE)
    add_text_box(slide, cur_x, table_y + Inches(0.08), cw, row_h - Inches(0.1),
                 header, font_size=13, color=cc if j == 1 else LIGHT_GRAY,
                 bold=True, alignment=PP_ALIGN.CENTER)
    cur_x += cw

# 数据行
rows = [
    ["长期记忆",
     "✅ 三层架构 · 7类型\nEbbinghaus + 语义图谱",
     "⚠️ 短期上下文为主\n无独立记忆模块",
     "⚠️ 会话级记忆\n无跨会话持久化",
     "❌ 无状态\n每次对话从零开始"],
    ["知识库",
     "✅ 12种来源 · 结构化蒸馏\nCard/Fact/Node 三层产出",
     "⚠️ RAG 基础检索\n文档切块模式",
     "⚠️ 简单文档上传\n无结构化提取",
     "❌ 无知识库\n纯 Prompt 驱动"],
    ["自进化",
     "✅ Combee 经验蒸馏\n能力缺口检测 · 技能自修复",
     "❌ 无自进化机制\n固定能力集",
     "❌ 无自进化\n依赖人工配置",
     "❌ 无进化\n工具固定"],
    ["执行能力",
     "✅ SSH + GUI + 浏览器\n流程录制回放 · 多通道协作",
     "✅ Web 搜索为主\n有限的工具调用",
     "⚠️ 基础工具\n无远程操作",
     "⚠️ API 调用\n无操作系统级能力"],
    ["结构化数据",
     "✅ MIS 内置\n多租户 · 审计 · 质量管控",
     "❌ 无结构化数据管理",
     "❌ 无 MIS 能力",
     "❌ 需外部系统"],
    ["工作流",
     "✅ 19 种模板\n质量闭环 · 阶段审阅",
     "⚠️ 单轮对话为主",
     "⚠️ 基础对话",
     "⚠️ 简单 Chain"],
    ["部署形态",
     "✅ 桌面/IM/TUI/REST\n多通道统一体验",
     "☁️ 云端 Web",
     "☁️ 云端 Web",
     "☁️ API 服务"],
]

data_colors = [
    [MID_GRAY, GREEN, YELLOW, YELLOW, RED],
    [MID_GRAY, GREEN, YELLOW, YELLOW, RED],
    [MID_GRAY, GREEN, RED, RED, RED],
    [MID_GRAY, GREEN, YELLOW, YELLOW, YELLOW],
    [MID_GRAY, GREEN, RED, RED, RED],
    [MID_GRAY, GREEN, YELLOW, YELLOW, YELLOW],
    [MID_GRAY, GREEN, RED, RED, RED],
]

for i, (row_data, row_colors) in enumerate(zip(rows, data_colors)):
    cur_x = table_x
    y = table_y + (i + 1) * row_h
    bg = DARK_CARD if i % 2 == 0 else DARKER_CARD
    for j, (cell_text, cw, cc) in enumerate(zip(row_data, col_widths, row_colors)):
        cell = add_shape(slide, cur_x, y, cw, row_h,
                         fill_color=bg if j != 1 else RGBColor(0x0E, 0x1A, 0x40),
                         line_color=RGBColor(0x1E, 0x2A, 0x50),
                         shape_type=MSO_SHAPE.RECTANGLE)
        add_text_box(slide, cur_x + Inches(0.08), y + Inches(0.03), cw - Inches(0.16), row_h - Inches(0.06),
                     cell_text, font_size=9, color=cc if j > 0 else WHITE,
                     bold=(j == 0), alignment=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)
        cur_x += cw

print("竞品对比页完成")

# ================================================================
# 第11页：多形态交互 + 远程控制拓扑
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
             "MULTI-FORM INTERACTION", font_size=12, color=ACCENT_CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.8),
             "多形态交互——随时随地使用", font_size=36, color=WHITE, bold=True)
add_decorative_line(slide, Inches(0.8), Inches(1.55), Inches(2), ACCENT_BLUE)

# 左侧：三种形态
forms = [
    ("🖥️", "桌面 AI 面板", ACCENT_BLUE, "Wails 桌面应用",
     ["本地 GUI 完整体验", "文件操作 · 屏幕截图 · 进程管理",
      "即时响应，零网络延迟"]),
    ("💬", "IM 消息通道", GREEN, "飞书 / 微信 / QQ / 钉钉 / 企微",
     ["通过 Hub 消息代理", "随时随地远程操控",
      "语音消息自动 ASR 转文字"]),
    ("⌨️", "TUI 终端", ACCENT_PURPLE, "Bubble Tea 终端界面",
     ["开发者友好", "SSH 服务器环境可用",
      "REST API (MaClawSrv)"]),
]

form_w = Inches(3.7)
form_h = Inches(2.3)
for i, (icon, title, color, sub, lines) in enumerate(forms):
    x = Inches(0.8) + i * (form_w + Inches(0.3))
    y = Inches(2.0)
    add_card(slide, x, y, form_w, form_h, title, [f"• {l}" for l in lines],
             icon, title_color=color, accent=color)
    add_text_box(slide, x + Inches(0.2), y + form_h - Inches(0.45), form_w - Inches(0.4), Inches(0.3),
                 sub, font_size=9, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# 下方：远程控制拓扑图
topo_card = add_shape(slide, Inches(0.8), Inches(4.6), Inches(11.7), Inches(2.5),
                       fill_color=DARK_CARD, line_color=RGBColor(0x1E, 0x2A, 0x50))
add_text_box(slide, Inches(1.1), Inches(4.7), Inches(5), Inches(0.4),
             "远程控制拓扑", font_size=16, color=ACCENT_BLUE, bold=True)
add_decorative_line(slide, Inches(1.1), Inches(5.1), Inches(1.0), ACCENT_BLUE, Pt(2))

# 拓扑流程
topo_steps = [
    ("Claude Code", "桌面端", ACCENT_BLUE),
    ("PTY", "终端", MID_GRAY),
    ("Desktop", "桌面操作", GREEN),
    ("WSS", "WebSocket", ACCENT_PURPLE),
    ("Hub", "消息代理", ORANGE),
    ("PWA", "手机/平板", ACCENT_CYAN),
]

step_w = Inches(1.6)
step_h = Inches(0.9)
total_topo_w = len(topo_steps) * step_w.inches + (len(topo_steps) - 1) * 0.3
topo_start_x = Inches((13.333 - total_topo_w) / 2)

for i, (name, desc, color) in enumerate(topo_steps):
    x = topo_start_x + Inches(i * (step_w.inches + 0.3))
    y = Inches(5.4)
    shape = add_shape(slide, x, y, step_w, step_h, fill_color=DARKER_CARD, line_color=color)
    add_text_box(slide, x, y + Inches(0.1), step_w, Inches(0.35),
                 name, font_size=12, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, y + Inches(0.45), step_w, Inches(0.3),
                 desc, font_size=9, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
    if i < len(topo_steps) - 1:
        add_text_box(slide, x + step_w, y + Inches(0.2), Inches(0.3), Inches(0.4),
                     "→", font_size=16, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

# 描述
add_text_box(slide, Inches(1.5), Inches(6.5), Inches(10.3), Inches(0.5),
             "IM 发消息 → Hub 路由 → MaClaw 桌面端执行 → 结果回传 → IM 收到回复，支持文件/截图/语音",
             font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

print("多形态交互页完成")

# ================================================================
# 第12页：意图理解 + 工具路由 + Steering
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
             "INTELLIGENT ROUTING", font_size=12, color=ACCENT_CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.8),
             "智能调度——精准理解意图，高效匹配工具", font_size=32, color=WHITE, bold=True)
add_decorative_line(slide, Inches(0.8), Inches(1.55), Inches(2), ACCENT_BLUE)

# 上方：三层意图分类
uic_card = add_shape(slide, Inches(0.8), Inches(2.0), Inches(11.7), Inches(2.0),
                      fill_color=DARK_CARD, line_color=RGBColor(0x1E, 0x2A, 0x50))
add_text_box(slide, Inches(1.1), Inches(2.1), Inches(5), Inches(0.4),
             "三层意图分类 (UIC)", font_size=16, color=ACCENT_BLUE, bold=True)
add_decorative_line(slide, Inches(1.1), Inches(2.5), Inches(1.0), ACCENT_BLUE, Pt(2))

uic_layers = [
    ("Layer 1", "关键词规则", "<1ms", "精确匹配高频模式", ACCENT_BLUE),
    ("Layer 2", "BM25 语义检索", "<5ms", "模糊匹配与工作流模板", GREEN),
    ("Layer 3", "LLM 多轮确认", "10-30s", "复杂意图理解与澄清", ORANGE),
]

for i, (name, method, latency, desc, color) in enumerate(uic_layers):
    x = Inches(1.0) + i * Inches(3.8)
    y = Inches(2.8)
    add_shape(slide, x, y, Inches(3.5), Inches(1.0), fill_color=DARKER_CARD, line_color=color)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.05), Inches(3.2), Inches(0.3),
                 f"{name}:  {method}", font_size=12, color=color, bold=True)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.35), Inches(1.5), Inches(0.3),
                 f"⏱ {latency}", font_size=10, color=YELLOW)
    add_text_box(slide, x + Inches(1.6), y + Inches(0.35), Inches(1.8), Inches(0.3),
                 desc, font_size=10, color=LIGHT_GRAY)
    if i < 2:
        add_text_box(slide, x + Inches(3.5), y + Inches(0.2), Inches(0.3), Inches(0.4),
                     "→", font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# 下方左右：工具路由 + Steering
left_card = add_shape(slide, Inches(0.8), Inches(4.3), Inches(5.7), Inches(2.7),
                       fill_color=DARK_CARD, line_color=RGBColor(0x1E, 0x2A, 0x50))
add_text_box(slide, Inches(1.1), Inches(4.4), Inches(5.2), Inches(0.4),
             "🛤️  工具路由系统", font_size=16, color=ACCENT_PURPLE, bold=True)
add_decorative_line(slide, Inches(1.1), Inches(4.8), Inches(1.0), ACCENT_PURPLE, Pt(2))
add_multi_para(slide, Inches(1.1), Inches(5.0), Inches(5.2), Inches(1.8),
    [f"•  按上下文关键词按需激活工具，不污染简单任务",
     f"•  渐进式暴露：核心工具始终可用，低频工具 discover_tool 按需加载",
     f"•  使用反馈闭环：成功/失败/重试记录回流到路由决策",
     f"•  高失败率工具自动降权 · 新工具渐进激活",
     f"•  DistillRoutingHints 聚合最近记录为路由提示",
    ], font_size=11, color=LIGHT_GRAY)

right_card = add_shape(slide, Inches(6.8), Inches(4.3), Inches(5.7), Inches(2.7),
                        fill_color=DARK_CARD, line_color=RGBColor(0x1E, 0x2A, 0x50))
add_text_box(slide, Inches(7.1), Inches(4.4), Inches(5.2), Inches(0.4),
             "⚙️  Steering 规则系统", font_size=16, color=ORANGE, bold=True)
add_decorative_line(slide, Inches(7.1), Inches(4.8), Inches(1.0), ORANGE, Pt(2))
add_multi_para(slide, Inches(7.1), Inches(5.0), Inches(5.2), Inches(1.8),
    [f"•  Markdown 文件声明行为规则，无需改代码",
     f"•  四种注入模式：始终 / 文件匹配 / 关键词匹配 / 手动引用",
     f"•  两级作用域：用户级 (~/.maclaw/steering/) 和项目级",
     f"•  Token 预算智能控制规则注入量",
     f"•  热加载：修改后 30 秒内自动生效",
    ], font_size=11, color=LIGHT_GRAY)

print("智能调度页完成")

# ================================================================
# 第13页：技术架构全景
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
             "TECH ARCHITECTURE", font_size=12, color=ACCENT_CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.8),
             "技术架构全景", font_size=36, color=WHITE, bold=True)
add_decorative_line(slide, Inches(0.8), Inches(1.55), Inches(2), ACCENT_BLUE)

# 架构层
arch_layers = [
    ("交互层", ACCENT_CYAN,
     ["Wails Desktop GUI", "IM 消息通道 (飞书/微信/QQ)", "Bubble Tea TUI", "REST API (MaClawSrv)"]),
    ("Agent 服务层", ACCENT_BLUE,
     ["统一 IMMessageHandler", "三层意图分类 (UIC)", "19 种工作流模板引擎",
      "工具路由 · Steering 规则", "AgentInstance 租户管理"]),
    ("核心能力层", ACCENT_PURPLE,
     ["记忆系统 (memory.Store)", "知识库 (knowledge.Store)", "结构化数据 (structureddata)",
      "经验蒸馏 (ExperienceDistiller)", "技能系统 (SkillHub)"]),
    ("执行层", GREEN,
     ["SSH 远程管理", "浏览器自动化 (CDP)", "桌面 GUI 自动化 (Accessibility+YOLO)",
      "文件系统操作", "进程管理 · 定时任务"]),
    ("存储层", ORANGE,
     ["SQLite WAL · FTS5", "PostgreSQL (Phase 2)", "向量索引 (Embedding)",
      "语义图谱 (SemanticGraph)", "文件系统 (~/.maclaw/)"]),
]

for i, (name, color, items) in enumerate(arch_layers):
    y = Inches(2.0) + i * Inches(1.05)
    # 层名
    label_bg = add_shape(slide, Inches(0.8), y, Inches(1.5), Inches(0.85),
                          fill_color=color, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text_box(slide, Inches(0.8), y + Inches(0.2), Inches(1.5), Inches(0.4),
                 name, font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # 项列表
    items_text = "  ·  ".join(items)
    add_shape(slide, Inches(2.5), y, Inches(9.9), Inches(0.85),
              fill_color=DARK_CARD, line_color=RGBColor(0x1E, 0x2A, 0x50),
              shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text_box(slide, Inches(2.7), y + Inches(0.15), Inches(9.5), Inches(0.55),
                 items_text, font_size=11, color=LIGHT_GRAY)

# 底部技术栈
add_shape(slide, Inches(0.8), Inches(7.0), Inches(11.7), Inches(0.35),
          fill_color=DARKER_CARD, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
add_text_box(slide, Inches(0.8), Inches(7.02), Inches(11.7), Inches(0.3),
             "Go 1.25  |  Wails v2  |  React 18  |  TypeScript  |  SQLite (WAL)  |  WebSocket  |  Bubble Tea (TUI)  |  WASM (wazero)  |  python-pptx",
             font_size=10, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

print("技术架构页完成")

# ================================================================
# 第14页：应用场景
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
             "SCENARIOS", font_size=12, color=ACCENT_CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.8),
             "应用场景——覆盖全行业工作场景", font_size=36, color=WHITE, bold=True)
add_decorative_line(slide, Inches(0.8), Inches(1.55), Inches(2), ACCENT_BLUE)

scenarios = [
    ("📝", "商业与战略", ACCENT_BLUE,
     ["商业计划书生成", "竞品分析报告", "项目提案撰写", "招投标文件", "创新方案设计"]),
    ("📊", "研发与技术", GREEN,
     ["软件开发全流程", "PRD 产品设计文档", "软件测试方案", "技术文档生成", "代码审查辅助"]),
    ("⚖️", "合规与法务", ORANGE,
     ["合同审查与比对", "尽职调查报告", "合规审计检查", "风险评估报告", "政策合规分析"]),
    ("🔬", "研究与学术", ACCENT_PURPLE,
     ["文献综述", "研究报告", "实验设计", "专利分析", "基金申请书"]),
    ("🖥️", "IT 运维", YELLOW,
     ["SSH 远程服务器管理", "日志分析与故障排查", "自动化部署脚本", "安全巡检", "监控告警处理"]),
    ("📊", "数据分析", ACCENT_CYAN,
     ["多源数据汇总分析", "经营分析报告", "KPI 监控与预警", "趋势分析与预测", "数据可视化"]),
]

sc_w = Inches(3.7)
sc_h = Inches(1.8)
for i, (icon, title, color, items) in enumerate(scenarios):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + col * (sc_w + Inches(0.3))
    y = Inches(2.0) + row * (sc_h + Inches(0.25))
    add_card(slide, x, y, sc_w, sc_h, title, [f"•  {it}" for it in items],
             icon, title_color=color, accent=color)

print("应用场景页完成")

# ================================================================
# 第15页：快速开始
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5),
             "QUICK START", font_size=12, color=ACCENT_CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.8),
             "四步上手，开箱即用", font_size=36, color=WHITE, bold=True)
add_decorative_line(slide, Inches(0.8), Inches(1.55), Inches(2), ACCENT_BLUE)

steps = [
    ("01", "安装", "下载 MaClaw 安装包\n一键安装桌面应用", ACCENT_BLUE),
    ("02", "配置 API", "配置 LLM API Key\n支持 OpenAI / Claude / DeepSeek", GREEN),
    ("03", "开始对话", "打开桌面应用或连接 IM\n自然语言描述你的需求", ACCENT_PURPLE),
    ("04", "获得成果", "工作流自动执行\n阶段产出等待你确认", ORANGE),
]

step_w = Inches(2.6)
step_h = Inches(3.5)
for i, (num, title, desc, color) in enumerate(steps):
    x = Inches(0.8) + i * (step_w + Inches(0.3))
    y = Inches(2.2)

    card = add_shape(slide, x, y, step_w, step_h, fill_color=DARK_CARD,
                     line_color=color)

    # 序号圆
    circle = add_shape(slide, x + step_w / 2 - Inches(0.35), y + Inches(0.3),
                        Inches(0.7), Inches(0.7), fill_color=color,
                        shape_type=MSO_SHAPE.OVAL)
    add_text_box(slide, x + step_w / 2 - Inches(0.35), y + Inches(0.38),
                 Inches(0.7), Inches(0.5),
                 num, font_size=24, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, x + Inches(0.2), y + Inches(1.2), step_w - Inches(0.4), Inches(0.4),
                 title, font_size=20, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_decorative_line(slide, x + Inches(0.5), y + Inches(1.65), step_w - Inches(1.0), color, Pt(2))
    add_text_box(slide, x + Inches(0.3), y + Inches(1.85), step_w - Inches(0.6), Inches(1.5),
                 desc, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # 连接箭头
    if i < 3:
        add_text_box(slide, x + step_w, y + step_h / 2 - Inches(0.2), Inches(0.3), Inches(0.4),
                     "→", font_size=20, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

print("快速开始页完成")

# ================================================================
# 第16页：结尾——核心理念
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, RGBColor(0x06, 0x09, 0x1C))

# 装饰线
add_decorative_line(slide, Inches(0), Inches(0), Inches(13.333), ACCENT_BLUE, Pt(4))

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.6),
             "MaClaw", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_decorative_line(slide, Inches(5.5), Inches(2.3), Inches(2.3), ACCENT_CYAN, Pt(3))

add_text_box(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(1.0),
             "不只是聊天，而是替你干活", font_size=32, color=ACCENT_CYAN,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(3.8), Inches(9), Inches(1.5),
             "内中外脑三层记忆  ·  19 种结构化工作流  ·  知识库蒸馏  ·  经验自进化\n"
             "SSH/GUI 全域操作  ·  多通道协作  ·  企业级结构化数据管理",
             font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# 底部标签
tags2 = ["开源", "自托管", "数据本地化", "多LLM支持", "零配置上手"]
tag_w2 = Inches(1.5)
total_w2 = len(tags2) * tag_w2.inches + (len(tags2) - 1) * 0.3
start_x2 = Inches((13.333 - total_w2) / 2)

for i, tag in enumerate(tags2):
    x = start_x2 + Inches(i * (tag_w2.inches + 0.3))
    add_shape(slide, x, Inches(5.5), Inches(tag_w2), Inches(0.45),
              fill_color=None, line_color=ACCENT_BLUE, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text_box(slide, x, Inches(5.52), Inches(tag_w2), Inches(0.4),
                 tag, font_size=12, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0), Inches(6.5), Inches(13.333), Inches(0.5),
             "github.com/RapidAI/CodeClaw",
             font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

print("结尾页完成")

# ========== 保存 ==========
output_dir = os.path.dirname(os.path.abspath(__file__))
output_path = os.path.join(os.path.dirname(output_dir), "MaClaw_产品介绍_CN.pptx")
prs.save(output_path)
print(f"\n✅ PPT 已保存: {output_path}")
print(f"共 {len(prs.slides)} 页")
